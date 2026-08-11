# -*- coding: utf-8 -*-
"""잇다 발표 — 「목록을 무엇으로 채울까」 한 장. (2026-08-11)

무엇을 담나
  슬롯 「갇힌 출력(enum)」의 **목록 자체가 어디서 왔는가**.
  사용자가 말한 흐름 그대로 — 찾았다 → 발견했다 → 막혀 있었기에 → 적용했다.

⚠ 톤 — 변명·자기비판을 넣지 않는다(사용자 지시, 2026-08-11).
  「사실 틀렸습니다」는 6분짜리 소개에 넣을 말이 아니다. **쓴 것만 쓴다.**
  안 쓴 것(Miller·Cowan)은 «같이 본 것»으로 한 줄, 그리고 그 자리를 우리 실측이 채웠다는
  사실만 적는다. 그건 변명이 아니라 우리가 가진 제일 센 숫자다.

원문 대조 (2026-08-11, 전부 1차 출처로 확인)
  · Prediger (1982) Dimensions underlying Holland's hexagon: Missing link between
      interests and occupations?  J. Vocational Behavior 21(3), 259–287
      인용문 확인: 흥미검사는 «작업 과제와 나란한 활동 선호»를 재기 때문에 작동한다
  · Holland (1959) A theory of vocational choice. J. Counseling Psychology 6(1), 35–45
      원래 명칭 확인: motoric·intellectual·esthetic·supportive·persuasive·conforming
  · Miller (1956) Psychological Review 63, 81–97  (psychclassics.yorku.ca 원문)
  · Cowan (2001) Behavioral and Brain Sciences 24, 87–114  (실제 주장 3~5 chunks)

⚠ 슬라이드에 «안» 쓰는 것
  · 「활동유형 9종은 RIASEC 기반」 — 코드에 기록이 없고 Conventional 이 통째로 빠졌다.
    「참고해 일상어로 재구성」까지가 말할 수 있는 최대다.
  · Miller 의 「7±2」를 칸 개수 근거로 대는 것 — 대상이 다르고(사람 기억 vs 스키마 필드),
    Cowan 이 3~5 로 내렸다. 우리 실측(7/7 → 4/7)이 더 정확하고 더 세다.

쓰는 법 (Backend/ 에서)
  python app/itda/scripts/build_ppt_itda_slotsrc.py
  → C:/Users/TFX255GS/Downloads/이음_잇다_목록출처.pptx   이후 verify_ppt.py 로 «열어» 확인
"""
import io
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r'C:\Users\TFX255GS\Downloads\이음_잇다_목록출처.pptx'

C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_EYE = RGBColor(0x8B, 0x84, 0xA0)
C_ACC = RGBColor(0x6B, 0x3F, 0xA0)
C_TXT = RGBColor(0x24, 0x1F, 0x2E)
C_RED = RGBColor(0xA8, 0x36, 0x4A)
C_DEEP = RGBColor(0x32, 0x16, 0x6F)
C_PANEL_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_BG = RGBColor(0xEF, 0xEB, 0xF6)
C_BG_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Malgun Gothic'


def _font(r, size, color, bold):
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.28, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        for (t, size, color, bold) in line:
            r = p.add_run()
            r.text = t
            _font(r, size, color, bold)
    return tb


def box(sl, x, y, w, h, fill=C_WHITE, line=C_PANEL_LN, lw=0.75):
    shp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.05
    except Exception:                                   # noqa: BLE001
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


B, N = True, False
T, S = 22, 17.5          # 본문 · 부연


def main():
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    text(sl, 1.00, 0.86, 12.0, 0.45,
         [[('PART 03 · 잇다 — 슬롯 목록의 출처', 18.75, C_EYE, False)]])
    text(sl, 14.0, 0.80, 5.0, 0.55, [[('발표자 김정준', 30, C_ACC, True)]],
         align=PP_ALIGN.RIGHT)
    text(sl, 1.00, 1.42, 18.1, 0.95,
         [[('목록을 무엇으로 채울까 — ', 48, C_TITLE, True),
           ('IT 밖에서 찾았습니다', 48, C_ACC, True)]])

    #  ⚠ 라벨은 3.35in 칸에 «한 줄»로 들어가야 한다 — 「② 그래서 어디를 찾았나」가
    #    두 줄로 접혔다(PNG 육안 검수, 2026-08-11). 21pt 한글은 11자쯤이 한계다.
    labels = ('①  막혀 있던 곳', '②  어디를 찾았나',
              '③  찾은 것', '④  그대로 써봤습니다')
    heights = (1.55, 1.55, 1.95, 1.77)
    fills = (C_WHITE, C_WHITE, C_WHITE, C_BG)
    lines_c = (C_PANEL_LN, C_PANEL_LN, C_PANEL_LN, C_BG_LN)

    bands = [
        #  ① 막혀 있던 곳
        [[('「갇힌 출력」은 정했습니다 — 모델이 ', T, C_TXT, N),
          ('목록 밖의 값은 낼 수 없게', T, C_ACC, B), ('.', T, C_TXT, N)],
         [('그런데 ', T, C_TXT, N), ('그 목록을 무엇으로 채울지', T, C_RED, B),
          ('가 남았습니다.', T, C_TXT, N)],
         [('가족돌봄청년을 대상으로 한 진로 데이터는 ', S, C_EYE, N),
          ('사업 자체가 최근이라 없었습니다', S, C_EYE, B), ('.', S, C_EYE, N)]],

        #  ② 어디를 찾았나
        [[('IT·NLP 문헌에는 답이 없었습니다 — 거기엔 ', T, C_TXT, N),
          ('「무엇을 물을지」가 없으니까요', T, C_TXT, B), ('.', T, C_TXT, N)],
         [('그래서 ', T, C_TXT, N), ('진로심리학', T, C_ACC, B),
          ('을 찾았습니다. 취약계층 진로상담이 수십 년 쌓인 분야입니다.', T, C_TXT, N)]],

        #  ③ 찾은 것
        [[('Prediger (1982) · Journal of Vocational Behavior 21(3), 259–287',
           19, C_ACC, B)],
         [('흥미 유형 여섯 개 ', T, C_TXT, N), ('«아래»에 두 축이 깔려 있다', T, C_DEEP, B),
          ('는 것을 밝힌 논문입니다.', T, C_TXT, N)],
         [('Data ↔ Ideas    ·    People ↔ Things', 25, C_ACC, B)],
         [('「흥미검사가 «작동하는» 이유는 그것이 작업 과제와 나란한 활동 선호를 재기 때문이다」',
           S, C_EYE, N)]],

        #  ④ 적용
        [[('다루는대상 6값 ', T, C_DEEP, B), ('← 그 네 자리를 검색에 쓸 만큼 세분화했습니다',
                                          T, C_TXT, N)],
         [('People=사람 · Things=기계·설비/자연·생물 · Data=컴퓨터·데이터/숫자·문서 · Ideas=창작물',
           S, C_EYE, N)],
         [('활동유형 9값 ', T, C_DEEP, B),
          ('← Holland(1959)의 여섯 유형을 ', T, C_TXT, N),
          ('일상어 「행위」로 풀어', T, C_ACC, B), (' 재구성했습니다', T, C_TXT, N)]],
    ]

    y = 2.78
    for i, (lb, h) in enumerate(zip(labels, heights)):
        box(sl, 1.0, y, 18.08, h, fill=fills[i], line=lines_c[i])
        text(sl, 1.42, y, 3.35, h, [[(lb, 21, C_ACC, True)]],
             anchor=MSO_ANCHOR.MIDDLE, ls=1.12)
        text(sl, 5.05, y, 13.7, h, bands[i], anchor=MSO_ANCHOR.MIDDLE, ls=1.30)
        y += h + 0.15

    #  ★ 2026-08-11 — 칸의 «개수»를 각주에서 «본문 띠»로 올렸다.
    #    처음엔 「Miller 를 뺀다」로 갔는데 사용자 지적이 옳았다 —
    #    지금 5개는 **Cowan(3~5) 범위 안**이라, 두 논문이 반박이 아니라 «둘 다 만족»이다.
    #  ⚠ 시작 숫자(7개? 8개?)는 «쓰지 않는다». 코드 주석으로는 8개(5+환경·시간·비용)로
    #    읽히는데 git 에 그 시점 기록이 없어 확인이 안 된다. 확인 안 된 숫자는 안 올린다.
    #    요지는 「몇 개로 시작했나」가 아니라 **「잴 때마다 줄였다」**다. 그쪽이 더 세다.
    box(sl, 1.0, 10.15, 18.08, 0.72, fill=C_BG, line=C_BG_LN)
    text(sl, 1.42, 10.15, 17.3, 0.72, [
        [('칸의 «개수» — ', 18.5, C_DEEP, B),
         ('인지 부하 연구가 이 근처를 봅니다. Miller(1956) 「7±2」 · Cowan(2001)이 다시 재서 「3~5」.',
          18.5, C_TXT, N)],
        [('저희는 잴 때마다 줄였습니다 — 제약 3종(직업을 안 가름) · 관심대분류(지목 7/7→4/7) · 강점성향(400발화 0회). ',
          17, C_EYE, N),
         ('지금 5개입니다.', 17, C_ACC, B)],
    ], anchor=MSO_ANCHOR.MIDDLE, ls=1.24)

    prs.save(OUT)
    print(f'저장: {OUT} · 1장')


if __name__ == '__main__':
    main()
