# -*- coding: utf-8 -*-
"""「슬롯을 채우기 전에 — 관문 7개」 구조 그림. (2026-08-07)

사용자 요청
  「슬롯이 채워지기 전에, 7개의 게이트로 막아낸다. 같은 형식으로 하나 만들면 될 듯.
   이건 좀 보기 쉽게. pre_check 같은 경우도 안에 있는 거 모두 보여줘야 하니까」

⚠ 주석이 아니라 «코드»에서 뽑았다. 오늘 pre_check 를 세 번 고쳤기 때문이다.
   · 관문 순서 : itda_core.py `_step()` 5343~5560행의 «실제 호출 순서»
   · pre_check : 2308행 함수 본문의 «실제 분기 순서»
   · 예시 결과 : 아래 값들은 전부 `pre_check()` 를 직접 «돌려서» 얻었다 (2026-08-07)
        'ㅎㅇ'→None · '...'→SILENT · 'ㅠㅠ'→SILENT · 'ㅁㄴㅇㄹ'→VAGUE
        '죽고 싶어요'→SELFHARM · '자살예방 상담사가 되고싶어요'→None
        '할머니가 저를 때려요'→HARM · '저 혼자 해요'→None
        '아 씨발 짜증나'→None · '장애인 활동지원사'→None
   · 낱말 116개 : SELF_HARM 27 + HARM_OTHERS 17 + ABUSE 4 + _ABUSE_VARIANT 35
                  + _SEXUAL_HARM 9 + _SEXUAL_HARD 13 + _SEXUAL_SOFT 11 = 116  (len() 으로 셈)

⚠ 인수인계장의 「통과시키는 예외 73개」는 «재현 못 했다».
   내가 셀 수 있었던 것: _SH_SAFE_WORD 4 + _SH_CAREER_CTX 22 + _HARM_IDIOM 11 + _JAMO_OK 8 = 45.
   나머지는 함수 안 조건문이라 목록으로 안 세진다. ⇒ 그림에 «숫자를 안 쓴다».

쓰는 법
  python make_gate_diagram.py [바탕.pptx] [출력.pptx]
"""
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                           # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v6.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v7.pptx'

IN = 914400
FONT = 'Arial'
C_HEAD    = RGBColor(0x8B, 0x84, 0xA0)
C_PRESENT = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE   = RGBColor(0x3B, 0x20, 0x63)
C_BODY    = RGBColor(0x55, 0x50, 0x6B)
C_CARD    = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_PANEL   = RGBColor(0xEF, 0xEB, 0xF6)
C_PANEL_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_MUTED   = RGBColor(0x8B, 0x84, 0xA0)
#  결과 색 — «차단»과 «차단 아님»을 한눈에 가른다. 이게 이 그림의 핵심이다.
C_PASS    = RGBColor(0x6E, 0x8B, 0x74)     # 통과 (차분한 초록)
C_ASK     = RGBColor(0x6B, 0x3F, 0xA0)     # 되묻기 (보라)
C_CONNECT = RGBColor(0xC2, 0x6B, 0x2E)     # ★ 차단 «안 하고» 잇는다 (주황)
C_BLOCK   = RGBColor(0xA8, 0x36, 0x4A)     # 차단 (붉은 자주)


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.25):
    tb = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                  Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.color.rgb = col
    return tb


def box(slide, x, y, w, h, fill, line):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Emu(int(x * IN)), Emu(int(y * IN)),
                                Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.02
    except Exception:
        pass
    sh.text_frame.text = ''
    return sh


# ── 관문 7개 — _step() 의 실제 호출 순서 ──────────────────────────
GATES = [
    ('①', '이미 잠긴 사람인가',   'abuse_limits',      '차단'),
    ('②', '「그거 말고」인가',     'rejects_last_card', '후보에서 뺀다'),
    ('③', '「괜찮아요」인가',      'declines_policy',   '안내를 멈춘다'),
    ('④', '불법 신호인가',        'illegal_signal',    '«세기만» 한다'),
    ('⑤', 'pre_check',          '← 오른쪽에 펼침',     ''),
    ('⑥', '「1번이요」인가',       'pick_from_options', '바로 카드로'),
    ('⑦', '공격인가',            'is_injection',      '되돌린다'),
]

# ── pre_check 안 — 실제 분기 순서 · 결과는 직접 돌려서 확인 ─────────
#    (라벨, 예시 입력, 결과, 결과색, 꼬리설명)
#  ⚠ 꼬리설명은 «짧게». 처음엔 길게 썼다가 두 줄로 넘쳐 다음 행을 덮었다.
#    긴 설명은 맨 아래 ★ 한 줄로 뺀다 — 거기가 제일 중요한 말이기도 하다.
PRE = [
    ('통과시킨다', [
        ('자주 쓰는 자모 축약 8개',   '“ㅎㅇ”',                  '통과',    C_PASS, ''),
        ('어절 «시작»에서만 본다',    '“저 혼자 해요”',            '통과',    C_PASS, '어절 중간은 안 봄'),
        ('진로 이야기는 예외',        '“자살예방 상담사가 되고싶어요”', '통과',    C_PASS, ''),
        ('혼잣말 좌절은 안 막는다',   '“아 씨발 짜증나”',           '통과',    C_PASS, '상대 향할 때만'),
    ]),
    ('말 대신 쓰는 표현', [
        ('점만 찍고 보냈다',         '“…”',                    'SILENT', C_ASK, '사다리로'),
        ('감정 표현만 있다',         '“ㅠㅠ”  “😭”  “^^”',       'SILENT', C_ASK, '오타냐 안 물음'),
        ('키보드 난타',             '“ㅁㄴㅇㄹ”',                'VAGUE',  C_ASK, '되묻는다'),
    ]),
    ('위험 — 여기서 «갈린다»', [
        ('자기 위해',              '“죽고 싶어요”',             'SELFHARM', C_CONNECT, '아래 ★'),
        ('폭력 · 성적 피해',        '“할머니가 저를 때려요”',      'HARM',     C_CONNECT, '2층에 물음'),
        ('명백한 희롱',            '',                        'UNSAFE',   C_BLOCK, '차단'),
        ('상대를 향한 욕설',        '',                        'UNSAFE',   C_BLOCK, '차단'),
    ]),
]


def build(prs):
    s = blank_like(prs)

    text(s, 1.00, 0.92, 8.0, 0.30, [('PART 03 · 잇다 — 예외 처리', 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRESENT)])
    text(s, 1.00, 1.50, 19.0, 0.95,
         [('슬롯을 채우기 «전에», 코드가 먼저 봅니다', 44, True, C_TITLE)])

    # ══ 좌: 관문 7개 ══════════════════════════════════════════
    LX, LW = 1.00, 6.30
    box(s, LX, 2.60, LW, 7.70, C_CARD, C_CARD_LN)
    text(s, LX + 0.40, 2.90, LW - 0.80, 0.40,
         [('관문 7개 — 전부 코드', 22, True, C_TITLE)])
    text(s, LX + 0.40, 3.36, LW - 0.80, 0.30,
         [('LLM 을 부르기 전에 순서대로 지납니다', 15, False, C_MUTED)])

    y = 3.92
    for mark, q, fn, res in GATES:
        hi = (mark == '⑤')
        text(s, LX + 0.40, y, 0.45, 0.34, [(mark, 19, True, C_PRESENT if hi else C_MUTED)])
        text(s, LX + 0.95, y, 3.05, 0.34,
             [(q, 17.5, True, C_TITLE)])
        text(s, LX + 0.95, y + 0.34, LW - 1.35, 0.28,
             [(fn, 13.5, False, C_MUTED)])
        if res:
            text(s, LX + 4.05, y, LW - 4.45, 0.34,
                 [('→ ' + res, 15, False, C_BODY)])
        y += 0.79

    text(s, LX + 0.40, 9.62, LW - 0.80, 0.62,
         [('여기서 걸리면 LLM 을 한 번도 안 부릅니다', 17, True, C_PRESENT),
          ('턴당 0원 · 약 1밀리초', 15, False, C_MUTED)])

    # ══ 우: pre_check 펼침 ════════════════════════════════════
    RX, RW = 7.70, 11.30
    box(s, RX, 2.60, RW, 7.70, C_PANEL, C_PANEL_LN)
    text(s, RX + 0.45, 2.90, RW - 0.90, 0.40,
         [('⑤ pre_check 안 — 무엇을 어디로 보내나', 22, True, C_TITLE)])
    text(s, RX + 0.45, 3.36, RW - 0.90, 0.30,
         [('낱말 116개는 «차단기»가 아니라, 누구에게 물어볼지 정하는 «갈림길»입니다',
           15, False, C_MUTED)])

    #  ⚠ 세로 예산 — 패널은 Y2.60~10.30 (7.70"). 처음에 행 0.44 · 틈 0.18 로 잡았다가
    #    마지막 그룹이 10.48 까지 내려가 «하단 ★ 문장과 겹쳤다». 0.40/0.12 로 줄인다.
    #      3.80 + (0.36+4×0.40+0.12) + (0.36+3×0.40+0.12) + (0.36+4×0.40) = 9.52
    y = 3.80
    for group, rows in PRE:
        text(s, RX + 0.45, y, RW - 0.90, 0.30,
             [(group, 15, True, C_PRESENT)])
        y += 0.36
        for label, ex, res, col, tail in rows:
            text(s, RX + 0.60, y, 2.90, 0.30, [(label, 15.5, False, C_BODY)])
            text(s, RX + 3.60, y, 3.55, 0.30, [(ex, 15.5, True, C_TITLE)])
            text(s, RX + 7.25, y, 1.75, 0.30, [(res, 15.5, True, col)])
            if tail:
                text(s, RX + 9.10, y, RW - 9.55, 0.30, [(tail, 12.5, False, C_MUTED)])
            y += 0.40
        y += 0.12

    text(s, RX + 0.45, 9.62, RW - 0.90, 0.62,
         [('★ 자해·폭력은 «차단하지 않습니다» — ☎109·129 로 «잇고»,', 16.5, True, C_CONNECT),
          ('남용 카운터에 «절대» 안 셉니다. 막는 자리가 아니라 «잇는» 자리입니다.',
           16.5, True, C_CONNECT)])

    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    build(prs)
    prs.save(OUT)
    W, H = prs.slide_width / IN, prs.slide_height / IN
    over = [sh for sh in Presentation(OUT).slides[-1].shapes
            if (sh.left + sh.width) / IN > W + 0.01 or (sh.top + sh.height) / IN > H + 0.01]
    print(f'  {n0}장 → {len(Presentation(OUT).slides)}장 (그림 1장 추가)')
    print(f'  슬라이드 밖으로 나간 도형: {len(over)}개  {"✅" if not over else "🔴"}')
    print(f'  저장 → {OUT}')


main()
