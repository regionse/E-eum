# -*- coding: utf-8 -*-
"""잇다 발표 — 「왜 이 층들인가」 1장. (2026-08-10)

왜 만드나
  발표에서 나올 두 질문에 «숫자»로 답하려는 것이다.
    ① 「k=60 은 왜죠?」            ② 「엔트로피까지 왜 필요해요? 리랭커면 되잖아요」
  둘 다 어제까지는 답이 「표준값이라서」·「그게 나아 보여서」뿐이었다.

숫자 출처 — 전부 오늘 실측
  checks/gate_ablation.py · 정답표 100건 · 원값 checks/_gate_ablation.json
  실측기록 §18 에 표와 한계를 그대로 적어 뒀다.

⚠ 자기완결로 만든다. build_ppt_itda_main.py 를 다른 세션이 편집 중이라
  같은 파일을 건드리면 서로 덮어쓴다(오늘 GATE_MODEL 이 실제로 그렇게 날아갔다).

쓰는 법 (Backend/ 에서)
  python app/itda/scripts/build_ppt_itda_why.py
  → C:/Users/TFX255GS/Downloads/이음_잇다_왜.pptx   이후 verify_ppt.py 로 «열어» 확인
"""
import io
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r'C:\Users\TFX255GS\Downloads\이음_잇다_왜.pptx'

C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_EYE = RGBColor(0x8B, 0x84, 0xA0)
C_ACC = RGBColor(0x6B, 0x3F, 0xA0)
C_TXT = RGBColor(0x24, 0x1F, 0x2E)
C_RED = RGBColor(0xA8, 0x36, 0x4A)
C_DEEP = RGBColor(0x32, 0x16, 0x6F)
C_PANEL_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_BG = RGBColor(0xEF, 0xEB, 0xF6)
C_BG_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Malgun Gothic'


def _font(r, size, color, bold):
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.18, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        for (t, size, color, bold) in line:
            r = p.add_run()
            r.text = t
            _font(r, size, color, bold)
    return tb


def box(sl, x, y, w, h, fill=C_WHITE, line=C_PANEL_LN, lw=0.75):
    shp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.06
    except Exception:                                   # noqa: BLE001
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def row(sl, x, y, cols, widths, size=17.5, bold=False, color=C_TXT):
    """표 한 줄 — 열 시작 x 를 누적해서 찍는다."""
    cx = x
    for c, w in zip(cols, widths):
        text(sl, cx, y, w, 0.36, [[(c, size, color, bold)]], align=PP_ALIGN.RIGHT)
        cx += w


def main():
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    text(sl, 1.00, 0.86, 12.0, 0.45, [[('PART 03 · 잇다 — 왜 이 층들인가', 18.75, C_EYE, False)]])
    text(sl, 14.0, 0.80, 5.0, 0.55, [[('발표자 김정준', 30, C_ACC, True)]], align=PP_ALIGN.RIGHT)
    text(sl, 1.00, 1.42, 18.1, 0.95,
         [[('층을 하나씩 얹어 «각각 무엇을 벌었나» 쟀습니다', 46, C_TITLE, True)]])
    #  2026-08-10 — 아래 「⚠ 한계」 상자를 걷고 조건을 이 줄로 올렸다.
    #    ⚠ 로 사과하듯 다는 건 발표에서 변명으로 읽힌다. 조건은 «측정 방법»으로 적으면 된다.
    text(sl, 1.00, 2.32, 18.1, 0.42,
         [[('측정 조건 — 정답표 100건 · 질의 1개 · 검색 단계만 (생성 LLM 0회) · 실측기록 §18',
            19, C_EYE, False)]])

    #  ── 왼쪽: 층별 어블레이션 ──────────────────────────────────
    box(sl, 1.0, 2.92, 10.4, 5.35)
    text(sl, 1.42, 3.14, 9.6, 0.5, [[('왜 리랭커인가 · 왜 엔트로피인가', 29, C_TITLE, True)]])

    W = [3.9, 1.5, 2.3, 1.9]
    row(sl, 1.42, 3.86, ['조건', '답한율', '답했을 때 정답률', '막기 적중'],
        W, size=16, bold=True, color=C_EYE)
    y = 4.28
    for name, cov, acc, blk, hl in [
        ('RRF만', '100%', '53.0%', '—', False),
        ('+ 리랭커', '100%', '75.0%', '—', False),
        ('+ TOP% 만', '65%', '83.1%', '40%', False),
        ('+ 엔트로피 만', '76%', '82.9%', '50%', False),
        ('+ TOP% & 엔트로피', '57%', '87.7%', '42%', True),
    ]:
        c = C_ACC if hl else C_TXT
        row(sl, 1.42, y, [name, cov, acc, blk], W, size=18.5, bold=hl, color=c)
        y += 0.50
    text(sl, 1.42, y + 0.06, 9.6, 1.5, [
        [('리랭커가 ', 19, C_TXT, False), ('+22점', 22, C_ACC, True),
         ('   엔트로피가 «답할 때의 정확도»를 ', 19, C_TXT, False), ('+12.7점', 22, C_ACC, True)],
        [('대신 답하는 비율이 57% 로 준다 — 확신 없으면 «되묻는다»', 17.5, C_EYE, False)],
        [('막은 43건 중 18건은 리랭커도 틀렸을 것이었다', 17.5, C_EYE, False)],
    ], ls=1.34)

    #  ── 오른쪽: k 쓸기 ───────────────────────────────────────
    box(sl, 11.8, 2.92, 7.3, 5.35)
    text(sl, 12.22, 3.14, 6.5, 0.5, [[('왜 k = 1 인가', 29, C_TITLE, True)]])
    KW = [1.5, 1.6, 1.6]
    row(sl, 12.22, 3.86, ['k', '@1', '@8'], KW, size=16, bold=True, color=C_EYE)
    y = 4.28
    for k, a1, a8, mark in [('0', '68', '95', ''), ('1', '63', '95', '채택'),
                            ('2', '64', '94', ''), ('60', '53', '86', '이전')]:
        c = C_ACC if mark == '채택' else (C_RED if mark == '이전' else C_TXT)
        row(sl, 12.22, y, [k, a1, a8], KW, size=18.5, bold=bool(mark), color=c)
        if mark:
            text(sl, 17.1, y, 1.9, 0.36, [[('← ' + mark, 16, c, True)]])
        y += 0.50
    text(sl, 12.22, y + 0.06, 6.5, 1.7, [
        [('상위 8개를 리랭커가 다시 줄 세운다.', 17.5, C_TXT, True)],
        [('그래서 «1위냐(@1)»보다 «풀에 들어왔나(@8)»가 중요하다.', 17.5, C_TXT, False)],
        [('@8 이 ', 18, C_TXT, False), ('86 → 95', 22, C_ACC, True),
         ('  — 이게 실익이다.', 18, C_TXT, False)],
    ], ls=1.34)

    #  ── 아래: 왜 60이 나빴나 + 한계 ───────────────────────────
    box(sl, 1.0, 8.48, 18.1, 2.28, fill=C_BG, line=C_BG_LN)
    text(sl, 1.5, 8.86, 17.1, 1.7, [
        [('왜 60 이 아니라 1 인가 — ', 21, C_DEEP, True),
         ('k 가 클수록 「몇 개 목록에 나왔나」가 「어디서 1등이냐」를 이깁니다.', 21, C_TXT, True)],
        [('Cormack 2009 의 60 은 «수십 개» 시스템을 합칠 때 나온 값입니다. 잇다는 «둘»(벡터·글자)만 합치니, '
          '표를 셀 재료가 없는데 표 세기 모드로 쓰고 있었던 셈입니다.', 18, C_EYE, False)],
        [('2023년 ACM TOIS 도 같은 결론입니다 — 작은 k 가 9개 데이터셋 중 8개에서 이겼습니다.',
          18, C_EYE, False)],
    ], ls=1.46)

    prs.save(OUT)
    print(f'저장: {OUT} · 1장')


if __name__ == '__main__':
    main()
