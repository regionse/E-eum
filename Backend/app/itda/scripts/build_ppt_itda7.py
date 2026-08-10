# -*- coding: utf-8 -*-
"""잇다 발표 — **7분 · 7장**으로 다시 짓는다 (2026-08-09).

왜 다시 짓나
  · 발표 총 37분 · 3명 → 잇다 몫이 **7분**이다. 지금 10장은 장당 42초라 못 맞춘다.
  · 그리고 **숫자가 낡았다**(2026-08-09 실측으로 정정):
      턴당 LLM 1.1회 → «2~3회»  ·  턴당 1.2원 → «2.42원»  ·  캐시 60% → «38%»
      「① 입력 방어 — LLM 0회」 → 유해·위기 게이트가 «매 턴» 돈다
  · 오늘 만든 것이 하나도 없다 — 「종류」(원함/해봤음/못함) · 유해·위기 게이트 ·
    정체 게이지 · 먼저 보여주기.

만드는 방식 (build_ppt_itda.py 와 같다)
  이 덱은 레이아웃 자리표시자가 없고 슬라이드마다 도형을 직접 놓은 구조다.
  python-pptx 로 «새로» 만들면 디자인이 통째로 날아간다.
  ⇒ 기존 슬라이드 XML 을 복제하고 **글자만** 바꾼다.

⚠ notesSlide 는 복제하지 않는다 — PowerPoint 가 파일을 거부한다(2026-08-07 실측).
⚠ 만들고 나면 반드시 `verify_ppt.py` 로 열어 본다. python-pptx 가 통과해도 PowerPoint 는 거부한다.

쓰는 법
  python -m app.itda.scripts.build_ppt_itda7 <원본.pptx> <출력.pptx>
"""
import sys
import io
import copy

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else ''
OUT = sys.argv[2] if len(sys.argv) > 2 else ''

#  현재 덱에서 «잇다» 구간 (1-based). 이 구간을 통째로 갈아끼운다.
ITDA_FROM, ITDA_TO = 18, 27
#  본보기 (1-based, 원본 기준)
TPL_COVER = 18      # 표지 — 큰 제목 + 항목 3개
TPL_FLOW = 19       # 4단 흐름 (제목+본문 × 4)
TPL_TROUBLE = 26    # 2단 카드 (TROUBLE)


def clone(prs, idx0):
    """idx0(0-based) 슬라이드를 복제해 «맨 뒤»에 붙인다."""
    src = prs.slides[idx0]
    dst = prs.slides.add_slide(src.slide_layout)
    for sh in list(dst.shapes):                 # add_slide 가 넣은 빈 자리표시자 제거
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(sh._element))
    for k, rel in src.part.rels.items():
        #  ⚠ notesSlide 는 절대 따라가지 않는다 — PowerPoint 가 파일을 거부한다.
        if rel.reltype.endswith('slideLayout') or rel.reltype.endswith('notesSlide'):
            continue
        if rel.is_external:
            dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            dst.part.rels.get_or_add(rel.reltype, rel._target)
    return dst


def put(shape, text):
    """도형의 글자를 바꾼다 — 첫 문단 서식을 유지한 채 줄만 갈아 끼운다."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    lines = str(text).split('\n')
    p0 = tf.paragraphs[0]
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    if not p0.runs:
        p0.add_run()
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    p0.runs[0].text = lines[0]
    from pptx.text.text import _Paragraph
    for ln in lines[1:]:
        newp = copy.deepcopy(p0._p)
        p0._p.getparent().append(newp)
        np = _Paragraph(newp, tf)
        for r in list(np.runs[1:]):
            r._r.getparent().remove(r._r)
        np.runs[0].text = ln


def fill(slide, mapping):
    shapes = list(slide.shapes)
    for i, txt in mapping.items():
        if txt is None or i >= len(shapes):
            continue
        put(shapes[i], txt)


def reorder(prs, order):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for e in ids:
        lst.remove(e)
    for i in order:
        lst.append(ids[i])


# ── 본보기별 채우기 ────────────────────────────────────────────
#  FLOW(19): 0 헤더 · 1 제목 · 3/6/9/12 소제목 · 4/7/10/13 본문 · 14 발표자
def slide_flow(prs, head, title, blocks):
    s = clone(prs, TPL_FLOW - 1)
    m = {0: head, 1: title}
    for (ti, bi), (t, b) in zip(((3, 4), (6, 7), (9, 10), (12, 13)), blocks):
        m[ti], m[bi] = t, b
    fill(s, m)
    return s


#  TROUBLE(26): 0 헤더 · 2 제목 · 4 번호 · 5 소제목 · 7,8,9 본문
#                              11 번호 ·12 소제목 ·14,15,16 본문
def slide_trouble(prs, head, title, left, right):
    s = clone(prs, TPL_TROUBLE - 1)
    fill(s, {0: head, 2: title,
             4: left[0], 5: left[1], 7: left[2], 8: left[3], 9: left[4],
             11: right[0], 12: right[1], 14: right[2], 15: right[3], 16: right[4]})
    return s


H = 'PART 03 · 잇다'


def build(prs):
    made = []

    # ── 1. 표지 (30초) ────────────────────────────────────────
    s = clone(prs, TPL_COVER - 1)
    fill(s, {2: '말을 «칸»으로 바꿔서 찾습니다',
             6: '1 · 「돌보느라 못 했다」를 「돌봄을 원한다」로 안 읽기',
             7: '2 · LLM 을 믿지 않고 코드가 다시 보는 구조',
             8: '3 · 100턴을 견디게 — 그래서 제일 싼 모델로 턴당 2.4원'})
    made.append(s)

    # ── 2. 들어온 말이 먼저 지나는 곳 (1분) ────────────────────
    made.append(slide_flow(
        prs, f'{H} — 관문',
        '들어온 말이 먼저 지나는 곳',
        [('① 코드가 먼저 — LLM 0회 · 0원 · 1밀리초',
          '거부·번복 / 정책거절 / 남용 / 칩 선택 / 인젝션.\n'
          '낱말 1,010개는 «차단기»가 아니라 «어디로 보낼지» 정합니다.\n'
          '★ 자해·폭력은 막지 않고 ☎109·129 로 «잇습니다». 남용에 안 셉니다.'),
         ('② 안전 게이트 — LLM 이 «매 턴» 한 번 더',
          '낱말표가 「총 개조하는 법」·「대 마 관련 일자리」를 4건 중 «0건» 잡았습니다.\n'
          '한국어 유해 표현은 자모까지 쪼개져서 목록으로 못 닫습니다.\n'
          '⇒ 유해 · 위기 · 착지요청 · 사정 — 네 가지를 «한 호출»에 판정합니다.'),
         ('③ 본문과 «나란히» 돌립니다',
          '순서대로 돌리면 매 턴 +0.9초입니다. 나란히 던져 지연을 없앴습니다.\n'
          '유해면 본문 결과를 «쓰지 않고 버립니다» — 슬롯도 안 담깁니다.'),
         ('실측 — 공격 6/6 차단 · 오탐 0/10',
          '대마판매 · 총개조 · 「대 마」 · 「여자 직원만」 · 「필ㄹㅗ폰」 · 사제총기 → 전부 차단.\n'
          '용접 · 총포화약안전기술 · 경찰 · 장례지도사 · 술집 → 전부 통과.\n'
          '에두른 위기(「다들 저 없는 게 나을 거예요」) 8건 → 8건 모두 ☎109.')]))

    # ── 3. 슬롯 — 왜 넣었나 (1분 30초) ────────────────────────
    made.append(slide_flow(
        prs, f'{H} — 슬롯',
        '한 문장을 «칸»으로 바꿉니다',
        [('왜 칸으로 바꾸나',
          '🧑「할머니 돌보느라 학교를 그만뒀어요. 손으로 만드는 건 재밌었는데…」\n'
          '검색은 「뭘 해야 할지」를 못 읽습니다. 문장 그대로는 못 찾습니다.\n'
          '그리고 칸으로 바꾸면 **100턴이 가도 상태가 7칸**입니다.'),
         ('무엇을 받나 — 값 · 근거 · «종류»',
          '값 「돕기·돌봄」 / 근거 「돌보느라」 / 종류 «못함»\n'
          '근거는 원문에서 그대로 인용하게 하고, 없으면 그 칸을 버립니다.\n'
          '★ 「종류」가 이 서비스의 핵심입니다 — 원함 · 해봤음 · «못함».'),
         ('왜 「종류」가 필요했나',
          '「할머니 돌보느라 못 했어요」를 「돌봄을 원한다」로 읽으면 요양보호사가 나옵니다.\n'
          '가족돌봄청년에게 그건 이 서비스가 하면 안 되는 오독입니다.\n'
          '「~하고 싶어요」와 「~해야 해요」는 정반대입니다. 어미를 봅니다.'),
         ('「종류」가 네 곳에 걸립니다',
          '① 검색어에서 «못함»을 뺍니다  ② 순위 가중에서도 뺍니다\n'
          '③ 되묻는 선택지로 «다시 권하지 않습니다»  ④ 다음 턴 프롬프트에 함께 보냅니다\n'
          '실측 — 「못함」 재권유 2/4 → «0/5» · 돌봄 직업 1개 → «0개» · 제빵 2위 → «1위».')]))

    # ── 4. 찾기 (1분) ─────────────────────────────────────────
    made.append(slide_flow(
        prs, f'{H} — 찾기',
        '칸이 찼습니다. 이제 «찾습니다»',
        [('① 칸을 이어 붙여 질의를 만듭니다',
          '관심분야 · 활동유형 · 다루는대상 → 「제과제빵 만들기 사람」\n'
          '★ 이때 «못함»으로 표시된 값은 «안 들어갑니다».\n'
          '  「돌봄을 못 했다」로 요양보호사를 찾으면 안 되니까요.'),
         ('② 두 갈래로 «따로» 찾습니다',
          '벡터 — Pinecone · gemini-embedding-2 · 3,072차원 · 코사인 (뜻이 가까운 것)\n'
          'FULLTEXT — MySQL ngram (글자가 맞는 것)\n'
          'NCS 직업 1,094개 · 자격증 613종 · 강좌 8,273개'),
         ('③ RRF 로 순위를 합칩니다 (k=60)',
          '점수가 아니라 «몇 등이었나»를 더합니다.\n'
          '⚠ 무게를 주려 했다가 되돌렸습니다 — k=60 이면 1/61 과 1/66 이 거의 같아\n'
          '  RRF 는 사실상 «표 세기»입니다. 무게로 이길 수 있는 구조가 아닙니다.'),
         ('④ 크로스인코더로 다시 줄 세웁니다',
          'Jina reranker v3.5 — 질의와 후보를 «붙여서 함께» 읽습니다.\n'
          '임베딩은 따로 벡터로 만들어 거리를 재고, 크로스인코더는 함께 읽습니다.\n'
          '정확한 대신 느립니다. 그래서 8개 후보에만 씁니다.')]))

    # ── 5. 답 내기 (1분) ──────────────────────────────────────
    made.append(slide_flow(
        prs, f'{H} — 답 내기',
        '바로 줄까, 한 번 더 물을까',
        [('1단 TOP% — 방향이 정해졌나',
          '후보를 NCS 중분류로 묶고 1등 묶음의 점유율을 봅니다.\n'
          '0.90 미만 → 서로 다른 동네가 섞였다 → «방향 칩»으로 되묻습니다.'),
         ('2단 엔트로피 — 그 안에서 또 갈리나',
          '한 동네로 좁혀져도 그 안에서 고르게 흩어져 있을 수 있습니다.\n'
          '0.70 이상 → «세부 칩»으로 한 번 더. 그 아래면 카드를 냅니다.'),
         ('★ 그런데 «부정만 하는 사람»이 있습니다',
          '「저랑 안 맞아서」·「부담스럽고」·「다 힘들 것 같아요」 — 말할수록 무게가 0에 가까워집니다.\n'
          '18턴을 성실히 답한 사람이 카드를 «한 번도» 못 봤습니다.\n'
          '⇒ 대화가 3턴 동안 안 나가면 먼저 보여줍니다. 「아니면 왜 아닌지만 알려주세요」'),
         ('카드 — 없으면 «없다»고 말합니다',
          '직업 + 자격증 ≤3 + K-MOOC 강좌 + 국비훈련 링크 + «이 방향이 아니라면».\n'
          '자격증이 없는 직업이면 억지로 안 붙이고 내일배움카드를 안내합니다.\n'
          '⚠ 응시 자격·시험 일정·급여는 «말하지 않습니다». 우리가 아는 게 아닙니다.')]))

    # ── 6. 문제와 해결 ① (1분 30초) ───────────────────────────
    made.append(slide_trouble(
        prs, f'{H} — 문제와 해결',
        '사용자를 막고, 마음을 연 순간에 상품을 내밀었습니다',
        ('TROUBLE 01', '맞고 있는 사람을 가해자로 만들었다',
         '원인 — 「치매 할머니가 자꾸 저를 때려요」가 차단되고 남용 카운트까지 올랐다. '
         '폭력 낱말표는 «누구의 행동인지»를 모른다.',
         '해결 — 낱말을 «차단»이 아니라 «라우팅»으로. 걸리면 가해·피해·제3자만 묻고, '
         '피해면 ☎129 와 나누다 지도로 «잇는다».',
         '배운 것 — 실측 21케이스 중 9건이 오차단. 낱말이 할 수 있는 건 「여기 봐야 한다」까지다.'),
        ('TROUBLE 02', '마음을 연 순간에 직업 카드를 내밀었다',
         '원인 — 🤖「이건 좀 아니다 싶은 일이 있으세요?」 🧑「술 냄새요.」(아버지가 알코올 의존) '
         '→ 🤖 [게임 카드]. 그 대화는 「그만해요」로 끝났다.',
         '해결 — 사정을 말한 턴에는 카드를 내지 않는다. 게이트가 «사정»을 함께 판정하고, '
         '8자 이하 짧은 답은 판정하지 말고 «미룬다».',
         '배운 것 — 판정하려 들지 말고 미루면 된다. 다음 턴에 다시 기회가 온다.')))

    # ── 7. 숫자 (1분) ─────────────────────────────────────────
    made.append(slide_flow(
        prs, f'{H} — 잰 것',
        '전부 «재서» 정했습니다',
        [('비용 — 턴당 2.42원',
          'gemini-3.1-flash-lite · 입력 $0.25 / 출력 $1.50 / 캐시읽기 $0.025 per 1M\n'
          '돈은 «본문 한 호출»에 77% 가 몰립니다. 안전 게이트가 16%.\n'
          '참고 — 사람 상담사는 대화당 약 $6, 고객지원 챗봇은 약 $0.5로 봅니다.'),
         ('속도 — 일반 턴 1.3초 · 카드 턴 6.8초',
          '유해 게이트(1.0초)와 본문(1.2초)을 나란히 돌려 1.3초.\n'
          '카드는 검색·고르기가 순차라 6.8초. Pinecone 연결을 미리 데워 11.4→4.5초.'),
         ('★ 100턴이 가도 입력이 «평평»합니다',
          '1턴 8,131 토큰 → 6턴 8,219 토큰. 5턴 동안 88토큰만 늘었습니다.\n'
          '칸으로 압축하니까요. 대화 전체를 매 턴 보내면 턴에 비례해 커지고 큰 모델이 필요합니다.\n'
          '요즘 LLM 추천 챗봇은 칸을 안 씁니다. 저희는 «100턴을 전제»해서 유지했습니다.'),
         ('검증 — 골든셋 34/34 · 페르소나 11인',
          '골든셋은 «한 턴»만 봅니다. 그래서 사용자 11명을 만들어 LLM 이 연기하게 했습니다.\n'
          '골든셋이 34/34인 상태에서 페르소나가 «6건»을 더 잡았습니다.\n'
          '「전부 실패」가 나오면 대상이 아니라 «자»를 먼저 의심합니다 — 하루에 아홉 번 겪었습니다.')]))

    return made


def main():
    if not SRC or not OUT:
        print('쓰는 법: python -m app.itda.scripts.build_ppt_itda7 <원본.pptx> <출력.pptx>')
        return
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    made = build(prs)
    print(f'■ 원본 {n0}장 · 새로 만든 잇다 {len(made)}장')

    #  새 슬라이드는 «맨 뒤»에 붙었다. 원본의 잇다 구간(18~27)을 빼고 그 자리에 끼운다.
    new_idx = list(range(n0, n0 + len(made)))
    order = ([i for i in range(ITDA_FROM - 1)]          # 1 ~ 17
             + new_idx                                  # 새 잇다 7장
             + [i for i in range(ITDA_TO, n0)])         # 28 ~ 끝
    reorder(prs, order)
    prs.save(OUT)
    print(f'■ 저장 {OUT}  ·  총 {len(order)}장 (전 {n0}장)')
    print('■ ⚠ 반드시 verify_ppt.py 로 열어서 확인할 것')


if __name__ == '__main__':
    main()
