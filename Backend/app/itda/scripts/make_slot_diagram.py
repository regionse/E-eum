# -*- coding: utf-8 -*-
"""「한 문장 → 7개의 칸」 구조 그림을 만든다. (2026-08-07)

사용자 요청
  「PNG. 구조 그림을 하나 만들어주는 게 좋을 듯.
   붙여서 7축으로 보여주게, 어떻게 사용자의 말을 잘 커버할 수 있는지」

왜 이 그림이 필요한가
  지금 덱은 「7축 — 6개는 LLM 이, 1개는 코드가」라고 «이름만» 말하고 넘어간다.
  듣는 사람은 **7축이 뭘 덮는지**를 모른 채 다음 장으로 간다. 「갑자기 튀어나온다」가 이것이다.
  ⇒ 실제 발화 한 문장을 놓고, 그 문장의 «어느 조각»이 «어느 칸»에 들어가는지 보여준다.

설계 원칙
  ① 빈 칸을 «숨기지 않는다». 세부관심·강점성향은 이 발화에 근거가 없어서 비어 있다.
     그게 사양이다(required 아님 — 모르면 안 채운다). 빈 칸이야말로 설명거리다.
  ② 칸마다 «근거»(원문 인용)를 같이 보여준다. 그게 이 설계의 핵심이라서다.
  ③ 코드가 채우는 칸(대상세부)은 «색을 다르게» — LLM 이 한 일과 코드가 한 일을 안 섞는다.

⚠ 값·enum 은 itda_core.py 의 PROFILE_SCHEMA 에서 그대로 가져왔다. 지어내지 않았다.
   활동유형 9종 · 다루는대상 6종 · 제약 5종 · MULTI_MAX=4.

쓰는 법
  python make_slot_diagram.py [바탕.pptx] [출력.pptx]
  → 맨 뒤에 그림 슬라이드를 붙인다. PNG 는 export_slides.py 나 PowerPoint 로 뽑는다.
"""
import sys
import io
import copy

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                           # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v5.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v6.pptx'

IN = 914400

# ── 덱에서 뽑은 디자인 값 (슬라이드20 실측) ──────────────────────────
FONT      = 'Arial'
C_HEAD    = RGBColor(0x8B, 0x84, 0xA0)     # 헤더 회보라
C_PRESENT = RGBColor(0x6B, 0x3F, 0xA0)     # 발표자 / 강조 값
C_TITLE   = RGBColor(0x3B, 0x20, 0x63)     # 제목 · 카드 제목
C_BODY    = RGBColor(0x55, 0x50, 0x6B)     # 본문
C_CARD    = RGBColor(0xFF, 0xFF, 0xFF)     # 카드 배경
C_CARD_LN = RGBColor(0xE7, 0xE2, 0xF0)     # 카드 테두리
C_PANEL   = RGBColor(0xEF, 0xEB, 0xF6)     # 강조 패널 배경
C_PANEL_LN = RGBColor(0xDE, 0xD8, 0xEA)    # 강조 패널 테두리
C_MUTED   = RGBColor(0x8B, 0x84, 0xA0)     # 근거 인용 (흐리게)


def blank_like(prs):
    """덱의 마지막 슬라이드 레이아웃으로 «빈» 슬라이드를 만든다."""
    lay = prs.slides[0].slide_layout
    s = prs.slides.add_slide(lay)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    """runs = [(글자, pt, bold, color), ...] — 한 줄에 하나씩 문단으로 쌓는다."""
    tb = slide.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                                  Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.color.rgb = col
        p.line_spacing = 1.25
    return tb


def box(slide, x, y, w, h, fill, line, radius=0.02):
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Emu(int(x * IN)), Emu(int(y * IN)),
                                Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    try:                       # 모서리 둥글기 — 값이 작을수록 각지다
        sh.adjustments[0] = radius
    except Exception:
        pass
    if sh.has_text_frame:
        sh.text_frame.text = ''
    return sh


# ════════════════════════════════════════════════════════════════
#  내용 — 실제 발화 한 문장과, 그 문장이 어느 칸에 들어가는지
#  ⚠ 값은 PROFILE_SCHEMA 의 enum 에 «실제로 있는» 것만 쓴다.
# ════════════════════════════════════════════════════════════════
#  ⚠ 처음엔 «엄마»로 썼다가 실제로 돌려보고 고쳤다.
#    fill_obj_detail({}, '엄마 돌보느라…')  → None      ← '엄마' 는 인식 목록에 «없다»
#    fill_obj_detail({}, '할머니 돌보느라…') → '어르신'   ← 이건 실제로 채워진다
#    _OBJ_DETAIL_MARK 의 '어르신' 목록: 어르신·노인·할머니·할아버지·조부모·치매·요양·간병…
#  그림이 「코드가 채운다」고 주장하려면 «실제로 채워지는» 문장이어야 한다.
발화 = ('“할머니 돌보느라 학교를 그만뒀어요.\n'
        '   손으로 만드는 건 재밌었는데, 뭘 해야 할지 모르겠어요”')

#  (칸이름, 부가설명, 값, 근거인용, 코드가채우나)
SLOTS = [
    ('관심분야',  '일상어 그대로 · 최대 4개',  '만들기',
     '“손으로 만드는 건 재밌었”', False),
    ('활동유형',  '9종에서 고름 · 최대 4개',   '돕기·돌봄  /  만들기',
     '“돌보느라”   ·   “만드는”', False),
    ('다루는대상', '6종에서 고름 · 최대 4개',   '사람',
     '“할머니 돌보느라”', False),
    ('제약',      '5종에서 고름 · 최대 4개',   '학력부담',
     '“학교를 그만뒀어요”', False),
    ('세부관심',  '단일값',                   '(안 채움)',
     '이 발화엔 근거가 없다', False),
    ('강점성향',  '단일값',                   '(안 채움)',
     '이 발화엔 근거가 없다', False),
    #  ⚠ 값 6종은 job_attr.tsv 실측: 고객·손님 64 / 직원·동료 13 / 환자·장애인 11
    #    / 학습자 7 / 아동·청소년 6 / 어르신 1.  «가족·어르신» 같은 값은 «없다».
    ('대상세부',  '★ LLM 이 아니라 코드가',    '어르신',
     '‘할머니’ 를 읽고 코드가 보완', True),
]


def build(prs):
    s = blank_like(prs)

    # ── 헤더 ──────────────────────────────────────────────
    text(s, 1.00, 0.92, 6.0, 0.30,
         [('PART 03 · 잇다 — 슬롯', 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30,
         [('발표자 정준', 18.75, False, C_PRESENT)])
    text(s, 1.00, 1.50, 19.0, 0.95,
         [('한 문장을 «7개의 칸»으로 바꿉니다', 48, True, C_TITLE)])

    # ── ① 사용자 발화 ─────────────────────────────────────
    box(s, 1.00, 2.62, 12.6, 1.62, C_CARD, C_CARD_LN)
    text(s, 1.45, 2.86, 11.8, 1.2,
         [(발화, 25, True, C_TITLE)])

    #  오른쪽: 왜 칸으로 바꾸나
    box(s, 13.95, 2.62, 5.05, 1.62, C_PANEL, C_PANEL_LN)
    text(s, 14.35, 2.84, 4.3, 1.2,
         [('문장 그대로는 못 찾습니다', 19, True, C_TITLE),
          ('검색은 “뭘 해야 할지” 를 못 읽습니다.', 15.5, False, C_BODY),
          ('그래서 «칸»으로 바꿉니다.', 15.5, False, C_BODY)])

    # ── ② 화살표 띠 ───────────────────────────────────────
    ar = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,
                            Emu(int(6.95 * IN)), Emu(int(4.36 * IN)),
                            Emu(int(0.42 * IN)), Emu(int(0.46 * IN)))
    ar.fill.solid()
    ar.fill.fore_color.rgb = C_PANEL_LN
    ar.line.fill.background()
    ar.shadow.inherit = False
    text(s, 7.60, 4.44, 11.0, 0.34,
         [('LLM 한 번 —  7칸 + 답변 + 검색어를 «한꺼번에» 받습니다', 19, True, C_PRESENT)])

    # ── ③ 7개의 칸 (4열 × 2행) ────────────────────────────
    #  ⚠ 폭 계산 — 슬라이드 20" · 좌우 여백 1" → 쓸 수 있는 폭 18.0"
    #    4칸 + 3틈 = 4×4.16 + 3×0.45 = 17.99  (처음에 4.35/0.55 로 잡았다가
    #    오른쪽 끝이 20.05" 가 되어 «슬라이드 밖»으로 0.05" 넘어갔다)
    X0, Y0, BW, BH, GX, GY = 1.00, 5.02, 4.16, 2.52, 0.45, 0.30
    for i, (name, sub, val, ev, by_code) in enumerate(SLOTS):
        col, row = i % 4, i // 4
        x = X0 + col * (BW + GX)
        y = Y0 + row * (BH + GY)
        fill = C_PANEL if by_code else C_CARD
        line = C_PANEL_LN if by_code else C_CARD_LN
        empty = val.startswith('(')
        box(s, x, y, BW, BH, fill, line)
        text(s, x + 0.34, y + 0.30, BW - 0.68, 0.42,
             [(name, 21, True, C_TITLE)])
        text(s, x + 0.34, y + 0.78, BW - 0.68, 0.30,
             [(sub, 14, False, C_MUTED)])
        text(s, x + 0.34, y + 1.22, BW - 0.68, 0.50,
             [(val, 20, True, C_MUTED if empty else C_PRESENT)])
        text(s, x + 0.34, y + 1.80, BW - 0.68, 0.62,
             [(ev, 15, False, C_MUTED)])

    # ── ④ 마지막 칸 자리에 «근거 대조» 설명 ────────────────
    x = X0 + 3 * (BW + GX)
    y = Y0 + 1 * (BH + GY)
    box(s, x, y, BW, BH, C_PANEL, C_PANEL_LN)
    text(s, x + 0.34, y + 0.30, BW - 0.68, 0.42,
         [('값마다 «근거»를 받습니다', 21, True, C_TITLE)])
    text(s, x + 0.34, y + 0.84, BW - 0.68, 1.5,
         [('근거가 원문에 없으면', 16, False, C_BODY),
          ('그 칸을 «버립니다».', 16, True, C_TITLE),
          ('', 8, False, C_BODY),
          ('지어내려면 근거도 지어내야 하는데,', 15, False, C_MUTED),
          ('그건 원문 대조에서 걸립니다.', 15, False, C_MUTED)])

    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    build(prs)
    prs.save(OUT)
    print(f'  {n0}장 → {len(Presentation(OUT).slides)}장 (그림 1장 추가)')
    print(f'  저장 → {OUT}')
    print(f'  그림은 «맨 마지막» 슬라이드다.')


main()
