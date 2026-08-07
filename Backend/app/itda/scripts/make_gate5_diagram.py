# -*- coding: utf-8 -*-
"""관문 그림 — **코드 순서대로** 다시. (2026-08-07)

앞 판(make_gate4_diagram)이 뭘 틀렸나
  역할별로 묶느라 **①을 세 번째에 그렸다.** 그런데 코드에서 ①은 «맨 앞»이다.
  흐름도인데 순서를 바꿔 그렸으니 그냥 틀린 그림이다.

  그리고 순서를 고치다 «더 큰 것»을 알았다 —
    ① 은 이번 턴에 센 값을 보는 게 «아니다».
    bump_abuse 는 ①의 «아래»(④⑤⑦)에서 일어난다. 즉 ①이 읽는 _abuse 는
    **지난 턴들에 쌓인 값**이다. 그래서 ①이 맨 앞에 있어도 모순이 없다.
      이번 턴  ④⑤⑦ → bump_abuse → _abuse 가 «세션에» 쌓인다
      다음 턴  ①    → 그 값을 읽고 잠근다
  앞 판의 화살표는 «같은 턴 안»에서 도는 것처럼 그려져 있었다. 사실이 아니다.

이 판의 규칙
  · 세로 줄기는 **코드 실행 순서** ①~⑦ 그대로.
  · 역할(막는다·센다·기억한다·가른다)은 «색과 꼬리표»로만 표시한다.
  · 「센다 → 다음 턴의 ①」 화살표를 «턴을 건너뛰는» 것으로 그린다.

코드 근거 (itda_core.py)
  5343 ① return blocked          ← 유일하게 «끊는» 문 (이전 턴들의 _abuse 를 읽음)
  5371 ② _exclude 담기            return 없음
  5408 ③ _policy_declined=True    return 없음
  5428 ④ bump_abuse('불법신호')
  5430 ⑤ pre_check → 네 갈래 · 5470 harm · 5489 unsafe 도 bump_abuse
  5509 ⑥ pick_from_options        return 없음
  5560 ⑦ return redirect · 5561 bump_abuse('injection')
"""
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR            # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v12.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v13.pptx'

IN = 914400
FONT = 'Arial'
C_HEAD  = RGBColor(0x8B, 0x84, 0xA0)
C_PRES  = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_BODY  = RGBColor(0x55, 0x50, 0x6B)
C_MUTED = RGBColor(0x8B, 0x84, 0xA0)
C_CARD  = RGBColor(0xFF, 0xFF, 0xFF)
C_CARDL = RGBColor(0xE7, 0xE2, 0xF0)
C_LINE  = RGBColor(0xC9, 0xC1, 0xDC)

BLOCK = RGBColor(0xA8, 0x36, 0x4A)   # 막는다
COUNT = RGBColor(0xC2, 0x6B, 0x2E)   # 센다
MEMO  = RGBColor(0x4E, 0x7D, 0x5E)   # 기억한다
FORK  = RGBColor(0x6B, 0x3F, 0xA0)   # 가른다
BG = {BLOCK: RGBColor(0xF9, 0xEC, 0xEF), COUNT: RGBColor(0xFB, 0xF0, 0xE6),
      MEMO: RGBColor(0xEC, 0xF3, 0xEE), FORK: RGBColor(0xEF, 0xEB, 0xF6)}


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.2):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.color.rgb = col
    return tb


def box(sl, x, y, w, h, fill, line, lw=1.0):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False; sh.text_frame.text = ''
    try:
        sh.adjustments[0] = 0.02
    except Exception:
        pass
    return sh


def seg(sl, x1, y1, x2, y2, col=C_LINE, w=2.0, dash=False):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Emu(int(x1 * IN)), Emu(int(y1 * IN)),
                                Emu(int(x2 * IN)), Emu(int(y2 * IN)))
    c.line.color.rgb = col; c.line.width = Pt(w)
    if dash:
        from pptx.enum.dml import MSO_LINE_DASH_STYLE
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return c


#  (번호, 이름, 역할, 색, 결과)   ← 코드 실행 순서 그대로
GATES = [
    ('①', '이미 잠긴 사람인가', '막는다',   BLOCK, '차단 ⛔  — 여기서 대화가 «끝난다»'),
    ('②', '「그거 말고」인가',   '기억한다', MEMO,  '그 후보를 빼고 다시 찾는다'),
    ('③', '「괜찮아요」인가',    '기억한다', MEMO,  '정책 안내를 멈춘다'),
    ('④', '불법 신호인가',      '센다',    COUNT, '«막지 않는다». _abuse 만 +1'),
    ('⑤', 'pre_check',        '가른다',  FORK,  '통과 · 되묻기 · ★잇기 · 차단'),
    ('⑥', '「1번이요」인가',     '기억한다', MEMO,  '고른 값을 받아 카드로'),
    ('⑦', '공격인가',          '센다·막는다', COUNT, '되돌린다 ⛔  +  _abuse +1'),
]

WHY = [
    ('네 가지 «다른 일»이 섞여 있습니다', [
        '이름은 다 「관문」인데 하는 일이 다릅니다.',
        '막는다 · 센다 · 기억한다 · 가른다 — 색으로 갈라 뒀습니다.',
    ]),
    ('실제로 «끊는» 건 ①⑤⑦ 뿐입니다', [
        '코드에서 return 으로 대화를 끝내는 건 셋입니다.',
        '②③⑥ 은 상태만 바꾸고 그냥 지나갑니다.',
    ]),
    ('★ ①은 «지난 턴»의 값을 읽습니다', [
        'bump_abuse 는 ① «아래»(④⑤⑦)에서 일어납니다.',
        '그래서 ①이 이번 턴에 센 값을 볼 수는 없습니다 —',
        '지난 턴들에 세션에 쌓인 _abuse 를 읽는 겁니다.',
        '⇒ 이번 턴에 센 것은 «다음 턴»의 ①이 씁니다.',
    ]),
    ('④를 빼면 ①이 «영원히» 안 걸립니다', [
        '레드팀 15턴에 카운터가 0 이었던 게 그 상태입니다.',
        'LLM 이 막긴 했는데 우리는 몰랐습니다.',
        '낱말로 막으면 과차단하니 «세기만» 합니다.',
    ]),
]


def build(prs):
    s = blank_like(prs)
    text(s, 1.00, 0.92, 9.0, 0.30, [('PART 03 · 잇다 — 예외 처리', 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRES)])
    text(s, 1.00, 1.50, 19.0, 0.95,
         [('이름은 다 «관문»인데, 하는 일이 넷입니다', 42, True, C_TITLE)])

    # ── 왼쪽 : 왜 ───────────────────────────────────────
    LX, LW = 1.00, 7.15
    y = 2.66
    for head, lines in WHY:
        text(s, LX, y, LW, 0.32, [(head, 16.5, True, C_TITLE)])
        y += 0.36
        for ln in lines:
            text(s, LX + 0.24, y, LW - 0.24, 0.28, [(ln, 13, False, C_BODY)])
            y += 0.27
        y += 0.22

    #  범례
    yl = y + 0.06
    text(s, LX, yl, 2.0, 0.28, [('색 =', 13.5, True, C_MUTED)])
    xx = LX + 0.70
    for nm, col in (('막는다', BLOCK), ('센다', COUNT), ('기억한다', MEMO), ('가른다', FORK)):
        box(s, xx, yl + 0.03, 0.20, 0.20, col, col)
        text(s, xx + 0.28, yl, 1.30, 0.28, [(nm, 13.5, True, col)])
        xx += 1.55

    # ── 오른쪽 : 코드 순서 ①~⑦ ──────────────────────────
    RX, RW = 8.75, 10.25
    SPINE = RX + 0.55
    box(s, RX + 1.90, 2.66, 4.40, 0.50, C_CARD, C_CARDL, 0.75)
    text(s, RX + 1.90, 2.79, 4.40, 0.28,
         [('🧑  사용자 발화 한 마디', 14.5, True, C_TITLE)], align=PP_ALIGN.CENTER)

    BH, GAP = 0.76, 0.22
    ys = [3.46 + i * (BH + GAP) for i in range(7)]
    seg(s, SPINE, 3.16, SPINE, ys[-1] + BH / 2)

    for (mark, name, role, col, res), yy in zip(GATES, ys):
        seg(s, SPINE, yy + BH / 2, RX + 1.10, yy + BH / 2, col, 2.2)
        box(s, RX + 1.10, yy, RW - 1.10, BH, BG[col], col, 1.1)
        text(s, RX + 1.42, yy + 0.22, 0.42, 0.30, [(mark, 16, True, col)])
        text(s, RX + 1.90, yy + 0.22, 3.10, 0.30, [(name, 15, True, C_TITLE)])
        text(s, RX + 5.05, yy + 0.24, 1.30, 0.26, [(role, 12.5, True, col)])
        text(s, RX + 6.45, yy + 0.24, RW - 6.75, 0.26, [(res, 13, False, C_BODY)])

    #  ★ 「센다」 → «다음 턴»의 ① — 턴을 건너뛰는 화살표라 점선으로
    ax = RX + RW + 0.10
    for i in (3, 4, 6):                       # ④⑤⑦
        seg(s, ax, ys[i] + BH / 2, ax + 0.30, ys[i] + BH / 2, COUNT, 1.75)
    seg(s, ax + 0.30, ys[3] + BH / 2, ax + 0.30, ys[6] + BH / 2, COUNT, 1.75)
    seg(s, ax + 0.30, ys[0] + BH / 2, ax + 0.30, ys[3] + BH / 2, COUNT, 1.75, dash=True)
    seg(s, ax + 0.30, ys[0] + BH / 2, ax, ys[0] + BH / 2, COUNT, 1.75, dash=True)
    #  ⚠ 여기 「_abuse 가 세션에 쌓여 다음 턴의 ①이 읽는다」 라벨을 뒀다가
    #    ②번 상자를 «덮었다». 오른쪽엔 0.6" 밖에 안 남아 라벨이 들어갈 자리가 없다.
    #    왼쪽 글의 「★ ①은 지난 턴의 값을 읽습니다」가 같은 말을 하므로 그냥 뺀다.
    #    점선 + 주황색만으로 「턴을 건너뛴다」가 보인다.

    #  끝
    ye = ys[-1] + BH + 0.30
    seg(s, SPINE, ys[-1] + BH / 2, SPINE, ye + 0.25)
    box(s, RX + 1.90, ye, 4.40, 0.50, BG[FORK], FORK, 1.25)
    text(s, RX + 1.90, ye + 0.13, 4.40, 0.28,
         [('통과한 것만 LLM 으로', 14.5, True, C_TITLE)], align=PP_ALIGN.CENTER)
    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    build(prs)
    prs.save(OUT)
    W, H = prs.slide_width / IN, prs.slide_height / IN
    over = [round((sh.left + sh.width) / IN, 2)
            for sh in Presentation(OUT).slides[-1].shapes
            if (sh.left + sh.width) / IN > W + 0.01 or (sh.top + sh.height) / IN > H + 0.01]
    print(f'  {n0}장 → {len(Presentation(OUT).slides)}장')
    print(f'  슬라이드 밖: {len(over)}개  {"✅" if not over else over}')
    print(f'  저장 → {OUT}')


main()
