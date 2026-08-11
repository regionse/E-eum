# -*- coding: utf-8 -*-
"""잇다 발표 — 「LLM 은 읽고, 코드는 정한다」 2장. (2026-08-10)

왜 만드나
  사용자 지적: 「이거 중요한 기술 같은데 서술이 하나도 없다」.
  맞다 — 덱 어디에도 «누가 무엇을 결정하는가»가 없었다. 각 장이 기능만 보여주고
  그 기능의 «주인»을 안 밝혔다. 이 프로젝트의 뼈대가 정확히 그 분업인데도.

⚠ **자기완결로 만든다.** build_ppt_itda_main.py 를 다른 세션이 편집 중이라
  같은 파일을 건드리면 서로 덮어쓴다(오늘 실제로 한 번 당했다 — GATE_MODEL).
  ⇒ 별도 pptx 로 뽑고, PowerPoint 「슬라이드 재사용」으로 끼운다. 슬라이드 복제 0회.

숫자 출처 (전부 오늘 실측)
  · 코드가 모델의 action 을 덮어쓰는 자리 8곳 — itda_core.py grep
  · 착지 문턱 2.0 · TOP% 0.90 · 엔트로피 0.70 — ItdaEngine 상수 직접 출력
  · 종류 판정 원함 97% · 해봤음 96% — 79케이스 재측정(2026-08-10)
  · Zero-Shot Belief (arXiv:2502.08777, 2025) — 「LLM 단독보다 섞은 쪽이 나았다」

쓰는 법 (Backend/ 에서)
  python app/itda/scripts/build_ppt_itda_split.py
  → C:/Users/TFX255GS/Downloads/이음_잇다_분업.pptx   이후 verify_ppt.py 로 «열어» 확인
"""
import io
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r'C:\Users\TFX255GS\Downloads\이음_잇다_분업.pptx'

#  ── 팔레트 · 좌표 (수정안(4) 덜다 구간 실측 — 본편과 «같은 값») ──
C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_EYE = RGBColor(0x8B, 0x84, 0xA0)
C_ACC = RGBColor(0x6B, 0x3F, 0xA0)
C_TXT = RGBColor(0x24, 0x1F, 0x2E)
C_RED = RGBColor(0xA8, 0x36, 0x4A)
C_DEEP = RGBColor(0x32, 0x16, 0x6F)
C_PANEL_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_BG = RGBColor(0xEF, 0xEB, 0xF6)
C_BG_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_CHIP = RGBColor(0xF0, 0xE9, 0xFB)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LINE = RGBColor(0xC9, 0xBE, 0xE0)
FONT = 'Malgun Gothic'


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _font(r, size, color, bold):
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.18, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        for (t, size, color, bold) in line:
            r = p.add_run()
            r.text = t
            _font(r, size, color, bold)
    return tb


def box(sl, x, y, w, h, fill=C_WHITE, line=C_PANEL_LN, lw=0.75):
    shp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.06
    except Exception:                                   # noqa: BLE001
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def arrow(sl, x1, y1, x2, y2):
    ln = sl.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = C_LINE
    ln.line.width = Pt(2.25)
    el = ln.line._get_or_add_ln()
    el.append(el.makeelement(qn('a:tailEnd'),
                             {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return ln


def head(sl, crumb, title):
    text(sl, 1.00, 0.86, 12.0, 0.45, [[(crumb, 18.75, C_EYE, False)]])
    text(sl, 14.0, 0.80, 5.0, 0.55, [[('발표자 김정준', 30, C_ACC, True)]],
         align=PP_ALIGN.RIGHT)
    text(sl, 1.00, 1.42, 18.1, 0.95, [[(title, 48, C_TITLE, True)]])


# ══════════════════════════════════════════════════════════════════
def s_principle(prs):
    """1장 — 원칙과 «실제로 뒤집는» 예."""
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다 — 분업', 'LLM 은 «읽고», 코드는 «정한다»')

    #  왼쪽 — 두 역할
    box(sl, 1.0, 2.65, 8.6, 4.15)
    text(sl, 1.45, 2.92, 7.8, 0.5, [[('LLM 이 하는 일', 30, C_ACC, True),
                                     ('   — 무슨 «뜻»인가', 18, C_EYE, False)]])
    text(sl, 1.45, 3.62, 7.8, 3.0, [
        [('말에서 칸을 뽑는다', 20, C_TXT, True),
         ('   값 · 종류 · 근거', 17, C_EYE, False)],
        [('위험·거부·착지요청을 읽는다', 20, C_TXT, True)],
        [('후보 중 하나를 고르고 «왜»를 쓴다', 20, C_TXT, True)],
        [('긴 대화를 접는다', 20, C_TXT, True)],
        [('공감·되묻기 — 사람 말투 전부', 20, C_TXT, True)],
    ], ls=1.42)

    box(sl, 10.5, 2.65, 8.6, 4.15)
    text(sl, 10.95, 2.92, 7.8, 0.5, [[('코드가 하는 일', 30, C_DEEP, True),
                                      ('   — 그래서 «무엇을 할 것인가»', 18, C_EYE, False)]])
    text(sl, 10.95, 3.62, 7.8, 3.0, [
        [('근거를 발화와 대조한다', 20, C_TXT, True),
         ('   없으면 버린다', 17, C_EYE, False)],
        [('언제 결과를 낼지 정한다', 20, C_TXT, True),
         ('   점수 분포로', 17, C_EYE, False)],
        [('검색하고 순위를 합친다', 20, C_TXT, True), ('   RRF', 17, C_EYE, False)],
        [('사실 문장을 «직접» 쓴다', 20, C_TXT, True),
         ('   시험일 · 「없어요」', 17, C_EYE, False)],
        [('나가는 말을 문장 단위로 거른다', 20, C_TXT, True)],
    ], ls=1.42)

    #  아래 — 실제로 뒤집는 자리
    #  ⚠ 패널 높이 3.7 안에 제목 + 4줄 + 맺음 이 들어가야 한다.
    #    처음 판은 줄 간격 0.62 라 맺음줄이 패널 밖으로 나갔다(PNG 검수에서 잡음).
    box(sl, 1.0, 7.05, 18.1, 3.55, fill=C_BG, line=C_BG_LN)
    text(sl, 1.5, 7.22, 17.1, 0.5,
         [[('★ 모델이 「이제 추천하자」고 해도 — 코드가 «8곳»에서 뒤집는다', 26, C_TITLE, True)]])

    y = 7.92
    for a, b in [('🤖 모델', 'SEARCH — 이제 찾아도 되겠다'),
                 ('⚙ 코드', '축이 2개뿐이다(무게 2.0 미달) → 되묻기'),
                 ('⚙ 코드', '1위 동네가 흐리다(TOP% < 0.90) → 방향 칩'),
                 ('⚙ 코드', '동네 안이 갈린다(엔트로피 ≥ 0.70) → 세부 칩')]:
        col = C_ACC if '모델' in a else C_DEEP
        text(sl, 1.6, y, 2.2, 0.40, [[(a, 19, col, True)]])
        text(sl, 4.1, y, 14.5, 0.40, [[(b, 19, C_TXT, False)]])
        y += 0.54
    text(sl, 1.6, y + 0.06, 17.0, 0.5,
         [[('모델의 「충분히 알았다」를 안 믿는다 — «숫자»로 본다', 19, C_RED, True)]])


def s_map(prs):
    """2장 — 한 턴의 다섯 단계에 주인을 붙인다."""
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다 — 분업', '한 턴에서 주인이 다섯 번 바뀝니다')

    rows = [
        ('①', '게이트', '⚙ 코드', '낱말·표기 정규화 — 0원 급행',
         '🤖 LLM', '「누구 이야기인가」 판정'),
        ('②', '슬롯 뽑기', '🤖 LLM', '값 · 종류(원함/해봤음/못함) · 근거',
         '⚙ 코드', '근거가 발화에 «없으면 버린다»'),
        ('③', '검색', '🤖 LLM', '검색어를 만든다',
         '⚙ 코드', '두 갈래로 찾고 «등수»로 합친다(RRF)'),
        ('④', '착지 판정', '⚙ 코드', '무게 2.0 · TOP% 0.90 · 엔트로피 0.70',
         '—', 'LLM 은 여기 «관여하지 않는다»'),
        ('⑤', '답 내기', '🤖 LLM', '후보 고르기 · 문장 쓰기',
         '⚙ 코드', '사실 문장 작성 + 출력 가드 4체'),
    ]
    #  ⚠ 다섯 줄(1.32×5 + 간격) + 맺음 패널이 «11.25 안»에 들어가야 한다.
    #    처음 판은 1.55 간격이라 맺음 두 줄째가 화면 밖으로 나갔다(PNG 검수에서 잡음).
    y = 2.58
    for num, name, o1, d1, o2, d2 in rows:
        box(sl, 1.0, y, 18.1, 1.30)
        text(sl, 1.35, y + 0.38, 0.7, 0.52, [[(num, 25, C_ACC, True)]])
        text(sl, 2.15, y + 0.38, 2.5, 0.52, [[(name, 23, C_TITLE, True)]])
        c1 = C_ACC if 'LLM' in o1 else C_DEEP
        text(sl, 5.0, y + 0.14, 2.1, 0.40, [[(o1, 18, c1, True)]])
        text(sl, 5.0, y + 0.60, 6.6, 0.48, [[(d1, 17, C_TXT, False)]])
        c2 = C_ACC if 'LLM' in o2 else (C_EYE if o2 == '—' else C_DEEP)
        text(sl, 12.0, y + 0.14, 2.1, 0.40, [[(o2, 18, c2, True)]])
        text(sl, 12.0, y + 0.60, 6.7, 0.48, [[(d2, 17, C_TXT, False)]])
        arrow(sl, 11.35, y + 0.65, 11.85, y + 0.65)
        y += 1.42

    box(sl, 1.0, y + 0.02, 18.1, 1.42, fill=C_BG, line=C_BG_LN)
    text(sl, 1.5, y + 0.22, 17.1, 1.1, [
        [('왜 이렇게 나눴나 — ', 20, C_DEEP, True),
         ('모델은 «지시를 흘립니다».', 20, C_TXT, True)],
        [('12,615자 프롬프트가 금지한 질문을 72번 중 2번 했습니다. '
          '그래서 지키는 자리는 코드로 내렸습니다.', 18, C_EYE, False)],
    ], ls=1.3)


def main():
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    s_principle(prs)
    s_map(prs)
    prs.save(OUT)
    print(f'저장: {OUT} · 2장')


if __name__ == '__main__':
    main()
