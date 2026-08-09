# -*- coding: utf-8 -*-
"""데이터 출처 — 「어디서 무엇을 가져와 어떻게 바꿨나」. (2026-08-07)

사용자 요청
  「무슨 API에서 무슨 파라미터를 가져왔고, 그 파라미터가 어떤 상태여서 어떻게 정규화시켰고,
   우리에게 가져왔는지까지」 — 덜다에 있는 그 장이 잇다엔 «없다».
  「이건 맨 처음에 넣어야 함」

⚠ 값은 전부 적재 스크립트 원문에서 뽑았다. 지어낸 필드명이 없다.
  · load_certification.py:38  openapi.q-net.or.kr/.../getList
        jmcd·jmfldnm·seriesnm·obligfldnm·mdobligfldnm·qualgbnm → certification
  · load_exam_schedule.py:45  apis.data.go.kr/B490007/qualExamSchd/getQualExamSchdList
        implYy·implSeq·docRegStartDt…pracPassDt → exam_schedule (to_date 로 DATE 정규화)
  · load_course.py:37-38      apis.data.go.kr/B552881/kmooc_v2_0/courseList_v2_0 · courseDetail_v2_0
        id·name·classfy_name·middle_classfy_name·summary(HTML)·url → course
  · job_catalog              NCS_SUBD_CDNM·DUTY_DEF (컬럼 주석에 원본 필드명이 남아 있다)

⚠ DB 건수는 2026-08-07 RDS 실측: 직업 1,094 · 자격증 613 · 시험일정 2,655 · 강좌 8,273
"""
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR            # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v14.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v15.pptx'

IN = 914400
FONT, CODE = 'Arial', 'Consolas'
C_HEAD  = RGBColor(0x8B, 0x84, 0xA0)
C_PRES  = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_BODY  = RGBColor(0x55, 0x50, 0x6B)
C_MUTED = RGBColor(0x8B, 0x84, 0xA0)
C_CARD  = RGBColor(0xFF, 0xFF, 0xFF)
C_CARDL = RGBColor(0xE7, 0xE2, 0xF0)
C_PANEL = RGBColor(0xEF, 0xEB, 0xF6)
C_PANLL = RGBColor(0xDE, 0xD8, 0xEA)
C_CODEBG = RGBColor(0xF7, 0xF5, 0xFB)
C_ORANGE = RGBColor(0xC2, 0x6B, 0x2E)
C_GREEN = RGBColor(0x4E, 0x7D, 0x5E)


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.2, font=FONT):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls
        r = p.add_run(); r.text = t
        r.font.name = font; r.font.size = Pt(pt)
        r.font.bold = bold; r.font.color.rgb = col
    return tb


def box(sl, x, y, w, h, fill, line, lw=1.0):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False; sh.text_frame.text = ''
    try:
        sh.adjustments[0] = 0.02
    except Exception:
        pass
    return sh


def arrow(sl, x, y, w=0.30, h=0.24, col=C_PANLL):
    sh = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid(); sh.fill.fore_color.rgb = col
    sh.line.fill.background(); sh.shadow.inherit = False
    sh.text_frame.text = ''
    return sh


#  (출처, 기관, 형식, 가져온 것, 우리 테이블, 건수, 바꾼 것)
#  ⚠ 원본 필드명(jmcd·docRegStartDt…)은 «쓰지 않는다». 사용자 지시 —
#    「어떤 파라미터를 가져왔는지는 이름만. 강좌 요약, 이름 이런 식으로만」
#    발표에서 영문 필드명은 아무도 못 읽는다. 무엇을 받았는지만 보이면 된다.
ROWS = [
    ('NCS 직무',       '국가직무능력표준',
     'ncs.go.kr',
     '직무 코드 · 직무명 · 직무 정의 · 대/중/소분류',
     'job_catalog', '1,094',
     '직무 정의를 카드 설명에 «그대로» 씁니다'),
    ('자격증',          '한국산업인력공단 Q-Net',
     'XML',
     '종목 코드 · 종목명 · 등급 · 직무 분야 · 자격 구분',
     'certification', '613',
     '종목 코드를 키로 «덮어쓰기» — 다시 돌려도 안 쌓입니다'),
    ('시험일정',        '공공데이터포털',
     'JSON',
     '연도 · 회차 · 접수일 · 시험일 · 발표일',
     'exam_schedule', '2,655',
     '문자열로 오는 날짜를 «날짜 형식»으로 통일합니다'),
    ('K-MOOC 강좌',    '공공데이터포털',
     'JSON · 목록 + 상세',
     '강좌 이름 · 분류 · 요약 · 링크',
     'course', '8,273',
     '요약에 HTML 태그가 섞여 있어 «걷어냅니다»'),
]

#  ⚠ 한 줄은 «33자 이내». 상자 폭 5.32" · 11.5pt 기준이다.
#    처음에 45자짜리를 넣었더니 두 줄로 넘쳐 «아랫줄을 덮었다»
#    (줄마다 텍스트박스를 따로 놓고 0.30" 씩 내리는 구조라 넘치면 바로 겹친다).
NOTES = [
    ('공개 강좌만 받습니다', C_GREEN, [
        '상세를 17,417개 다 부르면 하루 한도(10,000)를 넘깁니다.',
        'public_yn="Y" 인 공개 강좌만 받아 한도 안에 듭니다.',
        '어차피 추천도 공개 강좌만 합니다.',
    ]),
    ('한도를 «세 겹»으로 막았습니다', C_PRES, [
        '0.15초 대기  ·  9,500회 상한  ·  넣은 건 건너뜀',
        '중간에 끊겨도 다시 돌리면 «이어받습니다».',
        '재실행이 안전해야 배치를 손으로 돌릴 수 있습니다.',
    ]),
    ('오래 못 찾았던 것 — 자격 구분 코드', C_ORANGE, [
        '기술자격 513종은 "T", 전문자격 100종은 "S" 입니다.',
        '"T" 로만 부르면 100종이 «0건»으로 옵니다.',
        '오류가 아니라 «빈 응답»이라 안 보였습니다.',
    ]),
]


def build(prs):
    s = blank_like(prs)
    text(s, 1.00, 0.92, 9.0, 0.30, [('PART 03 · 잇다 — 데이터', 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRES)])
    text(s, 1.00, 1.50, 19.0, 0.95,
         [('데이터는 «전부» 공공데이터입니다', 44, True, C_TITLE)])
    text(s, 1.00, 2.42, 14.0, 0.30,
         [('네 곳에서 받아 우리 테이블 네 개로 —  직접 만든 데이터는 «직업 태그» 하나뿐입니다',
           14.5, False, C_MUTED)])

    #  머리글
    X0 = 1.00
    COLS = [(0.00, 3.30, '어느 API'), (3.50, 6.20, '무엇을 가져와'),
            (10.00, 4.20, '어떻게 바꿔서'), (14.50, 3.50, '어디에 넣나')]
    for dx, w, name in COLS:
        text(s, X0 + dx, 2.90, w, 0.28, [(name, 13, True, C_MUTED)])

    y = 3.24
    RH = 1.34
    for src, org, url, fields, table, n, norm in ROWS:
        box(s, X0, y, 18.00, RH, C_CARD, C_CARDL, 0.9)
        #  ① 어디서
        text(s, X0 + 0.30, y + 0.22, 3.10, 0.30, [(src, 16, True, C_TITLE)])
        text(s, X0 + 0.30, y + 0.56, 3.10, 0.26, [(org, 12, False, C_MUTED)])
        text(s, X0 + 0.30, y + 0.84, 3.10, 0.26, [(url, 11, False, C_PRES)], font=CODE)
        arrow(s, X0 + 3.14, y + RH / 2 - 0.12)
        #  ② 필드
        text(s, X0 + 3.60, y + 0.34, 6.10, 0.60, [(fields, 13.5, False, C_TITLE)])
        arrow(s, X0 + 9.66, y + RH / 2 - 0.12)
        #  ③ 정규화
        text(s, X0 + 10.05, y + 0.34, 4.15, 0.66, [(norm, 12, False, C_BODY)])
        arrow(s, X0 + 14.20, y + RH / 2 - 0.12)
        #  ④ 테이블
        box(s, X0 + 14.55, y + 0.28, 3.20, 0.76, C_PANEL, C_PANLL, 1.0)
        text(s, X0 + 14.55, y + 0.40, 3.20, 0.28,
             [(table, 14.5, True, C_TITLE)], align=PP_ALIGN.CENTER, font=CODE)
        text(s, X0 + 14.55, y + 0.70, 3.20, 0.26,
             [(f'{n} 행', 12.5, False, C_MUTED)], align=PP_ALIGN.CENTER)
        y += RH + 0.16

    #  아래 — 배치를 돌리며 배운 것
    y += 0.10
    for i, (head, col, lines) in enumerate(NOTES):
        bx = X0 + i * 6.06
        box(s, bx, y, 5.88, 1.62, C_CODEBG, C_PANLL, 0.9)
        text(s, bx + 0.28, y + 0.20, 5.32, 0.30, [(head, 14, True, col)])
        yy = y + 0.56
        for ln in lines:
            text(s, bx + 0.28, yy, 5.32, 0.28, [(ln, 11.5, False, C_BODY)])
            yy += 0.30
    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    build(prs)
    prs.save(OUT)
    W, H = prs.slide_width / IN, prs.slide_height / IN
    over = [round((sh.left + sh.width) / IN, 2)
            for sh in Presentation(OUT).slides[-1].shapes
            if (sh.left + sh.width) / IN > W + 0.01 or (sh.top + sh.height) / IN > H + 0.01]
    print(f'  {n0}장 → {len(Presentation(OUT).slides)}장')
    print(f'  슬라이드 밖: {len(over)}개  {"✅" if not over else over}')
    print(f'  저장 → {OUT}')


main()
