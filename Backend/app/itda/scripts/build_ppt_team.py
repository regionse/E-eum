# -*- coding: utf-8 -*-
"""팀 발표 슬라이드 보강 — 없는 장을 «신설»하고, 틀린 숫자를 고친다. (2026-08-07)

왜 만들었나
  `build_ppt_itda.py` 는 잇다 구간(PART 03)만 갈아끼운다. 그런데 42장 덱을 실측해 보니
  **팀 전체를 설명하는 장이 통째로 비어 있었다.**
    · 기술 스택   — 축별(덜다 9~16 · 잇다 17~29 · 나누다 30~38)로 흩어져 있고
                    「전체 스택이 뭐냐」에 가리킬 장이 없다. React·Vite 는 어디에도 안 나온다.
    · 인증        — 45장본에는 있었으나(잇다 — 인증) 42장본에서 «사라졌다».
                    bcrypt·JWT 를 설명할 자리가 0장이다.
    · 관리자 콘솔  — 슬라이드 4 에서 「관리자 콘솔이 사용자 화면과 같은 비중으로 존재한다」고
                    선언해 놓고, 그걸 보여주는 장이 «없다». 제작 의도의 약속을 안 지킨 덱이다.

  그리고 잇다 기술스택(18)은 ⓐ 태깅 수가 «1,064» 로 낡았고(실측 1,094)
  ⓑ 우측 패널의 텍스트 상자가 삭제돼 **빈 상자만 남아 있다**.

어떻게 하나
  `build_ppt_itda.py` 와 같은 방식 — 이 덱은 레이아웃 자리표시자가 없고 슬라이드마다 도형을
  직접 놓은 구조라, 새 슬라이드를 «만들면» 디자인이 날아간다. 기존 슬라이드 XML 을 복제하고
  글자만 바꾼다.

쓰는 법
  python build_ppt_team.py <원본.pptx> <출력.pptx>

⚠ 숫자는 전부 이 레포에서 «실측»한 값이다. 근거는 아래 MEASURED 주석 참고.
"""
import sys
import io
import copy

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\.claude\uploads\bafa128f-ace7-4371-aad4-be71ca976526\205ef8ff-__________.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v3.pptx'

# ── 실측 근거 (2026-08-07) ─────────────────────────────────────────────
#  MEASURED  프론트  React 18.3.1 · react-router-dom 6.30.4 · Vite 6.4.3
#                    (package-lock.json 의 «설치본» 기준. package.json 선언은 ^18.3.1 등)
#            화면    사용자 27 · 관리자 10   ← frontend/src/pages/**/*.jsx 계수
#            상태    Context API 만. redux/zustand/recoil 사용 파일 «0건» (grep)
#            백엔드  Python 3.12.5 · FastAPI 0.136.3 · SQLAlchemy 2.0.50 · aiomysql 0.3.2
#            API     사용자 49 · 관리자 21 = 70개 (app.routes 계수, FastAPI 자동문서 4개 제외)
#            인증    bcrypt 5.0.0 · PyJWT 2.13.0 · HS256 · 24시간 (user/security.py)
#            LLM     덜다 LangGraph(policy_recommendation_graph.py 1파일)
#                    나누다 google-genai SDK · 잇다 urllib 로 REST 직접 (gemini_util.py:22)
#            태깅    job_attr.tsv 1,094행 · act_type/obj_type 전량 채움 · obj_detail 102
#                    ⇒ 「태깅 1,064」는 2026-08-07 재태깅 이전 숫자다. 틀렸다.

TPL_A = 20      # 3카드 + 우패널   (원본 1-based · 잇다 시스템 구조)
TPL_C = 39      # 좌 큰카드 + 우 3카드 (원본 1-based · 인프라와 배포)


# ── 슬라이드 복제 (build_ppt_itda.py 와 동일) ──────────────────────
def clone(prs, idx0):
    """idx0(0-based) 슬라이드를 복제해 «맨 뒤»에 붙이고 새 슬라이드를 준다."""
    src = prs.slides[idx0]
    dst = prs.slides.add_slide(src.slide_layout)
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(sh._element))
    for k, rel in src.part.rels.items():
        if rel.reltype.endswith('slideLayout'):
            continue
        #  🔴🔴 2026-08-07 — **notesSlide 를 복사하면 PowerPoint 가 파일을 «거부»한다.**
        #    발표자 노트 파트는 자기 안에 «원본 슬라이드로 되돌아가는» 참조를 갖고 있다.
        #    그대로 복사하면 노트 하나를 슬라이드 둘이 가리키게 되고, 그 노트는 여전히
        #    «원본»을 가리킨다 → 관계가 순환/모순이 되어 PowerPoint 가 못 연다.
        #
        #    ⚠ 이게 지독한 이유: **python-pptx 도, zip 무결성 검사도 전부 통과한다.**
        #      슬라이드 수도 맞고 텍스트도 다 들어 있다. 실제로 PowerPoint 로 열어보기
        #      전까지는 «완성»으로 보인다. 리허설 당일에 발견했다.
        #      ⇒ 생성한 pptx 는 반드시 PowerPoint COM 으로 «열어서» 확인할 것.
        #        (scripts/verify_ppt.py 가 그 일을 한다)
        #
        #    새 슬라이드에 노트는 필요 없다. 필요하면 빈 노트를 새로 만들어야지
        #    원본 것을 «공유»하면 안 된다.
        if rel.reltype.endswith('notesSlide'):
            continue
        if rel.is_external:
            dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            dst.part.rels.get_or_add(rel.reltype, rel._target)
    return dst


def reorder(prs, order):
    """order(0-based) 순서로 재배열."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for x in ids:
        lst.remove(x)
    for i in order:
        lst.append(ids[i])


_RID = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'


def drop(prs, idx0s):
    """슬라이드 삭제 (0-based). ⚠ 관계(rel)까지 끊어야 파일에서 실제로 빠진다."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for i in sorted(idx0s, reverse=True):
        prs.part.drop_rel(ids[i].get(_RID))
        lst.remove(ids[i])


def put(shape, text):
    """도형의 글자를 바꾼다. 첫 run 의 서식을 살리고 나머지 run 은 지운다."""
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    paras = tf.paragraphs
    p0 = paras[0]
    if not p0.runs:
        return
    lines = str(text).split('\n')
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    p0.runs[0].text = lines[0]
    for p in list(paras[1:]):
        p._p.getparent().remove(p._p)
    for ln in lines[1:]:
        newp = copy.deepcopy(p0._p)
        p0._p.getparent().append(newp)
        from pptx.text.text import _Paragraph
        np = _Paragraph(newp, tf)
        for r in list(np.runs[1:]):
            r._r.getparent().remove(r._r)
        np.runs[0].text = ln


def fill(slide, mapping):
    shapes = list(slide.shapes)
    for i, txt in mapping.items():
        if txt is None or i >= len(shapes):
            continue
        put(shapes[i], txt)


#  TPL_A(20): 0 헤더 · 1 발표자 · 2 제목
#             3/4/5 카드1(배경·제목·본문) · 6/7/8 카드2 · 9/10/11 카드3
#            12/13/14 우패널(배경·제목·본문)
#  ⚠ 우패널 제목은 «한 줄»이어야 한다 — W7.4" · 34pt 라 15자를 넘으면 두 줄이 되어 넘친다.
def slide_a(prs, head, title, c1, c2, c3, right):
    s = clone(prs, TPL_A - 1)
    fill(s, {0: head, 2: title,
             4: c1[0], 5: c1[1], 7: c2[0], 8: c2[1], 10: c3[0], 11: c3[1],
             13: right[0], 14: right[1]})
    return s


#  TPL_C(39): 0 헤더 · 1 발표자 · 2 제목
#             4 좌라벨(18pt) · 5 좌제목(28pt) · 6 좌본문
#             8/9 우카드1(H1.3 — 가장 길게 쓸 수 있다) · 11/12 우카드2 · 14/15 우카드3
def slide_c(prs, head, title, left, r1, r2, r3):
    s = clone(prs, TPL_C - 1)
    fill(s, {0: head, 2: title,
             4: left[0], 5: left[1], 6: left[2],
             8: r1[0], 9: r1[1], 11: r2[0], 12: r2[1], 14: r3[0], 15: r3[1]})
    return s


# ════════════════════════════════════════════════════════════════
#  내용
# ════════════════════════════════════════════════════════════════

def build(prs):
    made = {}

    # ── ★신설 05 · 기술 스택 ──────────────────────────────────
    #  「전체 스택이 뭐냐」에 가리킬 장이 없었다. React·Vite 가 덱 어디에도 안 나온다.
    made['stack'] = slide_a(
        prs, '05 · 기술 스택', '세 축이 같은 기반 위에 올라가 있습니다',
        ('프론트 — React 18 · Vite 6',
         'react-router-dom 6.30 으로 화면 37개(사용자 27 · 관리자 10)를 붙였다. '
         '상태는 Context API 로만 관리한다 — Redux·Zustand 를 안 썼다.'),
        #  ⚠ 「라우터 8개」로 썼다가 실측에서 틀린 걸 잡았다 — router.py 를 가진 도메인은 9개다
        #    (dashboard · delda · inquiry · itda · mypage · nanuda · notice · notifications · user).
        ('백엔드 — FastAPI 한 서버',
         'Python 3.12 · FastAPI · SQLAlchemy 2.0 async + aiomysql. '
         '도메인 9개를 한 서버에 모았다. API 70개 — 사용자 49 · 관리자 21.'),
        #  ⚠ 「셋 다 3.1 Flash-Lite」로 썼다가 실측에서 틀린 걸 잡았다.
        #    잇다만 gemini-3.1-flash-lite(.env COURSE_LLM_MODEL) 다.
        #    덜다 policy_recommendation_service.py:60 · 나누다 llm_scorer.py:17 은 gemini-2.5-flash.
        ('AI · 데이터 — 축마다 «다르게» 쓴다',
         '셋 다 Gemini 지만 버전도 부르는 법도 다르다 — 잇다 3.1 Flash-Lite · REST 직접(urllib), '
         '덜다·나누다 2.5 Flash · LangGraph/SDK.'),
        ('공통 기반은 하나입니다',
         '세 축이 각자 서버를 띄우지 않는다.\n'
         'API 서버 한 대 · DB 한 곳(RDS) · 인증 한 곳.\n'
         '벡터는 Pinecone, 정형은 MySQL 이다.\n'
         '\n'
         '※ 버전은 실제 설치본을 실측한 값이다.'))

    # ── 잇다 기술 스택 «재작성» ───────────────────────────────
    #  ⓐ 「태깅 1,064」가 낡았다(실측 1,094 · 전량)
    #  ⓑ 우측 패널의 텍스트 상자가 삭제돼 «빈 상자»만 남아 있었다 → 본보기에서 복원한다
    made['itda_stack'] = slide_a(
        prs, 'PART 03 · 잇다 — 기술 스택', '재료는 평범합니다',
        ('LLM · 임베딩',
         'Gemini 3.1 Flash-Lite · REST 직접(urllib) · responseSchema 로 JSON 강제. '
         '임베딩은 gemini-embedding-2 · 3,072차원.'),
        ('검색 · 융합 · 리랭크',
         'Pinecone(ns 3개 · 코사인) + MySQL FULLTEXT(ngram) → RRF k=60 → '
         'Jina v3.5 → Pinecone bge-v2-m3 폴백.'),
        #  ⚠ 강좌를 「8,371」로 적고 있었다 — 2026-08-07 RDS 실측은 8,273 이다(98개 적다).
        #    문서(인수인계·기술스택·발표원고)에도 8,371 이 남아 있으니 같이 고칠 것.
        ('서버 · 데이터',
         'FastAPI · SQLAlchemy async · EC2 + RDS. NCS 직업 1,094(전량 태깅) · '
         '자격증 613 + 시험일정 2,655 · 강좌 8,273.'),
        ('여기까지는 흔합니다',
         '하이브리드 검색도 RRF 도 누구나 씁니다.\n'
         '덜다도 씁니다.\n'
         '\n'
         '다음 장부터가 본론입니다.\n'
         '※ 잇다는 LangGraph 를 안 씁니다 — 덜다만.'))

    # ── ★신설 07 · 인증과 권한 ────────────────────────────────
    #  45장본의 「잇다 — 인증」이 42장본에서 사라졌다. 인증은 잇다가 아니라 «공통»이므로
    #  잇다 안이 아니라 독립 장으로 되살린다.
    made['auth'] = slide_c(
        prs, '07 · 인증과 권한', '토큰이 아니라 «소유권»이 문제였다',
        ('기본 설계', 'bcrypt + JWT(HS256)',
         '비밀번호는 bcrypt 로 해싱해 저장한다 — 평문은 안 남는다. '
         '로그인하면 사용자 번호와 관리자 여부를 담은 24시간 토큰을 발급한다.'),
        ('토큰이 유효해도 다시 본다',
         '정지·휴면·탈퇴 계정은 토큰이 아직 안 만료됐어도 막는다. '
         '발급 시점의 권한을 그대로 믿지 않고, 요청마다 계정 상태를 DB 에서 다시 확인한다.'),
        ('관리자는 «두 번» 확인한다',
         '로그인 확인을 통과한 뒤 관리자 여부를 한 번 더 본다. '
         '관리자 API 21개가 전부 이 관문을 지난다.'),
        #  ⚠ 이 상자는 H0.9" 다(우1은 H1.3"). 원본 68자가 상한선이라 짧게 쓴다.
        ('진짜 구멍은 «소유권»이었다',
         'session_id 를 프론트가 만든다. 남의 것을 알면 남의 대화를 받아갈 수 있었다. '
         'owner 를 찍어 막았다.'))

    # ── ★신설 08 · 관리자 콘솔 ────────────────────────────────
    #  슬라이드 4 가 「관리자 콘솔이 사용자 화면과 같은 비중으로 존재한다」고 선언하는데
    #  그걸 보여주는 장이 «0장»이었다. 제작 의도의 약속을 덱이 안 지키고 있었다.
    made['admin'] = slide_a(
        prs, '08 · 관리자 콘솔', '기관이 운영한다면, 운영 화면이 있어야 합니다',
        ('화면 10개 · API 21개',
         '대시보드 · 회원 · 관리자 계정 · 공지 · 문의, 그리고 덜다 정책과 잇다 진로 데이터의 '
         '최신화. 사용자 화면 27개에 대해 관리자 화면 10개를 뒀다.'),
        ('데이터 최신화를 «사람이» 누른다',
         '버튼을 누르면 202 로 바로 응답하고 백그라운드에서 돈다. '
         '화면은 진행 상황을 폴링해 단계별로 보여준다 — 응답을 기다리며 멈추지 않는다.'),
        ('화면이 «계산하지» 않는다',
         '현황은 배치가 돌 때 기록한 값을 그대로 읽는다(itda_sync_log · content_hash). '
         '그래야 「언제 무엇이 바뀌었나」가 남는다.'),
        ('제작 의도와 이어집니다',
         '「기관이 운영하는 돌봄 SaaS」라고 했습니다.\n'
         '개인용 앱이면 관리자 화면이 필요 없습니다.\n'
         '\n'
         '지금은 사람이 누릅니다.\n'
         '다음은 스케줄러입니다 — 확장 계획 참고.'))

    return made


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    print(f'  원본 {n0}장')

    # ── ① 기존 슬라이드 수정 (원본 0-based 로 «먼저» 손본다) ──────────
    #  잇다 표지 — 인수인계장 §8 「발표에서 하면 안 되는 말」 첫 줄이 그대로 들어가 있었다:
    #    ✗「LLM 을 안 믿는 구조로 짰습니다」 → 낱말표에 의존하는 층이 있어 첫 질문에 무너진다
    #  실제로 하는 일(근거 대조)로 바꾼다. 이건 코드로 방어된다(verify_slots).
    put(list(prs.slides[16].shapes)[6], '1 · 모델이 낸 값을 «원문과 대조»한다')

    #  섹션 번호 재정렬 — 신설 3장이 들어가면서 뒤가 밀린다
    put(list(prs.slides[6].shapes)[0], '06 · 역할 분담')       # was 05
    put(list(prs.slides[38].shapes)[0], '09 · 인프라와 배포')   # was 06
    put(list(prs.slides[39].shapes)[0], '10 · 확장 계획')       # was 07
    #  ⚠ 41번(제작 소감)은 «내용을 건드리지 말라» 는 지시가 있었다.
    #    헤더의 «번호만» 맞춘다. 세 사람의 소감 칸은 그대로 둔다.
    put(list(prs.slides[40].shapes)[0], '11 · 제작 소감')       # was 08

    # ── ② 신설 슬라이드 생성 (맨 뒤에 붙는다: 42,43,44,45 → 0-based) ──
    made = build(prs)
    print(f'  신설 {len(made)}장 생성 (기술스택 · 잇다기술스택 재작성 · 인증 · 관리자)')

    # ── ③ 낡은 잇다 기술스택(0-based 17) 제거 ────────────────────────
    #     우패널 텍스트 상자가 삭제돼 빈 상자만 남은 판이라 되살리는 대신 갈아끼운다.
    drop(prs, [17])
    #  drop 뒤 인덱스: 0..16 그대로 · 옛 18..41 → 17..40 · 신설분 → 41,42,43,44

    # ── ④ 순서 재배열 ────────────────────────────────────────────────
    order = (
        list(range(0, 6))       # 1~6   표지 ~ 04 시스템 구성
        + [41]                  # ★ 05 · 기술 스택
        + [6]                   #   06 · 역할 분담
        + [7]                   #   시연
        + list(range(8, 16))    #   PART 01 덜다
        + [16]                  #   잇다 표지
        + [42]                  # ★ 잇다 — 기술 스택 (재작성)
        + list(range(17, 28))   #   잇다 나머지
        + list(range(28, 37))   #   PART 02 나누다
        + [43]                  # ★ 07 · 인증과 권한
        + [44]                  # ★ 08 · 관리자 콘솔
        + [37]                  #   09 · 인프라와 배포
        + [38]                  #   10 · 확장 계획
        + [39]                  #   11 · 제작 소감
        + [40]                  #   THANK YOU
    )
    assert len(order) == len(set(order)) == n0 - 1 + len(made), \
        f'순서 목록이 어긋났다: {len(order)}개 · 중복 {len(order) - len(set(order))}'
    reorder(prs, order)

    prs.save(OUT)
    print(f'  저장 → {OUT}')
    print(f'  전체 {len(Presentation(OUT).slides)}장')


main()
