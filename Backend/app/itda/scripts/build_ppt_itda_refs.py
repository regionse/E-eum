# -*- coding: utf-8 -*-
"""잇다 발표 «부록» — 설계 결정별 참고문헌 1장. (2026-08-10)

왜 만드나
  발표 본편(16~25번)은 6분 안에 끝내야 해서 «문헌 이야기를 할 시간이 없다».
  그래서 부록으로 뺀다. 말하지 않는 슬라이드다 — Q&A 에서 질문이 나오면 그때 넘긴다.

출처
  uploads/6f800641-*.md 「부록 · 설계 결정별 참고문헌」 (사용자가 Claude-in-Chrome 로 조사)
  9개 결정 × 3~4편. 여기엔 «결정당 대표 1편»만 싣는다. 전문은 그 문서에 있다.

톤 (2026-08-10 사용자 지적으로 고침)
  처음엔 결정마다 「뒷받침/반박」 딱지를 붙였다. 걷어냈다.
  6분짜리 «소개»에 자기비판 열을 다는 건 변명이지 정보가 아니다.
  부록은 「이 결정은 이 문헌에 기댔다」만 보이면 된다.
  물어보면 그때 정직하게 답하면 되고, 조사 문서에 소견은 그대로 남아 있다.

쓰는 법 (Backend/ 에서)
  python app/itda/scripts/build_ppt_itda_refs.py
  → C:/Users/TFX255GS/Downloads/이음_잇다_부록_참고문헌.pptx   이후 verify_ppt.py 로 «열어» 확인
"""
import io
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r'C:\Users\TFX255GS\Downloads\이음_잇다_부록_참고문헌.pptx'

C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_EYE = RGBColor(0x8B, 0x84, 0xA0)
C_ACC = RGBColor(0x6B, 0x3F, 0xA0)
C_TXT = RGBColor(0x24, 0x1F, 0x2E)
C_RED = RGBColor(0xA8, 0x36, 0x4A)
C_DEEP = RGBColor(0x32, 0x16, 0x6F)
C_PANEL_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_BG = RGBColor(0xEF, 0xEB, 0xF6)
C_BG_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Malgun Gothic'


def _font(r, size, color, bold):
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.18, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        for (t, size, color, bold) in line:
            r = p.add_run()
            r.text = t
            _font(r, size, color, bold)
    return tb


def box(sl, x, y, w, h, fill=C_WHITE, line=C_PANEL_LN, lw=0.75):
    shp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.055
    except Exception:                                   # noqa: BLE001
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


#  결정당 대표 1편. 인용은 조사 문서 원문 그대로 옮긴 것이며 «지어낸 것이 없다».
CARDS = [
    ('1', '구조화 출력으로 슬롯 뽑기',
     'Usman (2026). PhantomFill. arXiv:2607.20492',
     '필수 필드로 답을 «강제»하면 모델이 값을 지어낸다.\n'
     '그래서 잇다는 모르면 빈 배열을 낼 수 있게 열어 뒀다.'),
    ('2', '값마다 근거(발화 인용)를 받기',
     'Slobodkin et al. (2024). Attribute First, then Generate. NAACL',
     '먼저 근거를 고르고 그다음 생성하면 사실성과 검증가능성이 함께 오른다.\n'
     '잇다는 값과 근거를 같이 받아 코드가 발화와 대조한다.'),
    ('3', '얕은 문자열 대조로 인용 검증',
     'McNamee & Mayfield (2004). Character N-Gram Tokenization. IR Journal',
     '굴절이 많은 언어는 형태소 분석 없이 문자 단위 대조만으로도 검색이 선다.\n'
     '잇다가 인용 검증에 임베딩 대신 글자를 쓰는 근거.'),
    ('4', '하이브리드 검색 + RRF',
     'Bruch, Gai & Ingber (2023). Fusion Functions for Hybrid Retrieval. ACM TOIS',
     '「기존 연구와 달리 RRF 는 파라미터에 민감하다」 — 작은 k 가 9개 중 8개에서 이겼다.\n'
     '잇다 실측도 같은 방향이라 k 를 60 에서 1 로 내렸다.'),
    ('5', '여러 질의를 만들어 합치기',
     'Abe, Kato, Takeoka & Oyamada (2025). LLM-based Query Expansion Fails…',
     '질의 확장은 «낯설고 모호한» 질의에서 오히려 해롭다.\n'
     '잇다가 검색어에 지금까지의 슬롯을 «닻»으로 함께 넣는 이유.'),
    ('6', '점수 «분포»로 되물을지 정하기',
     'Zendel, Liu, Culpepper & Scholer (2023). Entropy-Based QPP. SIGIR',
     '점수 분포의 엔트로피로 검색 품질을 «답하기 전에» 예측한다.\n'
     '잇다 2단 판정의 직접 기전. 실측 정확도 75 → 87.7%.'),
    ('7', '계층 요약으로 긴 대화 접기',
     'Mao et al. (2026). Agents Don\'t Just Agree, They Remember. arXiv',
     '알게 된 것을 «저장하는 순간» 출처 표시가 사라지는 게 공통 실패 모드.\n'
     '값에 출처와 확신도를 함께 붙이는 것이 잇다의 다음 설계.'),
    ('8', '안전 게이트를 2층으로',
     'Pichowicz, Kotas & Piotrowski (2025). Scientific Reports',
     '정신건강 챗봇 29종 중 5종이 낱말 인식 때문에 «발화 전송 자체»를 막았다.\n'
     '그래서 잇다는 낱말을 신호로만 쓰고 판정은 문맥으로 한다.'),
    ('9', '생성 «후» 출력 검증',
     'Röttger et al. (2024). XSTest. NAACL',
     '안전 규칙을 세게 걸면 «과잉 거부»가 생긴다 — 안전·대조 프롬프트를 짝지어 잰다.\n'
     '잇다가 정상 발화 35건을 대조군으로 함께 두는 이유.'),
]


def main():
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    sl = prs.slides.add_slide(prs.slide_layouts[6])

    text(sl, 1.00, 0.80, 12.0, 0.45,
         [[('부록 · PART 03 잇다 — 설계 결정별 참고문헌', 18.75, C_EYE, False)]])
    text(sl, 14.0, 0.74, 5.0, 0.55, [[('발표자 김정준', 30, C_ACC, True)]], align=PP_ALIGN.RIGHT)
    text(sl, 1.00, 1.32, 18.1, 0.9,
         [[('설계 결정 9개 — ', 43, C_TITLE, True),
           ('무엇에 기대어 정했나', 43, C_ACC, True)]])
    text(sl, 1.00, 2.18, 18.1, 0.42,
         [[('결정당 대표 1편입니다. 원문을 직접 읽고 옮겼고, 전문 27편은 별도 문서에 있습니다.',
            18, C_EYE, False)]])

    xs = (1.00, 7.37, 13.74)
    ys = (2.78, 5.42, 8.06)
    for i, (no, name, cite, note) in enumerate(CARDS):
        x, y = xs[i % 3], ys[i // 3]
        vc = C_ACC
        box(sl, x, y, 6.09, 2.46)
        text(sl, x + 0.34, y + 0.22, 5.45, 0.44,
             [[(f'결정 {no}', 17, C_EYE, True)]])
        text(sl, x + 0.34, y + 0.62, 5.45, 0.5,
             [[(name, 23, C_TITLE, True)]])
        text(sl, x + 0.34, y + 1.10, 5.45, 0.42,
             [[(cite, 14.5, vc, True)]], ls=1.16)
        text(sl, x + 0.34, y + 1.52, 5.45, 0.86,
             [[(ln, 14.5, C_TXT, False)] for ln in note.split('\n')], ls=1.28)

    text(sl, 1.0, 10.66, 18.1, 0.44,
         [[('IR · NLP · 정신건강 정보학 세 갈래에서 모았습니다. 학회는 SIGIR · ACM TOIS · NAACL · EMNLP · Scientific Reports.',
            16.5, C_EYE, False)]])

    prs.save(OUT)
    print(f'저장: {OUT} · 1장')


if __name__ == '__main__':
    main()
