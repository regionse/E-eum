# -*- coding: utf-8 -*-
"""관문을 «하는 일»로 다시 묶은 그림. (2026-08-07)

왜 다시 그리나 — 사용자 지적
  「사용자 자연어를 받자마자 7개를 한꺼번에 돌리는데, 그럼 4개는 의미가 없는 거네?
   3개만 돌려도 되는 거네?」

  코드로 확인해 보니 **「7중 게이트」라는 틀 자체가 틀렸다.** 실제로 `return` 해서
  대화를 끊는 건 ①⑤⑦ 셋뿐이고, 나머지는 «상태를 갱신»한다.
  그런데 3개만 남기면 안 된다 — **①은 혼자서 아무것도 못 하기 때문이다:**

      ④ 불법신호 · ⑤ 가해/욕설 · ⑦ 공격   →  bump_abuse()  →  _abuse +1
                                                   ↓
      ① abuse_limits 는 그 _abuse 를 «읽기만» 한다 → 문턱 넘으면 차단

  ④ 를 빼면 _abuse 가 안 올라 ① 이 «영원히» 안 걸린다.
  레드팀 15턴에 카운터가 0 이었던 게 정확히 그 상태였다.

  ⇒ 「막는 문 7개」가 아니라 **「막는 문 하나 + 그 문에 쓸 숫자를 세는 것들 + 메모 + 갈림길」**

코드 근거 (itda_core.py)
  5343  ①  if _abuse >= abuse_limits()[1]:  return blocked      ← 유일하게 «끊는» 문
  5371  ②  _exclude 에 담고 칩 강등.  return 없음
  5408  ③  profile['_policy_declined']=True.  return 없음
  5428  ④  bump_abuse('불법신호')          ← 세기만
  5430  ⑤  pre_check → SILENT/VAGUE/SELFHARM/HARM/UNSAFE (네 갈래)
  5470  ⑤  bump_abuse('harm')  ·  5489 bump_abuse('unsafe')
  5509  ⑥  pick_from_options.  return 없음
  5560  ⑦  is_injection → return redirect  ·  5561 bump_abuse('injection')

쓰는 법
  python make_gate4_diagram.py [바탕.pptx] [출력.pptx]
"""
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR            # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v11.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v12.pptx'

IN = 914400
FONT = 'Arial'
C_HEAD  = RGBColor(0x8B, 0x84, 0xA0)
C_PRES  = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_BODY  = RGBColor(0x55, 0x50, 0x6B)
C_MUTED = RGBColor(0x8B, 0x84, 0xA0)
C_CARD  = RGBColor(0xFF, 0xFF, 0xFF)
C_CARDL = RGBColor(0xE7, 0xE2, 0xF0)
C_LINE  = RGBColor(0xC9, 0xC1, 0xDC)

C_BLOCK = RGBColor(0xA8, 0x36, 0x4A)      # 잠근다
C_COUNT = RGBColor(0xC2, 0x6B, 0x2E)      # 센다
C_MEMO  = RGBColor(0x4E, 0x7D, 0x5E)      # 기억한다
C_FORK  = RGBColor(0x6B, 0x3F, 0xA0)      # 가른다
BG = {C_BLOCK: RGBColor(0xF9, 0xEC, 0xEF), C_COUNT: RGBColor(0xFB, 0xF0, 0xE6),
      C_MEMO: RGBColor(0xEC, 0xF3, 0xEE), C_FORK: RGBColor(0xEF, 0xEB, 0xF6)}


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.2):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.color.rgb = col
    return tb


def box(sl, x, y, w, h, fill, line, lw=1.0):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False; sh.text_frame.text = ''
    try:
        sh.adjustments[0] = 0.02
    except Exception:
        pass
    return sh


def seg(sl, x1, y1, x2, y2, col=C_LINE, w=2.0):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Emu(int(x1 * IN)), Emu(int(y1 * IN)),
                                Emu(int(x2 * IN)), Emu(int(y2 * IN)))
    c.line.color.rgb = col; c.line.width = Pt(w)
    return c


#  (묶음이름, 색, 한 줄 설명, 안에 든 것)
GROUPS = [
    ('기억한다', C_MEMO, '막지 않는다. 메모만 하고 대화는 그대로 간다',
     '②「그거 말고」→ 후보 제외    ③「괜찮아요」→ 안내 중단    ⑥「1번이요」→ 선택 수신'),
    ('센다', C_COUNT, '이것도 안 막는다. _abuse 만 올린다',
     '④ 불법 신호        ⑤ 가해 · 욕설        ⑦ 공격'),
    ('잠근다', C_BLOCK, '쌓인 _abuse 가 문턱을 넘으면 — 여기서 «끝»',
     '① 이미 잠긴 사람인가'),
    ('가른다', C_FORK, '낱말 116개가 «어디로 보낼지» 정한다',
     '⑤ pre_check   →   통과 · 되묻기 · ★잇기 · 차단'),
]

WHY = [
    ('막는 문은 «하나»뿐입니다', [
        '코드에서 return 으로 대화를 끊는 건 ①⑤⑦ 셋.',
        '나머지 넷은 상태만 바꾸고 그냥 지나갑니다.',
    ]),
    ('①은 «혼자서» 아무것도 못 합니다', [
        '① 은 _abuse 값을 읽기만 합니다.',
        '그 값을 올리는 건 ④⑤⑦ 입니다.',
        '④ 를 빼면 _abuse 가 안 올라 ① 이 «영원히» 안 걸립니다 —',
        '레드팀 15턴에 카운터가 0 이었던 게 그 상태입니다.',
    ]),
    ('그럼 왜 «세기만» 하나', [
        '낱말로 막으면 반드시 과차단합니다.',
        '실제로 「성폭력 상담사」를 막았다가 고쳤습니다.',
        '⇒ 판정은 LLM 에 맡기고, 우리는 «몇 번 있었나»만 셉니다.',
    ]),
    ('기억하는 셋도 빼면 사고가 납니다', [
        '② 없으면 「그거 말고」 해도 같은 카드가 또 나옵니다.',
        '③ 없으면 거절한 정책을 계속 권합니다.',
        '⑥ 없으면 「1번이요」가 근거 없다고 버려집니다.',
    ]),
]


def build(prs):
    s = blank_like(prs)
    text(s, 1.00, 0.92, 9.0, 0.30, [('PART 03 · 잇다 — 예외 처리', 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRES)])
    text(s, 1.00, 1.50, 19.0, 0.95,
         [('막는 문은 하나, 나머지는 «그 문에 쓸 숫자»를 셉니다', 40, True, C_TITLE)])

    # ── 왼쪽 : 왜 ───────────────────────────────────────
    LX, LW = 1.00, 7.35
    y = 2.72
    for head, lines in WHY:
        text(s, LX, y, LW, 0.32, [(head, 17, True, C_TITLE)])
        y += 0.38
        for ln in lines:
            text(s, LX + 0.26, y, LW - 0.26, 0.28, [(ln, 13.5, False, C_BODY)])
            y += 0.28
        y += 0.24

    # ── 오른쪽 : 네 묶음 ────────────────────────────────
    RX, RW = 8.95, 10.05
    SPINE = RX + 1.30

    #  발화
    box(s, RX + 3.00, 2.72, 4.10, 0.52, C_CARD, C_CARDL, 0.75)
    text(s, RX + 3.00, 2.85, 4.10, 0.30,
         [('🧑  사용자 발화 한 마디', 15, True, C_TITLE)], align=PP_ALIGN.CENTER)

    BH, GAP = 1.44, 0.30
    ys = [3.56 + i * (BH + GAP) for i in range(4)]
    seg(s, SPINE, 3.24, SPINE, ys[-1] + BH / 2)

    for (name, col, sub, items), yy in zip(GROUPS, ys):
        seg(s, SPINE, yy + BH / 2, RX + 1.86, yy + BH / 2, col, 2.5)
        box(s, RX + 1.86, yy, RW - 1.86, BH, BG[col], col, 1.25)
        text(s, RX + 2.24, yy + 0.22, 2.10, 0.40, [(name, 21, True, col)])
        text(s, RX + 4.45, yy + 0.28, RW - 4.75, 0.30, [(sub, 13.5, False, C_MUTED)])
        text(s, RX + 2.24, yy + 0.80, RW - 2.60, 0.34, [(items, 14.5, True, C_TITLE)])

    #  ★ 「센다」 → 「잠근다」 의존 화살표 (이 그림의 요점)
    ax = RX + RW + 0.06
    seg(s, ax, ys[1] + BH / 2, ax + 0.34, ys[1] + BH / 2, C_COUNT, 2.5)
    seg(s, ax + 0.34, ys[1] + BH / 2, ax + 0.34, ys[2] + BH / 2, C_COUNT, 2.5)
    seg(s, ax + 0.34, ys[2] + BH / 2, ax, ys[2] + BH / 2, C_COUNT, 2.5)
    text(s, ax - 1.30, ys[1] + BH + 0.02, 1.60, 0.28,
         [('_abuse 누적', 12.5, True, C_COUNT)], align=PP_ALIGN.RIGHT)

    #  끝 — LLM 으로
    ye = ys[-1] + BH + 0.34
    seg(s, SPINE, ys[-1] + BH / 2, SPINE, ye + 0.26)
    box(s, RX + 3.00, ye, 4.10, 0.52, BG[C_FORK], C_FORK, 1.25)
    text(s, RX + 3.00, ye + 0.13, 4.10, 0.30,
         [('통과한 것만 LLM 으로', 15, True, C_TITLE)], align=PP_ALIGN.CENTER)
    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    build(prs)
    prs.save(OUT)
    W, H = prs.slide_width / IN, prs.slide_height / IN
    over = [round((sh.left + sh.width) / IN, 2)
            for sh in Presentation(OUT).slides[-1].shapes
            if (sh.left + sh.width) / IN > W + 0.01 or (sh.top + sh.height) / IN > H + 0.01]
    print(f'  {n0}장 → {len(Presentation(OUT).slides)}장')
    print(f'  슬라이드 밖: {len(over)}개  {"✅" if not over else over}')
    print(f'  저장 → {OUT}')


main()
