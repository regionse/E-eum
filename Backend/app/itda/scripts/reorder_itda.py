# -*- coding: utf-8 -*-
"""잇다 구간 «순서 정리» + 「무엇을 만들었나」 한 장. (2026-08-07)

왜 필요한가
  오늘 만든 잇다 슬라이드 11장이 전부 «맨 뒤»(43~54)에 몰려 있다. 그 상태로는
  발표 순서가 안 맞고 대사도 못 쓴다.

골격은 덜다에서 가져온다 (실측 — v15 슬라이드 10~17)
    표지 → 전체 시스템 구조 → 데이터 수집·정규화 → 핵심 기술 5장
  잇다도 같은 골격으로 세운다:
    표지 → 무엇을 만들었나 → 데이터 → 관문 → 칸 → 찾기 → 답내기 → 문제와 해결

버리는 8장 — 새 장이 «같은 내용을 더 낫게» 말한다
    19 기능            사용자가 이미 지웠다고 함
    20 시스템 구조      관문 + 7칸 + 3-a 로 쪼개져 대체됨
    21 잇다 기술 스택   3-a 가 벡터·FULLTEXT·RRF·리랭커를 다 말한다. 중복
    22 슬롯            → 44 「한 문장을 7개의 칸으로」
    23 검색            → 52 「어떻게 찾나」
    24 착지 판정        → 53 「바로 줄까, 한 번 더 물을까」
    45 관문 (구판)      → 51 로 대체 (순서를 코드대로 고친 판)
    50 관문 네 묶음(구판) → 51 로 대체

부록으로 미루는 3장 — «측정» 이야기라 소개 흐름에 안 맞는다
    43 검사 결과 · 48 착지 무게 · 49 신호 9개 비교

⚠ 덜다·나누다·공통 장은 «한 장도» 안 건드린다.
"""
import sys
import io
import copy

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                           # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v15.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v16.pptx'

IN = 914400
FONT = 'Arial'
C_HEAD  = RGBColor(0x8B, 0x84, 0xA0)
C_PRES  = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_BODY  = RGBColor(0x55, 0x50, 0x6B)
C_MUTED = RGBColor(0x8B, 0x84, 0xA0)
C_CARD  = RGBColor(0xFF, 0xFF, 0xFF)
C_CARDL = RGBColor(0xE7, 0xE2, 0xF0)
C_PANEL = RGBColor(0xEF, 0xEB, 0xF6)
C_PANLL = RGBColor(0xDE, 0xD8, 0xEA)

_RID = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.24):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls
        r = p.add_run(); r.text = t
        r.font.name = FONT; r.font.size = Pt(pt)
        r.font.bold = bold; r.font.color.rgb = col
    return tb


def box(sl, x, y, w, h, fill, line, lw=1.0):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False; sh.text_frame.text = ''
    try:
        sh.adjustments[0] = 0.02
    except Exception:
        pass
    return sh


def reorder(prs, order):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for x in ids:
        lst.remove(x)
    for i in order:
        lst.append(ids[i])


def drop(prs, idx0s):
    """⚠ 관계(rel)까지 끊어야 파일에서 실제로 빠진다."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for i in sorted(idx0s, reverse=True):
        prs.part.drop_rel(ids[i].get(_RID))
        lst.remove(ids[i])


# ════════════════════════════════════════════════════════════════
#  「무엇을 만들었나」 — 매우 간략하게
# ════════════════════════════════════════════════════════════════
CARDS = [
    ('누가 쓰나', '가족돌봄청년', [
        '돌봄이 학업·취업을 밀어냅니다.',
        '진로를 설계할 시간과 정보가',
        '«가장 먼저» 사라집니다.',
    ]),
    ('무엇을 하나', '대화로 방향을 잡습니다', [
        '자격증·강좌를 붙여',
        '«미래설계지도»로 남깁니다.',
        '다음에 이어서 대화할 수 있습니다.',
    ]),
    ('왜 «대화»인가', '검색창은 안 맞습니다', [
        '검색창은 «무엇을 찾을지 아는 사람»의',
        '도구입니다. 우리 사용자는 그걸',
        '모르는 상태로 옵니다.',
    ]),
]


def slide_intro(prs):
    s = blank_like(prs)
    text(s, 1.00, 0.92, 9.0, 0.30, [('PART 03 · 잇다', 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRES)])
    text(s, 1.00, 1.50, 19.0, 0.95,
         [('돌봄 뒤의 «다음»을 같이 찾습니다', 44, True, C_TITLE)])

    W, GAP = 5.86, 0.21
    for i, (label, head, lines) in enumerate(CARDS):
        x = 1.00 + i * (W + GAP)
        box(s, x, 2.96, W, 3.30, C_CARD if i < 2 else C_PANEL,
            C_CARDL if i < 2 else C_PANLL, 1.1)
        text(s, x + 0.44, 3.28, W - 0.88, 0.28, [(label, 14, False, C_MUTED)])
        text(s, x + 0.44, 3.66, W - 0.88, 0.46, [(head, 24, True, C_TITLE)])
        yy = 4.34
        for ln in lines:
            text(s, x + 0.44, yy, W - 0.88, 0.32, [(ln, 15.5, False, C_BODY)])
            yy += 0.34

    #  한 줄 요약
    box(s, 1.00, 6.70, 18.00, 1.20, C_PANEL, C_PANLL, 1.25)
    text(s, 1.00, 7.06, 18.00, 0.46,
         [('“할머니 돌보느라 학교를 그만뒀어요”  한 마디에서 시작합니다',
           26, True, C_TITLE)], align=PP_ALIGN.CENTER)
    text(s, 1.00, 8.14, 18.00, 0.34,
         [('다음 장부터 — 그 한 문장이 «직업 카드»가 되기까지', 16, False, C_MUTED)],
         align=PP_ALIGN.CENTER)
    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    slide_intro(prs)                       # 새 장 → 인덱스 n0 (=54)
    NEW = n0

    #  ── 최종 순서 (0-based · v15 기준) ─────────────────────────
    KEEP = (
        list(range(0, 9))          # 공통 앞 9장 (표지 ~ 시연 구분자)
        + list(range(9, 17))       # 덜다 8장
        + [17]                     # 잇다 표지
        + [NEW]                    # ★ 무엇을 만들었나
        + [53]                     # 데이터 출처
        + [50]                     # 관문 — 하는 일이 넷
        + [46]                     # 코드 — 관문
        + [43]                     # 한 문장을 7개의 칸으로
        + [45]                     # 코드 — 7개의 칸
        + [51]                     # 3-a 어떻게 찾나
        + [52]                     # 3-b 어떻게 답을 내나
        + [24, 25, 26]             # 문제와 해결 3장
        + list(range(27, 36))      # 나누다 9장
        + list(range(36, 42))      # 공통 뒤 6장 (인증 ~ THANK YOU)
        + [42]                     # 부록 · 검사 결과
        + [47, 48]                 # 부록 · 착지 무게 · 신호 비교
    )
    DROP = [18, 19, 20, 21, 22, 23, 44, 49]

    assert len(set(KEEP)) == len(KEEP), '중복'
    assert set(KEEP) | set(DROP) == set(range(n0 + 1)), \
        f'빠진 인덱스: {set(range(n0 + 1)) - set(KEEP) - set(DROP)}'

    #  ⚠ 지울 것을 «맨 뒤»로 몰아 두고 한 번에 지운다 — drop 이 인덱스를 흔들지 않게.
    reorder(prs, KEEP + DROP)
    drop(prs, list(range(len(KEEP), len(KEEP) + len(DROP))))

    prs.save(OUT)
    out = Presentation(OUT)
    print(f'  {n0}장 → {len(out.slides)}장   (신설 1 · 버림 {len(DROP)})')
    print(f'  잇다 구간: 18~29 (12장)')
    print(f'  저장 → {OUT}')


main()
