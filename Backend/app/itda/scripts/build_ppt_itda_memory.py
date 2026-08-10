# -*- coding: utf-8 -*-
"""잇다 발표 — «기억» 한 장. (2026-08-10 신설)

왜 따로 만드나
  · build_ppt_itda_main.py 를 지금 다른 작업이 편집 중이라 그 파일을 건드리지 않는다.
    import 로 헬퍼를 빌려 오는 것도 피한다 — 그 파일이 바뀌면 이것도 같이 깨진다.
    그래서 디자인 상수·헬퍼를 **그대로 베껴 자기완결**로 둔다(수정안(4) 실측값 동일).
  · 완성되면 PowerPoint 「슬라이드 재사용」으로 본편에 끼워 넣는다 —
    main 빌더 머리말이 정한 방식 그대로다. 슬라이드 복제는 **0회**.

왜 이 장이 필요한가
  덱 8장 ①이 「대화 접기」를 **앞으로 할 일**로 적어 두었는데, 계층 요약은
  2026-08-06 에 **이미 들어가 있다**(실측기록 §12-4). 그래서 「그거 이미 하시는 거 아닌가요」가
  나올 수 있다. 지금 «돌아가는 것»을 그리는 장을 따로 둔다.

숫자 출처 (전부 실측 — 지어낸 값 없음)
  · +17.4% / +15.8%  200턴 원시기록 `checks/_persona_long_{1,2}.json` 의 turn 별 `used.in`
                     1~20턴 평균 대 181~200턴 평균 (2026-08-10 재집계)
  · 사실 30개        같은 파일의 최종 `facts` 길이 — 두 페르소나 모두 30
  · 접기 14회        같은 파일 `folded` 합계 — 두 페르소나 모두 14
  · 요약 210~240자   같은 파일 최종 `summary` 길이 (238자 · 210자)
                     ⚠ 「300자」로 적을 뻔했다 — 그건 «옛 방식» 주석의 값이다.
                       지금 코드 상한은 `summarize()` 의 `[:800]` 이고, 실측은 210~240자다.
                       주석을 근거로 슬라이드 숫자를 쓰면 안 된다는 걸 여기서 또 확인했다.
  · 상한값           itda_core.HISTORY_TURNS=8 · SUMMARIZE_AFTER=20 ·
                     FACTS_HEAD=12 · FACTS_TAIL=18 · controllers.py:205 의 keep

쓰는 법 (Backend/ 에서)
  python app/itda/scripts/build_ppt_itda_memory.py
  → C:/Users/TFX255GS/Downloads/이음_잇다_기억.pptx
  이후 verify_ppt.py 로 «실제로 열어» 확인한다.
"""
import io
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r'C:\Users\TFX255GS\Downloads\이음_잇다_기억.pptx'

#  ── 팔레트 (수정안(4) 실측 — build_ppt_itda_main.py 와 동일) ──────────
C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_EYE = RGBColor(0x8B, 0x84, 0xA0)
C_ACC = RGBColor(0x6B, 0x3F, 0xA0)
C_TXT = RGBColor(0x24, 0x1F, 0x2E)
C_RED = RGBColor(0xA8, 0x36, 0x4A)
C_DEEP = RGBColor(0x32, 0x16, 0x6F)
C_PANEL_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_BG = RGBColor(0xEF, 0xEB, 0xF6)
C_BG_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_CHIP = RGBColor(0xF0, 0xE9, 0xFB)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LINE = RGBColor(0xC9, 0xBE, 0xE0)
FONT = 'Malgun Gothic'


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_font(r, size, color, bold):
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    #  한글 폰트는 East-Asian 속성(ea)에도 넣어야 실제로 먹는다
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.18, anchor=MSO_ANCHOR.TOP):
    """runs = [줄, 줄, ...] · 줄 = [(글, size, color, bold), ...]"""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        for (t, size, color, bold) in line:
            r = p.add_run()
            r.text = t
            _set_font(r, size, color, bold)
    return tb


def box(sl, x, y, w, h, fill=C_WHITE, line=C_PANEL_LN, lw=0.75, round_=True):
    shp = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try:
            shp.adjustments[0] = 0.06
        except Exception:                                  # noqa: BLE001
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def _arrow_head(ln):
    el = ln.line._get_or_add_ln()
    head = el.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    el.append(head)


def arrow_v(sl, x, y1, y2):
    ln = sl.shapes.add_connector(2, Inches(x), Inches(y1), Inches(x), Inches(y2))
    ln.line.color.rgb = C_LINE
    ln.line.width = Pt(2.25)
    _arrow_head(ln)
    return ln


def head(sl, crumb, title):
    text(sl, 1.00, 0.86, 12.0, 0.45, [[(crumb, 18.75, C_EYE, False)]])
    text(sl, 14.0, 0.80, 5.0, 0.55, [[('발표자 김정준', 30, C_ACC, True)]],
         align=PP_ALIGN.RIGHT)
    text(sl, 1.00, 1.42, 18.1, 0.95, [[(title, 48, C_TITLE, True)]])


def flowbox(sl, x, y, w, h, title, sub=None, fill=C_WHITE, tcol=C_DEEP, tsz=22, ssz=17):
    """제목 + (선택) 부제 한 줄이 든 흐름 상자 — 세로 중앙 정렬."""
    box(sl, x, y, w, h, fill=fill)
    lines = [[(title, tsz, tcol, True)]]
    if sub:
        lines.append([(sub, ssz, C_EYE, False)])
    text(sl, x + 0.18, y, w - 0.36, h, lines, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, ls=1.05)


# ══════════════════════════════════════════════════════════════════════
def s_memory(prs):
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다 — 기억', '기억 — 200턴을 지나도 «쌓이지 않는다»')

    #  ── 왼쪽 : 구조 ────────────────────────────────────────────────
    box(sl, 1.0, 2.65, 10.9, 8.1)
    text(sl, 1.45, 2.92, 10.0, 0.5,
         [[('세 겹으로 «접는다»', 30, C_TITLE, True),
           ('   — 버리는 게 아니다', 17, C_EYE, False)]])

    flowbox(sl, 1.6, 3.66, 9.7, 1.20, '최근 대화 — 8턴을 «원문 그대로»',
            '각 300자에서 자른다 · 8턴에서 멈춘다 — 무한정 늘리는 게 답이 아니라서')

    arrow_v(sl, 6.45, 4.86, 5.30)
    text(sl, 6.72, 4.88, 4.6, 0.38,
         [[('40개가 넘으면 — 앞부분을 «접는다»', 16.5, C_EYE, False)]])

    flowbox(sl, 1.6, 5.30, 4.72, 1.38, '확인된 사실',
            '최대 30개 · «코드»가 누적한다\n다시 압축하지 않는다', tsz=23, ssz=15.5)
    flowbox(sl, 6.58, 5.30, 4.72, 1.38, '서사 요약',
            '매번 다시 쓴다 — 흐름·말투를 잇는다\n200턴 뒤에도 210~240자', tsz=23, ssz=15.5)

    #  ── 왼쪽 아래 : 왜 «가운데»를 버리나 ────────────────────────────
    box(sl, 1.6, 7.05, 9.7, 3.35, fill=C_BG, line=C_BG_LN)
    text(sl, 1.95, 7.28, 9.0, 0.5,
         [[('넘칠 땐 «가운데»를 버린다', 26, C_TITLE, True)]])
    text(sl, 1.95, 7.95, 9.0, 2.3, [
        [('앞 12개 ', 20, C_ACC, True),
         ('초반은 정체성·형편이라 «안 변한다»', 18, C_TXT, False)],
        [('           「할머니 간병하느라 학교를 못 다녔다」', 16.5, C_EYE, False)],
        [('뒤 18개 ', 20, C_ACC, True),
         ('최근은 지금의 관심이라 «선명해야 한다»', 18, C_TXT, False)],
        [('가운데 ', 20, C_EYE, True),
         ('흐려도 되는 자리 — 그래서 여기를 버린다', 18, C_TXT, False)],
    ], ls=1.30)
    text(sl, 1.95, 9.92, 9.0, 0.42,
         [[('옛 방식은 매번 다시 압축해서, 가장 안 변하는 초반부터 흐려졌다',
            16.5, C_RED, False)]])

    #  ── 오른쪽 : 실측 4개 ──────────────────────────────────────────
    cells = [
        ('+17%', '200턴 동안 늘어난 입력',
         '페르소나 둘 +17.4% · +15.8%\n선형으로 쌓였다면 몇 배가 된다'),
        ('30개', '사실 목록의 상한',
         '앞 12 + 뒤 18. 200턴 실측이\n정확히 이 상한에 닿았다'),
        ('14회', '200턴 동안 접은 횟수',
         '12턴마다 한 번 — 매 턴 부르면\n그게 더 비싸다'),
        ('0초', '사용자가 기다린 시간',
         '접기는 응답을 보낸 «뒤»\n백그라운드로 돈다'),
    ]
    x, h, gap = 12.2, 1.86, 0.22
    y = 2.65
    for big, cap, sub in cells:
        box(sl, x, y, 6.9, h)
        text(sl, x + 0.40, y + 0.16, 6.1, 0.72, [[(big, 40, C_ACC, True)]])
        text(sl, x + 0.40, y + 0.86, 6.1, 0.40, [[(cap, 19, C_TXT, True)]])
        text(sl, x + 0.40, y + 1.26, 6.1, 0.55, [[(sub, 14.5, C_EYE, False)]], ls=1.20)
        y += h + gap


def main():
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    s_memory(prs)
    prs.save(OUT)
    print(f'저장: {OUT} · {len(prs.slides._sldIdLst)}장')


if __name__ == '__main__':
    main()
