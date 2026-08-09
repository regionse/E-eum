# -*- coding: utf-8 -*-
"""「낱말표는 차단기가 아니라 갈림길」 — 그림판. (2026-08-07)

사용자 지적
  「저런 식 말고 그림 같은 걸로 보여주면 좋겠는데. 저건 너무 다 직설적이야」
  「6번은 그냥 다른 거잖아. 빼버리는 게 어떰? 애당초 6번이 왜 필요함?」

무엇을 바꿨나
  ① 표 → «갈림길». pre_check 의 11줄짜리 목록을 «네 갈래»로 줄인다.
     낱말 하나하나는 발표자가 말로 한다. 슬라이드는 「네 갈래로 갈린다」만 보여준다.
  ② 관문에서 ⑥ 칩 선택을 «뺀다» → 6개.
     ①②③④⑤⑦ 은 «막거나 돌려보내는» 관문인데 ⑥ 만 «답을 받는» 동작이라 성질이 다르다.
     ⚠ 인수인계장 남은과제 ⑨ 에 「칩 선택을 관문 목록에서 빼고 되묻기→답변수신 별도
       흐름으로」가 이미 적혀 있다. 코드 정리는 «안» 했고, 여기선 그림만 그렇게 그린다.
     ⇒ 그래서 아래쪽에 «별도 흐름»으로 따로 그린다. 숨기는 게 아니라 자리를 옮기는 것이다.

⚠ 값은 전부 코드 실측이다(make_gate_diagram.py 머리말 참고).

쓰는 법
  python make_fork_diagram.py [바탕.pptx] [출력.pptx]
"""
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR            # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v6.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v8.pptx'

IN = 914400
FONT = 'Arial'
C_HEAD    = RGBColor(0x8B, 0x84, 0xA0)
C_PRESENT = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE   = RGBColor(0x3B, 0x20, 0x63)
C_BODY    = RGBColor(0x55, 0x50, 0x6B)
C_CARD    = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_PANEL   = RGBColor(0xEF, 0xEB, 0xF6)
C_PANEL_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_MUTED   = RGBColor(0x8B, 0x84, 0xA0)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

#  네 갈래의 색 — 「막는 것」과 「막지 않는 것」이 한눈에 갈려야 한다
C_PASS    = RGBColor(0x4E, 0x7D, 0x5E)
C_ASK     = RGBColor(0x6B, 0x3F, 0xA0)
C_CONNECT = RGBColor(0xC2, 0x6B, 0x2E)
C_BLOCK   = RGBColor(0xA8, 0x36, 0x4A)
C_PASS_BG    = RGBColor(0xEC, 0xF3, 0xEE)
C_ASK_BG     = RGBColor(0xEF, 0xEB, 0xF6)
C_CONNECT_BG = RGBColor(0xFB, 0xF0, 0xE6)
C_BLOCK_BG   = RGBColor(0xF9, 0xEC, 0xEF)


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, ls=1.2):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.color.rgb = col
    return tb


def shape(sl, kind, x, y, w, h, fill, line=None, lw=0.75):
    sh = sl.shapes.add_shape(kind, Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.text = ''
    return sh


def line(sl, x1, y1, x2, y2, col, w=2.0):
    cn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Emu(int(x1 * IN)), Emu(int(y1 * IN)),
                                 Emu(int(x2 * IN)), Emu(int(y2 * IN)))
    cn.line.color.rgb = col
    cn.line.width = Pt(w)
    return cn


# ── 네 갈래 ──────────────────────────────────────────────────────
#   (제목, 결과이름, 무슨 일이 일어나나, 예시, 글자색, 배경색)
FORKS = [
    ('그냥 지나간다', '통 과',
     '슬롯 채우기로 그대로 넘어간다',
     '“ㅎㅇ”   ·   “저 혼자 해요”   ·   “자살예방 상담사가 되고싶어요”',
     C_PASS, C_PASS_BG),
    ('말 대신 쓴 표현', '되묻기',
     '차단하지 않고 «다르게» 물어본다',
     '“…”   ·   “ㅠㅠ”   ·   “ㅁㄴㅇㄹ”',
     C_ASK, C_ASK_BG),
    ('★ 힘든 이야기', '잇 기',
     '☎109 · 129 로 «잇고», 대화는 그대로 이어간다',
     '“죽고 싶어요”   ·   “할머니가 저를 때려요”',
     C_CONNECT, C_CONNECT_BG),
    ('공격 · 희롱', '차 단',
     '되돌려 보낸다',
     '상대를 향한 욕설   ·   명백한 희롱',
     C_BLOCK, C_BLOCK_BG),
]

# ── 관문 6개 (칩 선택은 «뺐다») ────────────────────────────────────
GATES6 = ['이미 잠긴 사람인가', '「그거 말고」인가', '「괜찮아요」인가',
          '불법 신호인가', 'pre_check  ★', '공격인가']


def build(prs):
    s = blank_like(prs)

    text(s, 1.00, 0.92, 8.0, 0.30, [('PART 03 · 잇다 — 예외 처리', 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRESENT)])
    text(s, 1.00, 1.50, 19.0, 0.95,
         [('낱말표는 «차단기»가 아니라 «갈림길»입니다', 44, True, C_TITLE)])

    # ══ 관문 6개 — 작은 칩으로 한 줄 ═══════════════════════════
    text(s, 1.00, 2.68, 3.0, 0.30,
         [('LLM 앞 관문 6개', 16, True, C_MUTED)])
    x = 1.00
    for i, g in enumerate(GATES6):
        hot = g.endswith('★')
        w = 2.62 if not hot else 2.62
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, x, 3.08, w, 0.56,
              C_PANEL if hot else C_CARD, C_PRESENT if hot else C_CARD_LN,
              1.5 if hot else 0.75)
        text(s, x, 3.22, w, 0.30,
             [(g, 14.5, hot, C_TITLE if hot else C_BODY)], align=PP_ALIGN.CENTER)
        if i < len(GATES6) - 1:
            text(s, x + w, 3.20, 0.34, 0.30,
                 [('›', 17, False, C_MUTED)], align=PP_ALIGN.CENTER)
        x += w + 0.34
    text(s, 1.00, 3.78, 12.0, 0.30,
         [('전부 코드 · 턴당 0원 · 약 1밀리초 — 여기서 걸리면 LLM 을 한 번도 안 부릅니다',
           14.5, False, C_MUTED)])

    # ══ 갈림길 ════════════════════════════════════════════════
    #  발화 말풍선 → pre_check 원 → 세로 줄기 → 네 갈래
    BY = 4.55
    ROW_H, ROW_GAP = 1.34, 0.22
    rows_y = [BY + i * (ROW_H + ROW_GAP) for i in range(4)]
    #  ⚠ 노드를 «가운데»에 두면 왼쪽 아래가 통째로 빈다. 첫 갈래 높이에 맞춰 «위»로 올리고,
    #    비는 자리에 설명 글을 넣는다 — 사용자 요청: 「글자로 왼쪽에 적어주면 그게 소개 같음」
    node_cy = rows_y[0] + ROW_H / 2

    #  ① 발화 말풍선
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, 1.00, node_cy - 0.56, 2.80, 1.12, C_CARD, C_CARD_LN)
    text(s, 1.00, node_cy - 0.38, 2.80, 0.32,
         [('사용자 발화', 14, False, C_MUTED)], align=PP_ALIGN.CENTER)
    text(s, 1.00, node_cy - 0.04, 2.80, 0.42,
         [('한 마디', 22, True, C_TITLE)], align=PP_ALIGN.CENTER)

    #  ② pre_check 노드
    shape(s, MSO_SHAPE.OVAL, 4.22, node_cy - 0.74, 2.45, 1.48, C_PANEL, C_PRESENT, 2.0)
    text(s, 4.22, node_cy - 0.34, 2.45, 0.38,
         [('pre_check', 20, True, C_TITLE)], align=PP_ALIGN.CENTER)
    text(s, 4.22, node_cy + 0.04, 2.45, 0.32,
         [('낱말 116개', 13.5, False, C_MUTED)], align=PP_ALIGN.CENTER)

    shape(s, MSO_SHAPE.RIGHT_ARROW, 3.86, node_cy - 0.15, 0.32, 0.30, C_PANEL_LN)

    #  ③ 줄기 + 네 갈래
    STEM = 7.62
    line(s, STEM, node_cy, STEM, rows_y[-1] + ROW_H / 2, C_PANEL_LN, 2.5)
    line(s, 6.67, node_cy, STEM, node_cy, C_PANEL_LN, 2.5)

    # ══ 왼쪽 설명 — «세는 것»과 «문턱» ════════════════════════
    #  코드 실측: bump_abuse 호출부 4곳 · abuse_limits() 를 직접 돌려 얻은 값
    LX, LY, LW = 1.00, 6.42, 5.95
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, LX, LY, LW, 3.88, C_CARD, C_CARD_LN)
    y = LY + 0.30
    text(s, LX + 0.38, y, LW - 0.76, 0.34,
         [('세는 것 — 네 가지', 18, True, C_TITLE)])
    y += 0.44
    text(s, LX + 0.38, y, LW - 0.76, 0.32,
         [('불법 신호  ·  가해 위협  ·  상대 향한 욕설  ·  탈옥 시도', 14.5, False, C_BODY)])
    y += 0.38
    text(s, LX + 0.38, y, LW - 0.76, 0.30,
         [('⚠  피해 진술과 혼잣말 좌절은 «안» 셉니다', 13.5, False, C_MUTED)])
    y += 0.32
    text(s, LX + 0.38, y, LW - 0.76, 0.32,
         [('★  자해는 «절대» 세지 않습니다', 15, True, C_CONNECT)])

    y += 0.50
    text(s, LX + 0.38, y, LW - 0.76, 0.34,
         [('문턱 — 대화가 길수록 넉넉하게', 18, True, C_TITLE)])
    y += 0.42
    #  ⚠ 한 줄에 두 공식을 넣었더니 «두 줄로 넘쳐» 아래와 겹쳤다. 한 줄에 하나씩.
    text(s, LX + 0.38, y, LW - 0.76, 0.92,
         [('좋은 턴 = 전체 턴 − 남용 턴', 14.5, False, C_BODY),
          ('경고 = 3 + 좋은턴÷5    (최대 10)', 14.5, False, C_BODY),
          ('종료 = 6 + 좋은턴÷3    (최대 20)', 14.5, False, C_BODY)])
    y += 1.02
    text(s, LX + 0.38, y, LW - 0.76, 0.30,
         [('6턴에 6번 떠들면 → 잠김        40턴에 6번이면 → 안 잠김',
           13.5, True, C_PRESENT)])

    BOX_X, BOX_W = 8.30, 10.70
    for (title, res, what, ex, col, bg), y in zip(FORKS, rows_y):
        cy = y + ROW_H / 2
        line(s, STEM, cy, BOX_X - 0.42, cy, col, 2.5)
        shape(s, MSO_SHAPE.ISOSCELES_TRIANGLE, BOX_X - 0.42, cy - 0.13, 0.30, 0.26, col).rotation = 90

        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, BOX_X, y, BOX_W, ROW_H, bg, col, 1.25)
        #  왼쪽: 결과 이름 (크게)
        text(s, BOX_X + 0.36, y + 0.30, 1.85, 0.50,
             [(res, 24, True, col)])
        #  가운데: 무슨 일이 일어나나
        text(s, BOX_X + 2.35, y + 0.24, BOX_W - 2.75, 0.34,
             [(title, 14, False, C_MUTED)])
        text(s, BOX_X + 2.35, y + 0.58, BOX_W - 2.75, 0.34,
             [(what, 17, True, C_TITLE)])
        text(s, BOX_X + 2.35, y + 0.94, BOX_W - 2.75, 0.30,
             [(ex, 13.5, False, C_MUTED)])

    #  ④ 아래 한 줄 — 이 그림이 하려는 말
    #    ⚠ 왼쪽 설명 상자가 10.30 까지 내려오므로 «오른쪽 갈래 폭 안»에만 놓는다.
    #      전체 폭으로 두면 상자와 겹친다(앞 판에서 실제로 겹쳤다).
    text(s, 8.30, 10.42, 10.70, 0.44,
         [('낱말이 하는 일은 «막는 것»이 아니라, 어느 길로 보낼지 «정하는 것»입니다',
           17, True, C_PRESENT)])

    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    build(prs)
    prs.save(OUT)
    W, H = prs.slide_width / IN, prs.slide_height / IN
    over = [sh for sh in Presentation(OUT).slides[-1].shapes
            if (sh.left + sh.width) / IN > W + 0.01 or (sh.top + sh.height) / IN > H + 0.01]
    print(f'  {n0}장 → {len(Presentation(OUT).slides)}장')
    print(f'  슬라이드 밖으로 나간 도형: {len(over)}개  {"✅" if not over else "🔴"}')
    print(f'  저장 → {OUT}')


main()
