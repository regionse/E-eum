# -*- coding: utf-8 -*-
"""데이터 출처 — «표 한 장». 앞 판을 갈아끼운다. (2026-08-07)

앞 판(make_data_slide.py)이 왜 안 읽혔나 — 사용자 지적
  「가독성 너무 별로야. 위에서 아래로. 표처럼 보여줘. 간략하게.
   어떤 API를 / 어떤 파라미터를 / 어떻게 바꿔서 / 어디에 넣냐 —
   이런 식으로만. 그냥 글자 크게 보일 수 있게.」

무엇을 바꿨나
  ① 아래 «배운 것» 세 칸을 통째로 뺐다. 표만 남긴다.
  ② 그 자리를 4행에 나눠 줘서 행 높이를 1.70" 로 키웠다.
  ③ 글자를 키웠다 —  API 이름 22pt · 파라미터 17pt · 변환 16pt (앞 판은 12~13.5pt)
  ④ 열 제목을 사용자가 말한 «그 문장»으로:  어떤 API를 / 어떤 파라미터를 / 어떻게 바꿔서 / 어디에

⚠ 이 스크립트는 «갈아끼운다» — 새로 붙이지 않는다.
  v16 의 잇다 데이터 장(20번, 0-based 19)을 새 것으로 바꾸고 옛 것을 지운다.
"""
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                           # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v16.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v17.pptx'
TARGET = 19          # 0-based — v16 의 「데이터는 전부 공공데이터입니다」

IN = 914400
FONT, CODE = 'Arial', 'Consolas'
C_HEAD  = RGBColor(0x8B, 0x84, 0xA0)
C_PRES  = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_BODY  = RGBColor(0x55, 0x50, 0x6B)
C_MUTED = RGBColor(0x8B, 0x84, 0xA0)
C_CARD  = RGBColor(0xFF, 0xFF, 0xFF)
C_CARDL = RGBColor(0xE7, 0xE2, 0xF0)
C_PANEL = RGBColor(0xEF, 0xEB, 0xF6)
C_PANLL = RGBColor(0xDE, 0xD8, 0xEA)

_RID = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.22, font=FONT):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = ls
        r = p.add_run(); r.text = t
        r.font.name = font; r.font.size = Pt(pt)
        r.font.bold = bold; r.font.color.rgb = col
    return tb


def box(sl, x, y, w, h, fill, line, lw=1.0):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False; sh.text_frame.text = ''
    try:
        sh.adjustments[0] = 0.02
    except Exception:
        pass
    return sh


def arrow(sl, x, y, w=0.34, h=0.26):
    sh = sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid(); sh.fill.fore_color.rgb = C_PANLL
    sh.line.fill.background(); sh.shadow.inherit = False
    sh.text_frame.text = ''
    return sh


#  (API, 기관, 파라미터 2줄, 바꾼 것 2줄, 테이블, 건수)
ROWS = [
    ('NCS 직무', '국가직무능력표준',
     ['직무 코드 · 직무명', '직무 정의 · 대/중/소분류'],
     ['직무 정의를 카드 설명에', '«그대로» 씁니다'],
     'job_catalog', '1,094'),
    ('Q-Net 자격증', '한국산업인력공단',
     ['종목 코드 · 종목명 · 등급', '직무 분야 · 자격 구분'],
     ['종목 코드를 키로 «덮어쓰기»', '다시 돌려도 안 쌓입니다'],
     'certification', '613'),
    ('Q-Net 시험일정', '공공데이터포털',
     ['연도 · 회차', '접수일 · 시험일 · 발표일'],
     ['문자열로 오는 날짜를', '«날짜 형식»으로 통일'],
     'exam_schedule', '2,655'),
    ('K-MOOC 강좌', '공공데이터포털',
     ['강좌 이름 · 분류', '요약 · 링크'],
     ['요약에 HTML 태그가 섞여', '있어 «걷어냅니다»'],
     'course', '8,273'),
]


def build(prs):
    s = blank_like(prs)
    text(s, 1.00, 0.92, 9.0, 0.30, [('PART 03 · 잇다 — 데이터', 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRES)])
    text(s, 1.00, 1.50, 19.0, 0.95,
         [('데이터는 «전부» 공공데이터입니다', 44, True, C_TITLE)])

    #  열 (x, 폭)
    C1, C2, C3, C4 = (1.40, 3.90), (5.95, 4.80), (11.45, 4.30), (16.35, 2.65)
    A1, A2, A3 = 5.42, 10.92, 15.87

    for x, w, name in ((C1[0], C1[1], '어떤 API 를'), (C2[0], C2[1], '어떤 파라미터를'),
                       (C3[0], C3[1], '어떻게 바꿔서'), (C4[0], C4[1], '어디에 넣나')):
        text(s, x, 2.84, w, 0.30, [(name, 15, True, C_MUTED)])

    y, RH, GAP = 3.26, 1.70, 0.14
    for api, org, params, norm, table, n in ROWS:
        box(s, 1.00, y, 18.00, RH, C_CARD, C_CARDL, 1.0)
        #  ① 어떤 API 를
        text(s, C1[0], y + 0.42, C1[1], 0.40, [(api, 22, True, C_TITLE)])
        text(s, C1[0], y + 0.94, C1[1], 0.30, [(org, 13.5, False, C_MUTED)])
        arrow(s, A1, y + RH / 2 - 0.13)
        #  ② 어떤 파라미터를
        text(s, C2[0], y + 0.48, C2[1], 0.76,
             [(params[0], 17, False, C_TITLE), (params[1], 17, False, C_TITLE)])
        arrow(s, A2, y + RH / 2 - 0.13)
        #  ③ 어떻게 바꿔서
        text(s, C3[0], y + 0.52, C3[1], 0.72,
             [(norm[0], 16, False, C_BODY), (norm[1], 16, False, C_BODY)])
        arrow(s, A3, y + RH / 2 - 0.13)
        #  ④ 어디에 넣나
        box(s, C4[0], y + 0.36, C4[1], 0.98, C_PANEL, C_PANLL, 1.0)
        text(s, C4[0], y + 0.52, C4[1], 0.34,
             [(table, 18, True, C_TITLE)], align=PP_ALIGN.CENTER, font=CODE)
        text(s, C4[0], y + 0.94, C4[1], 0.30,
             [(f'{n} 행', 14.5, False, C_MUTED)], align=PP_ALIGN.CENTER)
        y += RH + GAP

    text(s, 1.00, y + 0.16, 18.00, 0.34,
         [('직접 만든 데이터는 «직업 태그» 하나뿐입니다', 16, True, C_PRES)],
         align=PP_ALIGN.CENTER)
    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    build(prs)                                  # 새 장 → 인덱스 n0

    #  옛 장을 «맨 뒤»로 보내고 새 장을 그 자리에 — 그다음 맨 뒤를 지운다
    order = [n0 if i == TARGET else i for i in range(n0)] + [TARGET]
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for x in ids:
        lst.remove(x)
    for i in order:
        lst.append(ids[i])
    #  이제 맨 뒤(마지막 하나)가 옛 장이다
    ids = list(lst)
    prs.part.drop_rel(ids[-1].get(_RID))
    lst.remove(ids[-1])

    prs.save(OUT)
    out = Presentation(OUT)
    got = ' '.join(sh.text_frame.text for sh in out.slides[TARGET].shapes
                   if sh.has_text_frame)[:40]
    print(f'  {n0}장 → {len(out.slides)}장  (갈아끼움)')
    print(f'  {TARGET+1}번 장: {got!r}')
    print(f'  저장 → {OUT}')


main()
