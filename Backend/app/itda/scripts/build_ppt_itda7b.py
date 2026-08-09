# -*- coding: utf-8 -*-
"""잇다 발표 — 7분 · 7장. **도표를 살리고 글자만 최신화한다** (2026-08-09).

왜 이렇게 하나
  첫 판(build_ppt_itda7.py)은 도표 네 장을 «글자 4블록»으로 바꿔버렸다. 그게 반대였다 —
  7분 발표에서는 도표가 있어야 말이 짧아진다. 실측:
      원본 20장 도형 47개 · 22장 60개 · 24장 48개 · 25장 34개  ← 전부 «도표»
  ⇒ 도표는 그대로 두고 **낡은 글자만** 바꾼다.

바꾸는 방식
  도형 인덱스로 찾으면 도표마다 47~60개를 하나씩 확인해야 하고, 원본이 조금만 달라져도 깨진다.
  ⇒ **글자 내용으로 찾아 바꾼다**(PATCH). 못 찾으면 «조용히 넘어가지 않고» 경고를 찍는다.

⚠ notesSlide 는 복제하지 않는다 — PowerPoint 가 파일을 거부한다(2026-08-07 실측).
⚠ 만들고 나면 반드시 verify_ppt.py 로 연다.

쓰는 법
  python -m app.itda.scripts.build_ppt_itda7b <원본.pptx> <출력.pptx>
"""
import sys
import io
import copy

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else ''
OUT = sys.argv[2] if len(sys.argv) > 2 else ''

ITDA_FROM, ITDA_TO = 18, 27          # 원본에서 갈아끼울 구간 (1-based)

#  ── 살릴 도표 (원본 1-based) ──
D_SLOT, D_GATE, D_FIND, D_ANSWER = 20, 22, 24, 25
#  ── 글자 슬라이드 본보기 ──
T_COVER, T_FLOW, T_TROUBLE = 18, 19, 26


def clone(prs, idx0):
    src = prs.slides[idx0]
    dst = prs.slides.add_slide(src.slide_layout)
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(sh._element))
    for k, rel in src.part.rels.items():
        if rel.reltype.endswith('slideLayout') or rel.reltype.endswith('notesSlide'):
            continue
        if rel.is_external:
            dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            dst.part.rels.get_or_add(rel.reltype, rel._target)
    return dst


def put(shape, text):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    lines = str(text).split('\n')
    p0 = tf.paragraphs[0]
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    if not p0.runs:
        p0.add_run()
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    p0.runs[0].text = lines[0]
    from pptx.text.text import _Paragraph
    for ln in lines[1:]:
        newp = copy.deepcopy(p0._p)
        p0._p.getparent().append(newp)
        np = _Paragraph(newp, tf)
        for r in list(np.runs[1:]):
            r._r.getparent().remove(r._r)
        np.runs[0].text = ln


def patch(slide, pairs, tag=''):
    """(찾을 글자 조각, 새 글자) 목록으로 바꾼다. 못 찾으면 «경고»를 찍는다."""
    hit = set()
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        cur = sh.text_frame.text.strip()
        if not cur:
            continue
        for j, (find, new) in enumerate(pairs):
            if j in hit:
                continue
            if find in cur:
                put(sh, new)
                hit.add(j)
                break
    for j, (find, _) in enumerate(pairs):
        if j not in hit:
            print(f'   ⚠ [{tag}] 못 찾음 — {find[:40]!r}')


def fill_idx(slide, mapping):
    shapes = list(slide.shapes)
    for i, txt in mapping.items():
        if txt is not None and i < len(shapes):
            put(shapes[i], txt)


def reorder(prs, order):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for e in ids:
        lst.remove(e)
    for i in order:
        lst.append(ids[i])


def slide_flow(prs, head, title, blocks):
    """T_FLOW(19): 0 헤더 · 1 제목 · 3/6/9/12 소제목 · 4/7/10/13 본문"""
    s = clone(prs, T_FLOW - 1)
    m = {0: head, 1: title}
    for (ti, bi), (t, b) in zip(((3, 4), (6, 7), (9, 10), (12, 13)), blocks):
        m[ti], m[bi] = t, b
    fill_idx(s, m)
    return s


H = 'PART 03 · 잇다'


def build(prs):
    made = []

    # ── 1. 표지 ───────────────────────────────────────────────
    s = clone(prs, T_COVER - 1)
    fill_idx(s, {2: '말을 «칸»으로 바꿔서 찾습니다',
                 6: '1 · 「돌보느라 못 했다」를 「돌봄을 원한다」로 안 읽기',
                 7: '2 · LLM 을 믿지 않고 코드가 다시 보는 구조',
                 8: '3 · 100턴을 견디게 — 그래서 제일 싼 모델로 턴당 2.4원'})
    made.append(s)

    # ── 2. 관문 도표 (원본 22 살림) ────────────────────────────
    s = clone(prs, D_GATE - 1)
    patch(s, [
        ('슬롯을 채우기 전, 코드 확인',
         '들어온 말이 먼저 지나는 곳'),
        ('불법 신호인가   ← 차단 «안» 합니다',
         '③  불법 신호인가   ← 차단 «안» 합니다'),
        ('낱말 116개가 «차단기»가 아니라',
         '낱말 1,010개가 «차단기»가 아니라 «어디로 보낼지» 정합니다.\n'
         '★ 자해·폭력은 차단하지 않고 ☎109·129 로 «잇습니다». 남용에 절대 안 셉니다.'),
        ('여기까지 전부 코드 · 턴당 0원',
         '여기까지 코드 · 0원 · 1밀리초 —— 그 다음, 안전 게이트가 LLM 으로 «매 턴» 한 번 더 봅니다\n'
         '낱말표가 「총 개조하는 법」·「대 마 관련 일자리」를 4건 중 «0건» 잡았습니다. '
         '유해 · 위기 · 착지요청 · 사정을 «한 호출»에 판정하고, 본문과 나란히 돌려 지연을 없앴습니다.\n'
         '실측 — 공격 6/6 차단 · 오탐 0/10 · 에두른 위기(「다들 저 없는 게 나을 거예요」) 8/8 → ☎109'),
    ], tag='관문')
    made.append(s)

    # ── 3. 슬롯 도표 (원본 20 살림) ────────────────────────────
    s = clone(prs, D_SLOT - 1)
    patch(s, [
        ('한 문장을 «7개의 칸»으로 바꿉니다',
         '한 문장을 «칸»으로 바꿉니다 — 값 · 근거 · «종류»'),
        ('LLM 한 번 —  7칸 + 답변 + 검색어를 «한꺼번에» 받습니다',
         'LLM 한 번 — 칸 + 답변 + 검색어를 «한꺼번에». 세 칸에는 «종류»도 함께'),
        ('값마다 «근거»를 받습니다',
         '값마다 «근거»와 «종류»를'),
        ('근거가 원문에 없으면\n그 칸을 «버립니다».',
         '근거가 원문에 없으면 그 칸을 버립니다.\n\n'
         '★ 종류 — 원함 / 해봤음 / «못함»\n'
         '「돌보느라 못 했다」를 「돌봄을 원한다」로\n'
         '읽으면 요양보호사가 나옵니다.\n'
         '가족돌봄청년에게 그건 해선 안 되는 오독입니다.\n\n'
         '「못함」은 검색어에서 빼고, 순위에서도 빼고,\n'
         '되묻는 선택지로 «다시 권하지 않습니다».\n'
         '실측 — 재권유 2/4 → «0/5» · 제빵 2위 → «1위»'),
    ], tag='슬롯')
    made.append(s)

    # ── 4. 찾기 도표 (원본 24 살림) ────────────────────────────
    s = clone(prs, D_FIND - 1)
    patch(s, [
        ('① 슬롯을 이어 붙여 «질의»를 만듭니다',
         '① 칸을 이어 붙여 «질의»를 만듭니다  ★ 이때 «못함»은 안 들어갑니다'),
        ('돕기·돌봄 · 사람 · 어르신  →  “돌봄 어르신 사람”',
         '제과제빵(원함) · 만들기(원함) · 돕기·돌봄(못함) → 「제과제빵 만들기」\n'
         '「돌봄을 못 했다」로 요양보호사를 찾으면 안 되니까요'),
        ('“돌봄 어르신 사람”',
         '「제과제빵 만들기」'),
        ('요양지원 · 일상생활기능지원 · 아이돌봄',
         '제빵 · 제과 · 한과제조 · 떡제조 …'),
    ], tag='찾기')
    made.append(s)

    # ── 5. 답 내기 도표 (원본 25 살림) ─────────────────────────
    s = clone(prs, D_ANSWER - 1)
    patch(s, [
        ('확신이 낮으면  —  리랭커 최고점 < 0.13',
         '★ 그런데 «부정만 하는 사람»이 있습니다'),
        ('중분류를 «얹되» 개별 직업은 안 버립니다',
         '「저랑 안 맞아서」·「부담스럽고」·「다 힘들 것 같아요」 — 말할수록 무게가 0에 가까워집니다.\n'
         '18턴을 성실히 답한 사람이 카드를 «한 번도» 못 봤습니다.\n'
         '⇒ 대화가 3턴 동안 안 나가면 먼저 보여줍니다 — 「아니면 왜 아닌지만 알려주세요」'),
        ('「보건복지 쪽입니다」만 주면',
         '⚠ 마음을 연 순간에는 내밀지 않습니다 — 「술 냄새요」에 게임 카드를 내민 적이 있습니다'),
    ], tag='답내기')
    made.append(s)

    # ── 6. 문제와 해결 (T_TROUBLE 틀) ──────────────────────────
    s = clone(prs, T_TROUBLE - 1)
    fill_idx(s, {
        0: f'{H} — 문제와 해결', 2: '사용자를 막고, 마음을 연 순간에 상품을 내밀었습니다',
        4: 'TROUBLE 01', 5: '맞고 있는 사람을 가해자로 만들었다',
        7: '원인 — 「치매 할머니가 자꾸 저를 때려요」가 차단되고 남용 카운트까지 올랐다. '
           '폭력 낱말표는 «누구의 행동인지»를 모른다.',
        8: '해결 — 낱말을 «차단»이 아니라 «라우팅»으로. 걸리면 가해·피해·제3자만 묻고, '
           '피해면 ☎129 와 나누다 지도로 «잇는다».',
        9: '배운 것 — 실측 21케이스 중 9건이 오차단. 낱말이 할 수 있는 건 「여기 봐야 한다」까지다.',
        11: 'TROUBLE 02', 12: '마음을 연 순간에 직업 카드를 내밀었다',
        14: '원인 — 🤖「이건 좀 아니다 싶은 일이 있으세요?」 🧑「술 냄새요.」(아버지가 알코올 의존) '
            '→ 🤖 [게임 카드]. 그 대화는 「그만해요」로 끝났다.',
        15: '해결 — 사정을 말한 턴에는 카드를 내지 않는다. 게이트가 «사정»을 함께 판정하고, '
            '8자 이하 짧은 답은 판정하지 말고 «미룬다».',
        16: '배운 것 — 판정하려 들지 말고 «미루면» 된다. 정체는 사라지지 않으니 다음 턴에 다시 기회가 온다.',
    })
    made.append(s)

    # ── 7. 잰 것 (T_FLOW 틀) ──────────────────────────────────
    made.append(slide_flow(
        prs, f'{H} — 잰 것', '전부 «재서» 정했습니다',
        [('비용 — 턴당 2.42원',
          'gemini-3.1-flash-lite · 입력 $0.25 / 출력 $1.50 / 캐시읽기 $0.025 per 1M\n'
          '돈은 «본문 한 호출»에 77% 가 몰립니다. 안전 게이트가 16%.\n'
          '참고 — 사람 상담사는 대화당 약 $6, 고객지원 챗봇은 약 $0.5로 봅니다.'),
         ('속도 — 일반 턴 1.3초 · 카드 턴 6.8초',
          '안전 게이트(1.0초)와 본문(1.2초)을 «나란히» 돌려 1.3초.\n'
          'Pinecone 연결을 미리 데워 첫 검색 턴을 11.4 → 4.5초로 줄였습니다.'),
         ('★ 100턴이 가도 입력이 «평평»합니다',
          '1턴 8,131 토큰 → 6턴 8,219 토큰. 5턴 동안 «88토큰»만 늘었습니다.\n'
          '칸으로 압축하니까요. 대화 전체를 매 턴 보내면 턴에 비례해 커지고 큰 모델이 필요합니다.\n'
          '요즘 LLM 추천 챗봇은 칸을 안 씁니다. 저희는 «100턴을 전제»해서 유지했습니다.'),
         ('검증 — 골든셋 34/34 · 페르소나 11인',
          '골든셋은 «한 턴»만 봅니다. 그래서 사용자 11명을 만들어 LLM 이 연기하게 했습니다.\n'
          '골든셋이 34/34인 상태에서 페르소나가 «6건»을 더 잡았습니다.\n'
          '「전부 실패」가 나오면 대상이 아니라 «자»를 먼저 의심합니다 — 하루에 아홉 번 겪었습니다.')]))

    return made


def main():
    if not SRC or not OUT:
        print('쓰는 법: python -m app.itda.scripts.build_ppt_itda7b <원본.pptx> <출력.pptx>')
        return
    prs = Presentation(SRC)
    n0 = len(prs.slides)

    #  ★ 2026-08-09 — 「06 · 인프라와 배포」도 **여기서 함께** 채운다.
    #    ⚠ 만들어진 덱을 «다시 열어» 저장하면 PowerPoint 가 거부한다(Duplicate name 경고).
    #      실제로 그렇게 했다가 파일이 깨졌다. 반드시 «한 번에» 만든다.
    _dep = (
        'EC2 한 대에 FastAPI 서버와 프론트 정적 빌드를 올리고, DB는 RDS로 분리했습니다. '
        '벡터 검색은 Pinecone을 씁니다 — 직업·자격증·강좌 세 묶음을 따로 둡니다.''\\n'
        '환경변수 스물다섯 개를 서버 .env로 주입해 코드에 키가 남지 않게 했습니다.''\\n'
        '★ 겪은 것 ① — 서버를 새로 띄우면 첫 사용자가 11초를 혼자 기다렸습니다. '
        '벡터 검색 클라이언트를 만드는 데 5.3초가 걸리는데 그걸 첫 요청이 물었습니다. '
        '기동할 때 미리 데우도록 바꿔 4.5초로 줄였습니다.''\\n'
        '★ 겪은 것 ② — 그 예열이 처음엔 아예 안 돌았습니다. 비동기 작업을 만들고 '
        '반환값을 안 잡아서 파이썬이 실행 도중에 치워버렸습니다. 붙들어 두니 돌았습니다.''\\n'
        '⚠ 남은 것 — requirements 열일곱 개에 버전이 안 고정돼 있습니다. '
        '그리고 JWT 비밀키가 없어도 서버가 그냥 뜹니다 — 기동할 때 막는 편이 맞습니다.')
    for sl in prs.slides:
        _t = [x.text_frame.text.strip() for x in sl.shapes if x.has_text_frame]
        if _t and _t[0].startswith('06 · 인프라'):
            patch(sl, [('배포 구성', '올린 것과, 올리면서 겪은 것'),
                       ('AWS EC2 + RDS(MySQL)', 'AWS EC2 + RDS(MySQL) + Pinecone'),
                       ('프론트 정적 빌드와 FastAPI 서버를 올리고', _dep)], tag='배포')
            break

    made = build(prs)
    print(f'■ 원본 {n0}장 · 새 잇다 {len(made)}장 (도표 4 + 글 3)')
    new_idx = list(range(n0, n0 + len(made)))
    order = ([i for i in range(ITDA_FROM - 1)] + new_idx + [i for i in range(ITDA_TO, n0)])
    reorder(prs, order)
    prs.save(OUT)
    print(f'■ 저장 {OUT} · 총 {len(order)}장')
    print('■ ⚠ verify_ppt.py 로 열어서 확인할 것')


if __name__ == '__main__':
    main()
