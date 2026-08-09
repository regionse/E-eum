# -*- coding: utf-8 -*-
"""발표 자료 다듬기 — 문구 순화 · 오타 · 레이아웃 · 부록. (2026-08-07)

사용자 지적
  ① 「TROUBLE 02 가 너무 직설적이다. 프로젝트 소개인데 '발 기' 이런 소리 할 수는 없잖아」
  ② 「'좋아 보이던 셋을 껐다' 가 정확히 무슨 소리야?」
       → 🔴 내가 TROUBLE 05(V AND C)를 지우면서 «셋 중 하나»를 없앴다.
         슬라이드가 「셋」이라 써놓고 «둘»만 보여주고 있었다. 내 실수다.
  ③ 「23·24 에 테스트 결과를 붙였으면. 공간 없으면 맨 뒤도 좋다」

무엇을 하나
  ① TROUBLE 02 예시 교체 — 노골적인 말을 «우리 사용자가 실제로 하는 말»로
       ✗ 「웹 개발 기초」의 '발기'
       ○ 「저 혼자 해요」→ 혼[자 해]요 (자해로 차단) · 「할머니를 돌봤어요」→ 할[머니]
       word_boundary_test.py 가 실제로 검사하는 사례들이다. 더 깨끗하고 «더 세다» —
       가족돌봄청년이 실제로 쓰는 문장이 막혔다는 뜻이라 TROUBLE 01 과도 이어진다.
  ② TROUBLE 04 에 «세 번째»를 되살린다 — 두 신호 AND(14/14 만점을 되돌린 것).
       ⚠ 논문(Nature)은 «안» 쓴다. 이야기는 논문 없이도 완결된다.
  ③ 오타 「미래설계지묘」 → 「미래설계지도」
  ④ 「발표자 ○○」 상자가 좁아 «두 줄»로 깨지던 것 — 폭을 넓힌다 (덱 전체 31곳)
  ⑤ 부록 신설 — 오늘 실제로 돌린 검사 결과 (LLM 0회 · 0원)

쓰는 법
  python polish_ppt.py <입력.pptx> <출력.pptx>
"""
import sys
import io
import copy

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu                                        # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v4.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v5.pptx'

TPL_A = 20      # 3카드 + 우패널 (잇다 시스템 구조) — 부록 본보기
INCH = 914400


def clone(prs, idx0):
    src = prs.slides[idx0]
    dst = prs.slides.add_slide(src.slide_layout)
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(sh._element))
    for k, rel in src.part.rels.items():
        if rel.reltype.endswith('slideLayout'):
            continue
        #  🔴 notesSlide 를 복사하면 PowerPoint 가 파일을 «거부»한다 (build_ppt_team.py 참고)
        if rel.reltype.endswith('notesSlide'):
            continue
        if rel.is_external:
            dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            dst.part.rels.get_or_add(rel.reltype, rel._target)
    return dst


def put(shape, text):
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    p0 = tf.paragraphs[0]
    if not p0.runs:
        return
    lines = str(text).split('\n')
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    p0.runs[0].text = lines[0]
    for p in list(tf.paragraphs[1:]):
        p._p.getparent().remove(p._p)
    for ln in lines[1:]:
        newp = copy.deepcopy(p0._p)
        p0._p.getparent().append(newp)
        from pptx.text.text import _Paragraph
        np = _Paragraph(newp, tf)
        for r in list(np.runs[1:]):
            r._r.getparent().remove(r._r)
        np.runs[0].text = ln


def S(prs, n1):
    return list(prs.slides[n1 - 1].shapes)


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    print(f'  원본 {n0}장')

    # ══════════════════════════════════════════════════════════
    #  ① TROUBLE 02 — 노골적인 예시를 «우리 사용자의 말»로
    #     word_boundary_test.py 가 실제로 검사하는 사례에서 골랐다.
    # ══════════════════════════════════════════════════════════
    b = S(prs, 25)
    put(b[14], "원인 — 「저 혼자 해요」가 혼[자 해]요 로 갈려 «자해»로 차단됐다. "
               "「어려서부터 할머니를 돌봤어요」는 할[머니] 로 걸렸다. "
               "공백을 지운 문자열에서 낱말을 찾고 있었다.")
    put(b[15], "해결 — 낱말이 아니라 «매칭 방식»을 고쳤다. 우회는 어절 «사이»를 쪼개고 "
               "오탐은 어절 «가운데»서 붙는다 → 어절이 «시작»하는 자리에서만 인정한다.")
    put(b[16], "배운 것 — 낱말표를 한 글자도 안 건드렸는데 오차단 16건이 사라졌다. "
               "「이번 건 고쳤다」와 「이 종류가 안 나게 했다」는 다르다.")

    # ══════════════════════════════════════════════════════════
    #  ② TROUBLE 04 — «셋»이라 써놓고 둘만 있던 것. 세 번째를 되살린다.
    #     ⚠ 논문 없이 쓴다. 「만점을 받고도 정답표를 의심했다」가 이야기의 전부다.
    # ══════════════════════════════════════════════════════════
    c = S(prs, 26)
    put(c[14], "가중 RRF — 검색 두 갈래에 «무게»를 주려 했다. 순위가 «안 바뀌었다». "
               "k=60 이면 1/61 과 1/66 이 거의 같아서 RRF 는 사실상 «표 세기»다.")
    put(c[15], "자기일관성 N=3 — 같은 질문을 세 번 물어 다수결. "
               "흔들리던 1건이 «그대로»였고 비용만 3배(0.50 → 1.46원)였다.")
    put(c[16], "두 신호 AND — 14/14 «만점»을 받고도 되돌렸다. 정답표가 틀렸기 때문이다. "
               "「용접이요」를 카드로 라벨했는데, CO₂·로봇·피복아크 중 뭐냐고 «물어보는 게» 맞았다.")

    # ══════════════════════════════════════════════════════════
    #  ③ 오타 — 「미래설계지묘」
    # ══════════════════════════════════════════════════════════
    fixed = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame and '미래설계지묘' in sh.text_frame.text:
                put(sh, sh.text_frame.text.replace('미래설계지묘', '미래설계지도'))
                fixed += 1
    print(f'  오타 「미래설계지묘」 {fixed}곳 수정')

    # ══════════════════════════════════════════════════════════
    #  ④ 「발표자 ○○」이 두 줄로 깨지던 것 — 폭 1.35" 에 6글자였다.
    #     슬라이드 폭 20" · 오른끝 19.08" 라 오른쪽 여백 0.92" 가 남는다 → 넓힌다.
    #     ⚠ 발표 «내내» 모든 슬라이드에 보이는 자리다. 원본부터 깨져 있었다.
    # ══════════════════════════════════════════════════════════
    widened = 0
    for s in prs.slides:
        for sh in s.shapes:
            if not (sh.has_text_frame and sh.text_frame.text.strip().startswith('발표자')):
                continue
            if sh.width >= Emu(int(2.0 * INCH)):     # 표지의 넓은 상자는 건드리지 않는다
                continue
            sh.width = Emu(int(2.15 * INCH))
            widened += 1
    print(f'  「발표자」 상자 {widened}곳 폭 1.35" → 2.15"')

    # ══════════════════════════════════════════════════════════
    #  ⑤ 부록 — 오늘(2026-08-07) 실제로 돌린 검사 결과
    #     ⚠ 전부 «내가 돌려서 본» 숫자다. 문서에서 옮겨 적은 게 아니다.
    # ══════════════════════════════════════════════════════════
    ap = clone(prs, TPL_A - 1)
    put(list(ap.shapes)[0], '부록 · 검사 결과')
    put(list(ap.shapes)[2], '돌려본 것 — 전부 LLM 0회 · 0원')
    put(list(ap.shapes)[4], '안전층 — 오차단이 없는가')
    put(list(ap.shapes)[5],
        '어절 정렬 22,470건 오차단 0 · 정상 발화 오탐 0/35 · 정규화 26/26. '
        '「유령 위치」 8곳은 어절 정렬이 막고 있다.')
    put(list(ap.shapes)[7], '대화 — 망가져도 사는가')
    put(list(ap.shapes)[8],
        'LLM 이 «완전히 죽은» 상태에서 대화 생존 11/11. 출력 가드 14/14. '
        '슬롯 빼기 14/14. 남용 문턱 9건 + 자동해제 10턴.')
    put(list(ap.shapes)[10], '문턱 — 어디에 둬도 되는 «폭»')
    put(list(ap.shapes)[11],
        '확신 문턱을 0.05~0.30 으로 훑어도 전부 14/14, 0.32 에서 무너졌다. '
        '정확도 하나가 아니라 「폭」을 본다.')
    put(list(ap.shapes)[13], '검사를 «생성»합니다')
    put(list(ap.shapes)[14],
        '22,470건은 손으로 쓴 게 아니다.\n'
        'DB 의 직업·자격증·강좌명 3,745건 ×\n'
        '문장틀 6개로 «만들어» 낸다.\n'
        '\n'
        '내가 고른 문장이 아니라서,\n'
        '내가 못 떠올린 사례도 걸린다.')
    print('  부록 1장 신설 — 검사 결과')

    prs.save(OUT)
    print(f'  저장 → {OUT}')
    print(f'  전체 {len(Presentation(OUT).slides)}장')


main()
