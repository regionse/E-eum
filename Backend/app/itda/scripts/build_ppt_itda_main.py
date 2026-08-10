# -*- coding: utf-8 -*-
"""잇다 발표 본편 — «백지에서» 짓는다. (2026-08-10)

왜 새로 짓나
  · 팀 덱(프레젠테이션 수정안 (4).pptx)의 잇다 구간을 갈아엎기로 했다.
    덜다 구간의 문법 — 제목 48pt 한 줄 · 흰 패널 + 24pt 본문 · 번호 원 · 그림 —
    이 «눈에 들어온다»는 피드백이라, 그 문법(색·크기·좌표를 실측해서)을 그대로 쓴다.
  · **슬라이드 복제는 0회.** build_ppt_team.py 가 notesSlide 관계 복사로
    PowerPoint 에서 안 열렸던 사고(verify_ppt.py 머리말) 때문에, 완전히 새 파일로 만들고
    김정준이 PowerPoint 「슬라이드 재사용」으로 팀 덱에 끼워 넣는다.

원칙 (발표는 «소개»다 — 메모리 presentation-is-an-introduction)
  · 지금 «쓰이는» 기능만 그린다. 오늘 끈 것(태그 가중·강점성향)은 안 그린다.
  · 용어는 «역할과 함께» — RRF(순위로 합침)·리랭커(다시 줄 세움)처럼.
  · 숫자는 전부 실측 — 직업 1,094 · 자격증 613 · 강좌 8,273 은 2026-08-10 DB COUNT.
    나머지는 실측기록_잇다.md 의 값(§1-1 골든셋 · §2-1 오탐 · §2-2 우회 · §3-1c 마진 ·
    §6-3 종류 95% · §9-4 턴당 · §13 200턴).

디자인 값 (수정안(4) 실측 — 2026-08-10)
  캔버스 20×11.25in · 제목 48pt #3B2063 · 눈썹 18.75pt #8B84A0 · 발표자 30pt #6B3FA0
  흰 패널 #FFFFFF/#E7E2F0 0.75pt · 배경 패널 #EFEBF6/#DED8EA · 번호 원 #F0E9FB
  본문 24pt #241F2E · 강조 #6B3FA0 · 붉은 강조 #A8364A

쓰는 법 (Backend/ 에서)
  python app/itda/scripts/build_ppt_itda_main.py
  → C:/Users/TFX255GS/Downloads/이음_잇다_발표.pptx
  이후 verify_ppt.py 로 «실제로 열어» 확인한다.
"""
import io
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r'C:\Users\TFX255GS\Downloads\이음_잇다_발표.pptx'

#  ── 팔레트 (수정안(4) 실측) ─────────────────────────────────────────
C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_EYE = RGBColor(0x8B, 0x84, 0xA0)
C_ACC = RGBColor(0x6B, 0x3F, 0xA0)
C_TXT = RGBColor(0x24, 0x1F, 0x2E)
C_RED = RGBColor(0xA8, 0x36, 0x4A)
C_DEEP = RGBColor(0x32, 0x16, 0x6F)
C_PANEL_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_BG = RGBColor(0xEF, 0xEB, 0xF6)
C_BG_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_CHIP = RGBColor(0xF0, 0xE9, 0xFB)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LINE = RGBColor(0xC9, 0xBE, 0xE0)
FONT = 'Malgun Gothic'


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _set_font(r, size, color, bold):
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    #  한글 폰트는 East-Asian 속성(ea)에도 넣어야 실제로 먹는다
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.18, anchor=MSO_ANCHOR.TOP):
    """runs = [줄, 줄, ...] · 줄 = [(글, size, color, bold), ...]"""
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        for (t, size, color, bold) in line:
            r = p.add_run()
            r.text = t
            _set_font(r, size, color, bold)
    return tb


def box(sl, x, y, w, h, fill=C_WHITE, line=C_PANEL_LN, lw=0.75, round_=True):
    shp = sl.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if round_ else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    if round_:
        try:
            shp.adjustments[0] = 0.06
        except Exception:                                  # noqa: BLE001
            pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def chip(sl, x, y, d, label, size=19.5):
    """번호 원 — 덜다의 그 원."""
    c = sl.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = C_CHIP
    c.line.fill.background()
    c.shadow.inherit = False
    tf = c.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = label
    _set_font(r, size, C_ACC, True)
    return c


def arrow_v(sl, x, y1, y2):
    ln = sl.shapes.add_connector(2, Inches(x), Inches(y1), Inches(x), Inches(y2))
    ln.line.color.rgb = C_LINE
    ln.line.width = Pt(2.25)
    _arrow_head(ln)
    return ln


def arrow_h(sl, x1, x2, y):
    ln = sl.shapes.add_connector(2, Inches(x1), Inches(y), Inches(x2), Inches(y))
    ln.line.color.rgb = C_LINE
    ln.line.width = Pt(2.25)
    _arrow_head(ln)
    return ln


def _arrow_head(ln):
    el = ln.line._get_or_add_ln()
    head = el.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'})
    el.append(head)


def head(sl, crumb, title):
    text(sl, 1.00, 0.86, 12.0, 0.45, [[(crumb, 18.75, C_EYE, False)]])
    text(sl, 14.0, 0.80, 5.0, 0.55, [[('발표자 김정준', 30, C_ACC, True)]],
         align=PP_ALIGN.RIGHT)
    text(sl, 1.00, 1.42, 18.1, 0.95, [[(title, 48, C_TITLE, True)]])


def flowbox(sl, x, y, w, h, title, sub=None, fill=C_WHITE, tcol=C_DEEP, tsz=22, ssz=17):
    """제목 + (선택) 부제 한 줄이 든 흐름 상자 — 세로 중앙 정렬."""
    box(sl, x, y, w, h, fill=fill)
    lines = [[(title, tsz, tcol, True)]]
    if sub:
        lines.append([(sub, ssz, C_EYE, False)])
    text(sl, x + 0.18, y, w - 0.36, h, lines, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE, ls=1.05)


# ══════════════════════════════════════════════════════════════════════
def s1_overview(prs):
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다', '잇다 — 대화로 진로를 찾는 상담 챗봇')

    #  왼쪽 — 원칙
    box(sl, 1.0, 2.65, 7.7, 8.1)
    text(sl, 1.45, 2.95, 6.9, 0.5, [[('원칙 하나', 30, C_TITLE, True)]])
    text(sl, 1.45, 3.62, 6.9, 1.7, [
        [('「사실은 DB에서만,', 34, C_ACC, True)],
        [('  대화는 LLM답게」', 34, C_ACC, True)],
    ], ls=1.2)
    text(sl, 1.45, 5.75, 6.95, 4.7, [
        [('화면에 뜨는 «사실»', 22, C_TXT, True)],
        [('직업 1,094 · 자격증 613 · 강좌 8,273', 21, C_ACC, True)],
        [('— 전부 DB 값 그대로. 모델이 지어낼 자리가 없다', 18, C_EYE, False)],
        [('', 10, C_TXT, False)],
        [('문장·공감·되묻기', 22, C_TXT, True)],
        [('LLM(Gemini)이 만든다 — 상담다운 말투는 모델의 몫', 18, C_EYE, False)],
        [('', 10, C_TXT, False)],
        [('그 사이를 «코드»가 지킨다', 22, C_TXT, True)],
        [('검증 · 판정 · 출력 가드 — 이번 장들이 그 이야기', 18, C_EYE, False)],
    ], ls=1.22)

    #  오른쪽 — 한 턴의 길 (세로 흐름)
    box(sl, 9.0, 2.65, 10.1, 8.1, fill=C_BG, line=C_BG_LN)
    text(sl, 9.45, 2.92, 9.2, 0.5, [[('한 턴이 지나가는 길', 30, C_TITLE, True)]])
    steps = [
        ('사용자 발화', '「할머니 돌보느라 학교를 그만뒀어…」', C_WHITE),
        ('①  게이트', '위험·장난을 거른다 — 낱말은 신호, 판정은 LLM', C_WHITE),
        ('②  슬롯 추출', '말에서 «칸»을 채운다 — 구조화 출력', C_WHITE),
        ('③  검색', '두 갈래로 찾고 순위로 합친다 — Hybrid+RRF', C_WHITE),
        ('④  착지 판정', '더 물을까, 결과를 낼까 — 코드가 정한다', C_WHITE),
        ('⑤  답', '문장 + 직업 카드(자격증·시험일·강좌)', C_WHITE),
    ]
    y = 3.62
    bh, gap = 1.02, 0.20
    for i, (t, s, f) in enumerate(steps):
        box(sl, 9.5, y, 9.1, bh, fill=f)
        text(sl, 9.78, y + 0.10, 3.3, 0.8, [[(t, 23, C_DEEP, True)]],
             anchor=MSO_ANCHOR.MIDDLE)
        text(sl, 13.0, y + 0.10, 5.5, 0.86, [[(s, 16.5, C_EYE, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            arrow_v(sl, 10.7, y + bh, y + bh + gap)
        y += bh + gap


def s2_gate(prs):
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다 — 게이트', '게이트 — 낱말은 «신호», 판정은 LLM이')

    #  왼쪽 — 2층 사다리
    box(sl, 1.0, 2.65, 10.9, 8.1)
    text(sl, 1.45, 2.92, 10.0, 0.5, [[('두 층으로 거른다', 30, C_TITLE, True)]])

    flowbox(sl, 1.6, 3.70, 9.7, 1.30, '1층 — 코드 검사 (LLM 0회 · 0원)',
            '침묵·오타·자모 축약 구별 / 우회 표기 정규화 3종 (씨1발·필ㄹㅗ폰·가s슴)')
    arrow_v(sl, 6.45, 5.00, 5.32)
    text(sl, 6.7, 4.96, 4.3, 0.4, [[('위험 낱말이 «걸리면»', 16.5, C_EYE, False)]])
    flowbox(sl, 1.6, 5.32, 9.7, 1.30, '2층 — LLM 판정 「누구 이야기인가」',
            '본인 / 제3자 / 예방·직업 관심 / 불명 — 문맥을 읽는다')

    y = 7.0
    for t, s, col in [
        ('본인 위기  →  대화를 끊지 않는다', '전문 연락처(109)를 안내하고 «상담을 잇는다»', C_RED),
        ('제3자·예방·직업 관심  →  그대로 진행', '「상담사가 되고 싶어요」를 낱말로 막으면 안 되니까', C_DEEP),
        ('차단은 최후의 수단', '상대를 향한 욕설·성적 발화 · 반복 남용만 (턴 비례 문턱)', C_DEEP),
    ]:
        box(sl, 1.6, y, 9.7, 1.06)
        text(sl, 1.9, y + 0.08, 9.2, 0.5, [[(t, 20, col, True)]])
        text(sl, 1.9, y + 0.55, 9.2, 0.45, [[(s, 16.5, C_EYE, False)]])
        y += 1.22

    #  오른쪽 — 실측 3개
    x = 12.2
    for (big, cap1, cap2), yy in zip([
        ('0 / 35', '정상 발화를 잘못 막은 횟수', '「가슴이 답답해요」·「마약수사관이 되고 싶어요」 전부 통과'),
        ('29 / 32', '우회 표기 차단 (실측)', '남은 3건도 지어내지 않고 한계로 기록해 뒀다'),
        ('대화 유지', '위기 문장 뒤에도', '연락처만 던지고 끝내지 않는다 — 지속적 관여'),
    ], (2.65, 5.42, 8.19)):
        box(sl, x, yy, 6.9, 2.56)
        text(sl, x + 0.4, yy + 0.28, 6.1, 0.9, [[(big, 44, C_ACC, True)]])
        text(sl, x + 0.4, yy + 1.28, 6.1, 0.5, [[(cap1, 20, C_TXT, True)]])
        text(sl, x + 0.4, yy + 1.80, 6.1, 0.65, [[(cap2, 15.5, C_EYE, False)]])


def s3_slots(prs):
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다 — 슬롯', '슬롯 추출 — 모델은 «칸»에만 채운다 (구조화 출력)')

    #  왼쪽 — 칸 목록
    box(sl, 1.0, 2.65, 8.7, 8.1)
    text(sl, 1.45, 2.92, 7.9, 0.5, [[('칸 6개', 30, C_TITLE, True),
                                     ('   — 매 턴, 이번 발화에서 새로 안 것만', 17, C_EYE, False)]])
    rows = [
        ('관심분야', '좋아한다·해봤다고 말한 것 — 일상어 그대로', False),
        ('활동유형', '만들기·돕기돌봄·정비… — 9종 목록에서만', True),
        ('다루는대상', '사람·기계·데이터… — 6종 목록에서만', True),
        ('세부관심', '콕 집은 것 — 어르신 / 웹 / 제빵 / 병원', False),
        ('제약', '시간부족·비용부담·체력부담… 5종', True),
        ('관심대분류', 'NCS 24갈래 — 검색 범위를 넓게 가른다', True),
    ]
    y = 3.66
    for name, desc, enum in rows:
        box(sl, 1.45, y, 7.8, 0.98)
        text(sl, 1.75, y + 0.06, 3.0, 0.8, [[(name, 21, C_DEEP, True)]],
             anchor=MSO_ANCHOR.MIDDLE)
        text(sl, 4.45, y + 0.06, 4.7, 0.86, [[(desc, 15.5, C_EYE, False)]],
             anchor=MSO_ANCHOR.MIDDLE)
        if enum:
            text(sl, 7.85, y + 0.05, 1.3, 0.4, [[('갇힌 출력', 12.5, C_RED, True)]],
                 align=PP_ALIGN.RIGHT)
        y += 1.10
    text(sl, 1.45, y + 0.02, 7.9, 0.5,
         [[('「갇힌 출력」 = 목록(enum) 밖의 값은 출력 자체가 불가능', 16.5, C_RED, False)]])

    #  오른쪽 위 — 근거 대조
    box(sl, 10.0, 2.65, 9.1, 3.62)
    text(sl, 10.45, 2.93, 8.3, 0.5, [[('근거 없으면 «버린다»', 28, C_TITLE, True)]])
    text(sl, 10.45, 3.62, 8.3, 2.5, [
        [('모든 값에 사용자가 «실제로 한 말»을 근거로 함께 받는다', 19, C_TXT, True)],
        [('→ 코드가 발화와 글자로 대조 — 발화에 없으면 그 값은 폐기', 19, C_TXT, False)],
        [('', 8, C_TXT, False)],
        [('모델이 «지어낸 슬롯»이 프로필에 들어올 길을 막는다', 17, C_EYE, False)],
    ], ls=1.3)

    #  오른쪽 아래 — 종류
    box(sl, 10.0, 6.50, 9.1, 4.25)
    text(sl, 10.45, 6.78, 8.3, 0.5, [[('「종류」 — 사실과 선호를 가른다', 28, C_TITLE, True)]])
    text(sl, 10.45, 7.48, 8.3, 3.1, [
        [('원함 ', 21, C_ACC, True), ('· ', 21, C_EYE, False),
         ('해봤음 ', 21, C_DEEP, True), ('· ', 21, C_EYE, False),
         ('못함', 21, C_RED, True),
         ('   — 같은 값에도 처지가 다르다', 17, C_EYE, False)],
        [('', 8, C_TXT, False)],
        [('「돌봄은 힘든데 그래도 사람 만나는 건 좋아요」', 18.5, C_TXT, True)],
        [('→ 사람 만나기 = 원함  ·  돌봄 = 못함   (한 문장에서 갈라냄)', 18.5, C_TXT, False)],
        [('', 8, C_TXT, False)],
        [('「못함」은 검색어에서 뺀다 — 힘들다는 쪽으로 끌고 가지 않게', 17, C_EYE, False)],
        [('100케이스 실측 95%', 17, C_ACC, True)],
    ], ls=1.28)


def s4_search(prs):
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다 — 검색', '검색 — 두 갈래로 찾고, «순위»로 합친다')

    box(sl, 1.0, 2.65, 18.08, 8.1, fill=C_BG, line=C_BG_LN)

    #  맨 위 — 검색어
    flowbox(sl, 5.54, 2.95, 9.0, 1.06, '검색어 = 이번 발화 + 지금까지 쌓인 슬롯(닻)',
            '대화가 딴 데로 흘러도 «들은 방향»이 항상 함께 들어간다')
    arrow_v(sl, 7.2, 4.01, 4.45)
    arrow_v(sl, 12.9, 4.01, 4.45)

    #  두 갈래
    flowbox(sl, 1.7, 4.45, 8.0, 1.46, '벡터 검색 — «뜻»이 비슷한 것',
            'gemini-embedding 3,072차원 · Pinecone')
    text(sl, 1.7, 5.97, 8.0, 0.44,
         [[('뭉뚱그린 말에 강하다 — 「손으로 만드는 일이요」', 16.5, C_EYE, False)]],
         align=PP_ALIGN.CENTER)
    flowbox(sl, 10.4, 4.45, 8.0, 1.46, 'FULLTEXT — «글자»가 맞는 것',
            'MySQL ngram 전문 검색')
    text(sl, 10.4, 5.97, 8.0, 0.44,
         [[('콕 집은 말에 강하다 — 「용접이요」', 16.5, C_EYE, False)]],
         align=PP_ALIGN.CENTER)

    arrow_v(sl, 5.7, 6.45, 6.85)
    arrow_v(sl, 14.4, 6.45, 6.85)

    #  RRF → 리랭커 → 후보
    flowbox(sl, 3.1, 6.85, 13.9, 1.18, 'RRF — 점수가 아니라 «순위»로 합친다',
            '벡터 점수와 글자 점수는 기준이 달라 «직접 더할 수 없다» → 등수로 투표')
    arrow_v(sl, 10.05, 8.03, 8.43)
    flowbox(sl, 3.1, 8.43, 13.9, 1.06, '리랭커 — 대화 문장 그대로 다시 줄 세운다 (Jina)',
            '검색어로 찾고, «원문»으로 고쳐 앉힌다')
    arrow_v(sl, 10.05, 9.49, 9.89)
    flowbox(sl, 6.9, 9.89, 6.3, 0.72, '후보 직업 (상위 N)', None, fill=C_CHIP)


def s5_landing(prs):
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다 — 판정', '착지 판정 — 언제 결과를 낼지는 «코드»가 정한다')

    #  왼쪽 — 2단 트리
    box(sl, 1.0, 2.65, 10.9, 8.1)
    text(sl, 1.45, 2.92, 10.0, 0.5, [[('두 단으로 나눠 본다', 30, C_TITLE, True),
                                      ('   — 후보 점수의 «분포»를 읽는다', 17, C_EYE, False)]])

    #  본류는 세로(1단 → 2단 → 카드), 가지(칩)는 오른쪽으로 뺀다 —
    #  처음 판은 가지를 본류 위에 얹어서 «겹쳤다» (PNG 육안 검수에서 잡음, 2026-08-10)
    flowbox(sl, 1.5, 3.70, 6.1, 1.02, '1단 — 1위 «동네»가 확실한가', 'TOP% ≥ 0.90')
    arrow_h(sl, 7.6, 8.15, 4.21)
    text(sl, 7.55, 3.80, 1.4, 0.35, [[('아니오', 14, C_RED, True)]])
    flowbox(sl, 8.2, 3.46, 3.25, 1.5, '방향 칩', '「돌봄 / 만들기 / 사무\n— 어느 쪽?」',
            fill=C_CHIP, tsz=20, ssz=14.5)
    arrow_v(sl, 4.55, 4.72, 5.35)
    text(sl, 4.78, 4.83, 1.2, 0.35, [[('예', 14, C_ACC, True)]])
    flowbox(sl, 1.5, 5.35, 6.1, 1.02, '2단 — 그 동네 «안»이 갈리나', '이웃 점수 엔트로피 ≥ 0.70')
    arrow_h(sl, 7.6, 8.15, 5.86)
    text(sl, 7.72, 5.45, 1.2, 0.35, [[('예', 14, C_ACC, True)]])
    flowbox(sl, 8.2, 5.11, 3.25, 1.5, '세부 칩', '「CO₂ / 로봇 / 피복 용접\n— 어느 것?」',
            fill=C_CHIP, tsz=20, ssz=14.5)
    arrow_v(sl, 4.55, 6.37, 7.00)
    text(sl, 4.78, 6.52, 1.6, 0.35, [[('아니오', 14, C_DEEP, True)]])
    flowbox(sl, 2.8, 7.00, 3.5, 1.25, '카드', '직업 + 자격증 + 시험일', fill=C_WHITE,
            tcol=C_ACC, tsz=24, ssz=15.5)

    text(sl, 1.45, 8.85, 10.0, 1.4, [
        [('예 — ', 17, C_EYE, False), ('「빵 만들고 싶어요」→ 카드', 18, C_TXT, True),
         ('  ·  ', 17, C_EYE, False), ('「용접이요」→ 세부 칩', 18, C_TXT, True),
         ('  ·  ', 17, C_EYE, False), ('「컴퓨터요」→ 방향 칩', 18, C_TXT, True)],
        [('되묻기는 두 종류라서 한 단으로는 뭉개진다 — «어느 동네»와 «동네 어디»', 16.5, C_EYE, False)],
    ], ls=1.35)

    #  오른쪽 — 왜 튼튼한가
    box(sl, 12.2, 2.65, 6.9, 8.1)
    text(sl, 12.6, 2.93, 6.1, 0.5, [[('문턱이 튼튼한가', 28, C_TITLE, True)]])
    text(sl, 12.6, 3.66, 6.1, 1.3, [
        [('안전 구간 폭 ', 20, C_TXT, True), ('0.317', 30, C_ACC, True)],
        [('단일 신호(0.034)의 «10배»', 18, C_TXT, False)],
    ], ls=1.25)
    text(sl, 12.6, 5.15, 6.1, 1.2, [
        [('문턱이 조금 틀어져도 판정이', 17, C_EYE, False)],
        [('뒤집히지 않는 폭 — 14케이스 실측', 17, C_EYE, False)],
    ], ls=1.25)
    box(sl, 12.6, 6.45, 6.1, 4.0, fill=C_BG, line=C_BG_LN)
    text(sl, 12.9, 6.65, 5.6, 3.6, [
        [('모델의 자기보고를 안 믿는다', 19, C_DEEP, True)],
        [('「충분히 알았다」는 모델의 말이', 16.5, C_TXT, False)],
        [('아니라, 후보 점수의 분포라는', 16.5, C_TXT, False)],
        [('«숫자»로 판단한다.', 16.5, C_TXT, False)],
        [('', 8, C_TXT, False)],
        [('일찍 단정하는 버릇(성급한 착지)을', 16.5, C_TXT, False)],
        [('코드가 눌러 준다.', 16.5, C_TXT, False)],
    ], ls=1.3)


def s6_answer(prs):
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다 — 답 내기', '답 내기 — 문장은 LLM, 사실은 코드, 사이에 가드')

    cols = [
        ('①  하나를 «고른다»', [
            ('후보 여러 개 중 «대화 전체»에 비춰 하나를 고르고', 18, C_TXT, False),
            ('「왜」를 사용자의 말을 인용해 적는다', 18, C_TXT, False),
            ('', 9, C_TXT, False),
            ('이미 답한 걸 다시 묻지 않는다 — 후보를 가르는', 16, C_EYE, False),
            ('말이 대화에 있으면 그걸 쓴다', 16, C_EYE, False),
            ('', 9, C_TXT, False),
            ('정보가 모자라면 카드 대신 «확인 질문» 하나', 17, C_ACC, True),
        ]),
        ('②  사실은 «코드»가 쓴다', [
            ('시험 일정 — 오늘 이후 가장 가까운 회차 + D-day', 18, C_TXT, False),
            ('「목록에 없어요」 — 없다는 사실도 코드가 말한다', 18, C_TXT, False),
            ('자격증·강좌 — DB 값 그대로', 18, C_TXT, False),
            ('', 9, C_TXT, False),
            ('모델에게 맡기면 없는 자격증을 있다고 한다 —', 16, C_EYE, False),
            ('그래서 «사실 문장»은 모델의 손을 안 거친다', 16, C_EYE, False),
        ]),
        ('③  나가는 말을 «거른다»', [
            ('마지막 출구에서 문장 단위로 뗀다:', 18, C_TXT, False),
            ('', 6, C_TXT, False),
            ('· 응시요건 단정  「누구나 딸 수 있어요」', 17, C_TXT, False),
            ('· 근거 없는 속성  「바로 시작할 수 있는」', 17, C_TXT, False),
            ('· 카드에 없는 숫자  (지어낸 날짜·회차)', 17, C_TXT, False),
            ('· 금지된 되묻기  「어떤 일을 하고 싶으세요?」', 17, C_TXT, False),
            ('', 9, C_TXT, False),
            ('다 지워지면 안전 문구로 대체 — 침묵하지 않는다', 16, C_EYE, False),
        ]),
    ]
    x = 1.0
    for title, lines in cols:
        box(sl, x, 2.65, 5.9, 8.1)
        text(sl, x + 0.35, 2.95, 5.25, 0.9, [[(title, 25, C_TITLE, True)]])
        text(sl, x + 0.35, 3.85, 5.25, 6.6,
             [[(t, s, c, b)] for (t, s, c, b) in lines], ls=1.32)
        x += 6.09


def s7_measured(prs):
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다 — 잰 것', '잰 것 — 화면의 숫자는 전부 실측입니다')

    cells = [
        ('34 / 34', '골든셋 회귀', '기대 답을 정해 둔 34문항 — 고칠 때마다 돌린다'),
        ('1.8턴', '첫 결과 카드까지 평균', '콕 집어 말하면 1턴, 뭉뚱그리면 되묻고 2~3턴'),
        ('1.2~1.6초', '보통 턴 · 약 1.3원', 'LLM 1회 — 게이트는 의심스러울 때만 돈다'),
        ('0 / 35', '정상 발화 오차단', '위험 낱말이 «든» 정상 질문 35개 전부 통과'),
        ('200턴 × 2', '장기 대화 완주', '성격이 반대인 페르소나 둘 — 끝까지 밀어봤다'),
        ('재시작 생존', '세션이 DB에 산다', '서버가 죽어도 대화·프로필이 이어진다'),
    ]
    xs, ys = (1.0, 7.09, 13.18), (2.65, 6.20)
    for i, (big, cap, sub) in enumerate(cells):
        x, y = xs[i % 3], ys[i // 3]
        box(sl, x, y, 5.82, 3.30)
        text(sl, x + 0.4, y + 0.30, 5.1, 1.1, [[(big, 46, C_ACC, True)]])
        text(sl, x + 0.4, y + 1.55, 5.1, 0.6, [[(cap, 22, C_TXT, True)]])
        text(sl, x + 0.4, y + 2.20, 5.1, 0.95, [[(sub, 15.5, C_EYE, False)]], ls=1.22)

    text(sl, 1.0, 9.80, 18.0, 0.9, [
        [('검사 스크립트 36개 · 실측 기록 약 1,900줄', 19, C_DEEP, True),
         ('  — 대부분 LLM 없이(0원) 다시 돌릴 수 있게 만들었습니다', 17, C_EYE, False)],
    ])


def s8_next(prs):
    sl = blank(prs)
    head(sl, 'PART 03 · 잇다 — 다음 걸음', '여기서 더 간다면 — 이미 «재 둔» 방향들')

    text(sl, 1.0, 2.55, 18.0, 0.5,
         [[('넷 다 측정까지 끝내 둔 방향이라, 근거를 갖고 답할 수 있습니다', 19, C_EYE, False)]])

    cells = [
        ('대화 «접기»', ['수십 턴을 넘어가면 오래된 대화를 요약으로 접는 계층화.',
                       '200턴 시험이 남긴 숙제 — 슬롯·이력이 쌓이기만 하는 것.']),
        ('프롬프트 갈라 쓰기', ['상황(첫 턴·되묻기·카드)별로 나누면 호출당 토큰이 절반.',
                          '경계 사고가 없는지 상황별 재검증이 필요해 남겨 뒀다.']),
        ('합치기 상수(k) 조정', ['질의를 여러 개 합칠 때 더 좋은 값을 100건 실측으로 확보.',
                           '순위 산식의 심장이라, 회귀 전부를 다시 돌리고 바꾼다.']),
        ('경험만 말하는 사용자', ['「해봤다」밖에 말 못 하는 사용자도 결과에 닿도록',
                           '착지 무게를 조정 — 우리 주 사용자의 실제 모습이다.']),
    ]
    xs, ys = (1.0, 10.05), (3.30, 7.10)
    for i, (t, lines) in enumerate(cells):
        x, y = xs[i % 2], ys[i // 2]
        box(sl, x, y, 9.0, 3.55)
        chip(sl, x + 0.42, y + 0.42, 0.62, str(i + 1))
        text(sl, x + 1.30, y + 0.40, 7.4, 0.6, [[(t, 26, C_TITLE, True)]])
        text(sl, x + 1.30, y + 1.25, 7.35, 2.1,
             [[(ln, 17.5, C_TXT if j == 0 else C_EYE, j == 0)]
              for j, ln in enumerate(lines)], ls=1.3)


def main():
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    s1_overview(prs)
    s2_gate(prs)
    s3_slots(prs)
    s4_search(prs)
    s5_landing(prs)
    s6_answer(prs)
    s7_measured(prs)
    s8_next(prs)
    prs.save(OUT)
    print(f'저장: {OUT} · {len(prs.slides.__iter__.__self__._sldIdLst)}장')


if __name__ == '__main__':
    main()
