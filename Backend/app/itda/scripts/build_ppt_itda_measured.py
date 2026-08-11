# -*- coding: utf-8 -*-
"""잇다 발표 — 「왜 이 층들인가」 1장. (2026-08-10)

왜 만드나
  발표에서 나올 두 질문에 «숫자»로 답하려는 것이다.
    ① 「k=60 은 왜죠?」            ② 「엔트로피까지 왜 필요해요? 리랭커면 되잖아요」
  둘 다 어제까지는 답이 「표준값이라서」·「그게 나아 보여서」뿐이었다.

숫자 출처 — 전부 오늘 실측
  checks/gate_ablation.py · 정답표 100건 · 원값 checks/_gate_ablation.json
  실측기록 §18 에 표와 한계를 그대로 적어 뒀다.

⚠ 자기완결로 만든다. build_ppt_itda_main.py 를 다른 세션이 편집 중이라
  같은 파일을 건드리면 서로 덮어쓴다(오늘 GATE_MODEL 이 실제로 그렇게 날아갔다).

쓰는 법 (Backend/ 에서)
  python app/itda/scripts/build_ppt_itda_why.py
  → C:/Users/TFX255GS/Downloads/이음_잇다_왜.pptx   이후 verify_ppt.py 로 «열어» 확인
"""
import io
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = r'C:\Users\TFX255GS\Downloads\이음_잇다_잰것_수정.pptx'

C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_EYE = RGBColor(0x8B, 0x84, 0xA0)
C_ACC = RGBColor(0x6B, 0x3F, 0xA0)
C_TXT = RGBColor(0x24, 0x1F, 0x2E)
C_RED = RGBColor(0xA8, 0x36, 0x4A)
C_DEEP = RGBColor(0x32, 0x16, 0x6F)
C_PANEL_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_BG = RGBColor(0xEF, 0xEB, 0xF6)
C_BG_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = 'Malgun Gothic'


def _font(r, size, color, bold):
    r.font.name = FONT
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    rPr = r._r.get_or_add_rPr()
    ea = rPr.find(qn('a:ea'))
    if ea is None:
        ea = rPr.makeelement(qn('a:ea'), {})
        rPr.append(ea)
    ea.set('typeface', FONT)


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.18, anchor=MSO_ANCHOR.TOP):
    tb = sl.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        for (t, size, color, bold) in line:
            r = p.add_run()
            r.text = t
            _font(r, size, color, bold)
    return tb


def box(sl, x, y, w, h, fill=C_WHITE, line=C_PANEL_LN, lw=0.75):
    shp = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                              Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = 0.06
    except Exception:                                   # noqa: BLE001
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    shp.line.color.rgb = line
    shp.line.width = Pt(lw)
    shp.shadow.inherit = False
    return shp


def row(sl, x, y, cols, widths, size=17.5, bold=False, color=C_TXT):
    """표 한 줄 — 열 시작 x 를 누적해서 찍는다."""
    cx = x
    for c, w in zip(cols, widths):
        text(sl, cx, y, w, 0.36, [[(c, size, color, bold)]], align=PP_ALIGN.RIGHT)
        cx += w



def main():
    prs = Presentation()
    prs.slide_width = Inches(20)
    prs.slide_height = Inches(11.25)
    sl = prs.slides.add_slide(prs.slide_layouts[6])
    text(sl, 1.00, 0.86, 12.0, 0.45, [[('PART 03 · 잇다 — 잰 것', 18.75, C_EYE, False)]])
    text(sl, 14.0, 0.80, 5.0, 0.55, [[('발표자 김정준', 30, C_ACC, True)]], align=PP_ALIGN.RIGHT)
    text(sl, 1.00, 1.42, 18.1, 0.95,
         [[('잰 것 — 화면의 숫자는 전부 실측입니다', 48, C_TITLE, True)]])

    #  ★ 2026-08-10 — 「1.8턴」·「턴당 2.8원」을 고친 판.
    #    1.8턴: 실측기록 §12 가 「협조적인 대본을 내가 짰기 때문」이라고 스스로 부정한 숫자.
    #           동적 대화는 §12-2 6턴 · §13-1 첫 카드 4턴 ⇒ 「3~4턴(동적)」으로.
    #    2.8원: 게이트를 lite 로 되돌린 뒤 5턴 재측정 = 2.53원 ⇒ 「약 2.5원」으로.
    cells = [
        ('34 / 34', '골든셋 회귀', '기대 답을 정해 둔 34문항 — 고칠 때마다 돌린다'),
        ('3~4턴', '첫 결과 카드까지 (동적 대화)', '콕 집어 말하면 1턴 · 막연하면 되묻고 4~6턴'),
        ('턴당 약 2.5원', '보통 턴 · 1.2~1.6초', 'LLM 2회 (본문 + 안전 게이트) · 카드 턴은 4회'),
        ('0 / 35', '정상 발화 오차단', '위험 낱말이 «든» 정상 질문 35개 전부 통과'),
        ('200턴 × 2', '장기 대화 완주', '성격이 반대인 페르소나 둘 — 끝까지 밀어봤다'),
        ('재시작 생존', '세션이 DB에 산다', '서버가 죽어도 대화·프로필이 이어진다'),
    ]
    xs, ys = (1.0, 7.09, 13.18), (2.65, 6.20)
    for i, (big, cap, sub) in enumerate(cells):
        x, y = xs[i % 3], ys[i // 3]
        box(sl, x, y, 5.82, 3.30)
        text(sl, x + 0.4, y + 0.30, 5.1, 1.1, [[(big, 44, C_ACC, True)]])
        text(sl, x + 0.4, y + 1.55, 5.1, 0.6, [[(cap, 21, C_TXT, True)]])
        text(sl, x + 0.4, y + 2.20, 5.1, 0.95, [[(sub, 15.5, C_EYE, False)]], ls=1.22)

    text(sl, 1.0, 9.80, 18.0, 0.9, [
        [('검사 스크립트 36개 · 실측 기록 2,270줄', 19, C_DEEP, True),
         ('  — 대부분 LLM 없이(0원) 다시 돌릴 수 있게 만들었습니다', 17, C_EYE, False)],
    ])
    prs.save(OUT)
    print(f'저장: {OUT} · 1장')


if __name__ == '__main__':
    main()
