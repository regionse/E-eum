# -*- coding: utf-8 -*-
"""코드 슬라이드 2장 — 「이렇게 짰습니다」를 코드로 보여준다. (2026-08-07)

사용자 요청
  「이제 코드를 복붙하고 싶어. 정확히는 사진으로서. 덜다측에 코드 보이지? 저런 식으로.
   나도 «코드로서 이렇게 제안을 해뒀다» 라고 하고 싶어.
     1. 슬롯에 넣는 7개의 칸으로 바꿉니다 ← 코드로
     2. 관문 6개                      ← 코드로
     3. 각각 PPTX 슬라이드 1개씩 더 쓸 거임」

⚠ 코드는 itda_core.py 에서 «그대로» 가져왔다. 발표용으로 줄인 곳은 `…` 로 표시했고,
  줄인 것 말고는 한 글자도 안 바꿨다. 지어낸 줄이 없다.
    · _slot() / PROFILE_SCHEMA / verify_slots()  →  87~340행
    · _step() 의 관문 ①~⑥                        →  5343 · 5371 · 5408 · 5426 · 5430 · 5560행

만드는 법
  글자를 «도형»으로 넣는다(그림 파일이 아니다) → PowerPoint 에서 바로 고칠 수 있고,
  PNG 로 뽑으면 사진으로도 쓸 수 있다.

쓰는 법
  python make_code_slides.py [바탕.pptx] [출력.pptx]
"""
import sys
import io
import re

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE                           # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v9.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v10.pptx'

IN = 914400
FONT = 'Arial'
CODE_FONT = 'Consolas'

C_HEAD    = RGBColor(0x8B, 0x84, 0xA0)
C_PRESENT = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE   = RGBColor(0x3B, 0x20, 0x63)
C_MUTED   = RGBColor(0x8B, 0x84, 0xA0)
C_CODE_BG = RGBColor(0xF7, 0xF5, 0xFB)
C_CODE_LN = RGBColor(0xDE, 0xD8, 0xEA)

#  코드 색 — 배경이 밝으므로 대비를 충분히 준다
C_TXT = RGBColor(0x33, 0x30, 0x40)      # 보통 글자
C_CMT = RGBColor(0x8E, 0x88, 0x9E)      # 주석
C_KWD = RGBColor(0x8A, 0x2B, 0xE2)      # 키워드
C_STR = RGBColor(0x1F, 0x7A, 0x50)      # 문자열
C_NUM = RGBColor(0xC2, 0x6B, 0x2E)      # 숫자
C_HOT = RGBColor(0xA8, 0x36, 0x4A)      # ★ 강조 주석

KEYWORDS = {'def', 'return', 'if', 'else', 'elif', 'for', 'in', 'not', 'or', 'and',
            'None', 'True', 'False', 'import', 'from', 'class', 'async', 'await',
            'try', 'except', 'continue', 'break', 'lambda', 'is'}
TOKEN = re.compile(r"('[^']*'|\"[^\"]*\"|\bf'[^']*'|[A-Za-z_][A-Za-z_0-9]*|\d+|\s+|.)")


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.2):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.color.rgb = col
    return tb


def code_block(sl, x, y, w, h, lines, pt=11.5):
    """코드를 «색칠해서» 넣는다. 한 줄이 한 문단, 토큰마다 run 을 나눈다."""
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = C_CODE_BG
    sh.line.color.rgb = C_CODE_LN
    sh.line.width = Pt(0.75)
    sh.shadow.inherit = False
    sh.text_frame.text = ''
    try:
        sh.adjustments[0] = 0.02
    except Exception:
        pass

    tb = sl.shapes.add_textbox(Emu(int((x + 0.34) * IN)), Emu(int((y + 0.28) * IN)),
                               Emu(int((w - 0.68) * IN)), Emu(int((h - 0.56) * IN)))
    tf = tb.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = 1.28
        #  주석은 통째로 한 색 — ★ 가 있으면 강조색
        cut = ln.find('#')
        head, cmt = (ln, '') if cut < 0 else (ln[:cut], ln[cut:])
        for tok in TOKEN.findall(head):
            if not tok:
                continue
            if tok[0] in '\'"' or tok[:2] == "f'":
                col = C_STR
            elif tok in KEYWORDS:
                col = C_KWD
            elif tok.isdigit():
                col = C_NUM
            else:
                col = C_TXT
            r = p.add_run()
            r.text = tok
            r.font.name = CODE_FONT
            r.font.size = Pt(pt)
            r.font.color.rgb = col
        if cmt:
            r = p.add_run()
            r.text = cmt
            r.font.name = CODE_FONT
            r.font.size = Pt(pt)
            r.font.color.rgb = C_HOT if ('★' in cmt or '**' in cmt) else C_CMT
    return sh


def header(s, sub, title):
    text(s, 1.00, 0.92, 9.0, 0.30, [(sub, 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRESENT)])
    text(s, 1.00, 1.50, 19.0, 0.95, [(title, 42, True, C_TITLE)])


# ════════════════════════════════════════════════════════════════
#  ① 7개의 칸 — 스키마와 근거 대조
# ════════════════════════════════════════════════════════════════
CODE_SLOT_L = r"""
# ── 슬롯 정의 ── itda_core.py 87행
#  ★ 슬롯마다 '근거'를 함께 받는다. 프롬프트로
#    "지어내지 마세요" 하고 부탁하는 대신,
#    코드가 «대조»할 수 있게 만든 것이다.
def _slot(desc, enum=None):
    v = {'type': 'STRING'}
    if enum:
        v['enum'] = enum                # 목록에서만 고르게
    return {'type': 'OBJECT',
            'properties': {
                '값':   v,
                '근거': {'type': 'STRING',
                        'description':
                          '사용자가 실제로 한 말에서 그대로 인용'}},
            'required': ['값', '근거'],
            'description': desc}


# ── 근거 대조 ── itda_core.py 286행
def verify_slots(raw, user_msg):
    kept, dropped = {}, {}
    for k, v in raw.items():
        val, ev = v.get('값'), v.get('근거')
        if _grounded(ev, user_msg):
            kept[k] = val               # 발화에 있다 → 쓴다
        else:
            dropped[k] = val            # ★ 없다 → «버린다»
    return kept, dropped
""".strip('\n').split('\n')

CODE_SLOT_R = r"""
# ── 7개의 칸 ── itda_core.py 170행
PROFILE_SCHEMA = {'type': 'OBJECT', 'properties': {

  '관심분야':   _slot_multi('좋아한다·재밌다·해봤다고 말한 것'),

  '활동유형':   _slot_multi('무슨 행위를 좋아하나',
      ['만들기', '고치기·정비', '운전·조작', '돕기·돌봄',
       '가르치기', '분석·연구', '관리·운영',
       '표현·창작', '판매·설득']),

  '다루는대상': _slot_multi('일의 주 재료',
      ['사람', '기계·설비', '컴퓨터·데이터',
       '자연·생물', '창작물', '숫자·문서']),

  '세부관심':   _slot('관심을 좁히는 것 (어르신 / 웹 / 병원 …)'),

  '강점성향':   _slot('잘하거나 편한 것 (손재주 · 체력 …)'),

  '제약':      _slot_multi('사용자가 «말한» 현실적 부담',
      ['시간부족', '비용부담', '체력부담',
       '대인부담', '학력부담']),
}}
#  ↑ 6칸은 LLM 이. 7번째 '대상세부'는 코드가 채운다
#    ('할머니' 를 읽고 fill_obj_detail 이 '어르신')
""".strip('\n').split('\n')

# ════════════════════════════════════════════════════════════════
#  ② 관문 6개
# ════════════════════════════════════════════════════════════════
CODE_GATES = r"""
async def _step(self, db, profile, user_msg):        # itda_core.py 5331행

    # ①  이미 잠긴 사람인가 — 여기서 끝낸다. 트롤이 계속 보내도 «비용이 안 는다»
    if int(profile.get('_abuse') or 0) >= abuse_limits(profile)[1]:
        return {'kind': 'blocked', 'reply': ABUSE_STOP_REPLY, ...}

    # ②  「그거 말고」 · 「다 아니에요」 — 그 후보를 빼고 다시 찾는다
    _all_no   = none_of_these(user_msg)
    _rejected = rejects_last_card(user_msg) or _all_no

    # ③  「괜찮아요」 — 지원 정책을 «권했을 때만» 본다
    _off = profile.get('_policy_offered')
    if _off and declines_policy(user_msg):
        ...

    # ④  불법 신호 — ★ 차단하지 «않는다». 세기만 한다.
    #     이게 없으면 레드팀 15턴에서 그랬듯 카운터가 0으로 남아 «탐지조차 안 된다»
    _ill = illegal_signal(user_msg)
    if _ill:
        bump_abuse(profile, f'불법신호 {_ill[:2]}')

    # ⑤  pre_check — 네 갈래로 갈린다
    pc = pre_check(user_msg)
    if pc == 'SELFHARM':                        # ★ 차단하지 «않는다»
        _who = await self.crisis_who(user_msg, profile)   # 2층: 본인인가 제3자인가
        return {'kind': 'ask', ...}             # ☎109 를 건네고 대화는 «열어 둔다»

    # ⑥  공격 — 코드가 확정 차단. 문구도 코드가 쓴다 → 프롬프트가 샐 여지가 없다
    if is_injection(user_msg):
        _n = bump_abuse(profile, 'injection')
        return {'kind': 'redirect', ...}
""".strip('\n').split('\n')


def build(prs):
    #  ── 슬라이드 A ──────────────────────────────────────────
    s = blank_like(prs)
    header(s, 'PART 03 · 잇다 — 슬롯', '7개의 칸을 «이렇게» 정의했습니다')
    code_block(s, 1.00, 2.60, 8.85, 7.75, CODE_SLOT_L, pt=11.5)
    code_block(s, 10.15, 2.60, 8.85, 7.75, CODE_SLOT_R, pt=11.5)
    text(s, 1.00, 10.48, 18.0, 0.34,
         [('값과 «근거»를 한 묶음으로 받는 것이 핵심입니다 — '
           '값을 지어내려면 근거도 지어내야 하고, 그건 원문 대조에서 걸립니다',
           15, True, C_PRESENT)])

    #  ── 슬라이드 B ──────────────────────────────────────────
    s2 = blank_like(prs)
    header(s2, 'PART 03 · 잇다 — 예외 처리', '관문 6개는 «이렇게» 생겼습니다')
    #  ⚠ 12.5pt 로 넣었더니 «상자 밖으로 4줄이 삐져나와» 아래 문장과 겹쳤다.
    #    한글이 섞이면 줄 높이가 계산보다 커진다. 11.5pt 로 낮추고 끝의 두 줄을 뺐다.
    code_block(s2, 1.00, 2.60, 18.00, 7.75, CODE_GATES, pt=11.5)
    text(s2, 1.00, 10.48, 18.0, 0.34,
         [('여섯 관문 전부 «코드»입니다 — 턴당 0원 · 약 1밀리초. '
           '여기서 걸리면 LLM 을 한 번도 안 부릅니다', 15, True, C_PRESENT)])
    return s, s2


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    build(prs)
    prs.save(OUT)
    W, H = prs.slide_width / IN, prs.slide_height / IN
    out2 = Presentation(OUT)
    over = [(i, sh.shape_type) for i in (-2, -1)
            for sh in out2.slides[i].shapes
            if (sh.left + sh.width) / IN > W + 0.01 or (sh.top + sh.height) / IN > H + 0.01]
    print(f'  {n0}장 → {len(out2.slides)}장  (코드 슬라이드 2장 추가)')
    print(f'  슬라이드 밖: {len(over)}개  {"✅" if not over else over}')
    print(f'  저장 → {OUT}')


main()
