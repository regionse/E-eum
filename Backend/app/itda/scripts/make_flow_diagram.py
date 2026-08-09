# -*- coding: utf-8 -*-
"""「발화가 슬롯에 닿기까지」 — 위에서 아래로 흐르는 한 줄기. (2026-08-07)

사용자 지적 (앞 판이 왜 실패했나)
  「가독성이 없음. 그냥 위에서부터 아래로. 갈림길로 뭐가 뭔지 선으로 이어서 보여주는
   간단한 구조였으면 좋겠는데. 그 왼쪽에 글로 다 적으면 되고. 뭐가 뭐고 왜 필요한지」

  앞 판은 «가로 칩 한 줄» + «좌→우 갈림길» + «설명 상자» 를 한 장에 섞었다.
  눈이 세 번 방향을 바꾼다. 그래서 안 읽힌다.

이 판의 규칙
  ① 그림은 «위 → 아래» 한 줄기만. 가지는 오른쪽으로만 뻗는다.
  ② 글은 전부 «왼쪽 한 칸»에. 그림 안에는 이름과 결과만 둔다.
  ③ pre_check 만 네 갈래로 벌어진다 — 나머지는 한 줄 결과.

⚠ 값은 전부 코드 실측이다:
   · 관문 순서   itda_core.py `_step()` 5343~5560 실제 호출 순서 (칩 선택은 «뺐다» — 성질이 다르다)
   · 낱말 116개  SELF_HARM 27 + HARM_OTHERS 17 + ABUSE 4 + _ABUSE_VARIANT 35
                 + _SEXUAL_HARM 9 + _SEXUAL_HARD 13 + _SEXUAL_SOFT 11
   · 문턱        abuse_limits() 를 직접 돌려 얻음 (좋은턴 0→3/6 · 10→5/9 · 34→9/17 · 60→10/20)
   · 레드팀 15턴  illegal_signal 위 주석의 실측 기록

쓰는 법
  python make_flow_diagram.py [바탕.pptx] [출력.pptx]
"""
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR            # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v6.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v9.pptx'

IN = 914400
FONT = 'Arial'
C_HEAD    = RGBColor(0x8B, 0x84, 0xA0)
C_PRESENT = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE   = RGBColor(0x3B, 0x20, 0x63)
C_BODY    = RGBColor(0x55, 0x50, 0x6B)
C_CARD    = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_PANEL   = RGBColor(0xEF, 0xEB, 0xF6)
C_PANEL_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_MUTED   = RGBColor(0x8B, 0x84, 0xA0)
C_LINE    = RGBColor(0xC9, 0xC1, 0xDC)

C_PASS    = RGBColor(0x4E, 0x7D, 0x5E)
C_CONNECT = RGBColor(0xC2, 0x6B, 0x2E)
C_BLOCK   = RGBColor(0xA8, 0x36, 0x4A)


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.18):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.color.rgb = col
    return tb


def shape(sl, kind, x, y, w, h, fill, line=None, lw=0.75):
    sh = sl.shapes.add_shape(kind, Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.text = ''
    return sh


def seg(sl, x1, y1, x2, y2, col=C_LINE, w=2.0):
    cn = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                 Emu(int(x1 * IN)), Emu(int(y1 * IN)),
                                 Emu(int(x2 * IN)), Emu(int(y2 * IN)))
    cn.line.color.rgb = col
    cn.line.width = Pt(w)
    return cn


# ── 왼쪽 글 — «뭐가 뭐고 왜 필요한지» ────────────────────────────
NOTES = [
    #  ⚠ 처음엔 「경고 3+좋은턴÷5 · 종료 6+좋은턴÷3」 이라고 «공식»으로 적었다.
    #    사용자가 「좋은턴+3 이 무슨 소리냐」고 되물었다 — ÷ 가 «나누기»로 안 읽힌다.
    #    ⇒ 공식을 버리고 «말»로 쓴다. 발표 자료에 공식을 넣으면 읽는 데 시간이 든다.
    ('①  이미 잠긴 사람인가', [
        '딴 요청을 반복하면 잠급니다. 문턱은 «고정이 아닙니다» —',
        '시작은 경고 3 · 종료 6.  정상 대화 5턴마다 경고가, 3턴마다 종료가 1씩 늡니다.',
        '40턴쯤 대화했으면 경고 10 · 종료 19 까지 올라갑니다 (상한 10 · 20).',
        '왜 — 6턴에 딴소리 6번은 100%지만 40턴에 6번은 15%입니다. 같은 6이 다른 뜻입니다.',
        '⚠ 분모는 «정상 대화 턴»입니다. 전체 턴으로 세면 떠들수록 문턱이 올라 차단이 풀립니다.',
    ]),
    ('②  「그거 말고」인가', [
        '방금 준 카드를 거절한 것. 그 후보를 빼고 다시 찾습니다.',
    ]),
    ('③  「괜찮아요」인가', [
        '지원 정책을 권했을 때«만» 봅니다. 거절하면 그 안내를 멈춥니다.',
    ]),
    ('④  불법 신호인가   ← 차단 «안» 합니다', [
        '«세기만» 합니다. 레드팀 15턴 실측에서 로그가 «한 줄도» 안 찍혔습니다 —',
        'LLM 이 막긴 했는데 우리는 무슨 일이 있었는지 몰랐습니다. 모델이 바뀌면',
        '방어가 사라지는데 카운터가 0이라 «탐지조차» 안 됩니다.',
    ]),
    ('⑤  pre_check   ← 여기서 네 갈래', [
        '낱말 116개가 «차단기»가 아니라 «어디로 보낼지» 정합니다.',
        '★ 자해·폭력은 차단하지 않고 ☎109·129 로 «잇습니다». 남용에 절대 안 셉니다.',
    ]),
    ('⑥  공격인가', [
        '「이전 지시 무시하고…」 류. 되돌리고 남용 +1.',
    ]),
]

# ── 오른쪽 그림 — 관문과 결과 ────────────────────────────────────
GATES = [
    ('①', '이미 잠긴 사람인가',  '차단',           C_BLOCK),
    ('②', '「그거 말고」인가',    '후보에서 뺀다',    C_BODY),
    ('③', '「괜찮아요」인가',     '안내를 멈춘다',    C_BODY),
    ('④', '불법 신호인가',       '«세기만» 한다',   C_PRESENT),
]
FORKS = [
    ('통과',    '슬롯 채우기로',            C_PASS),
    ('되묻기',  '다르게 물어본다',           C_PRESENT),
    ('★ 잇기', '☎109·129 · 대화는 계속',   C_CONNECT),
    ('차단',    '되돌려 보낸다',            C_BLOCK),
]


def build(prs):
    s = blank_like(prs)

    text(s, 1.00, 0.92, 8.0, 0.30, [('PART 03 · 잇다 — 예외 처리', 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRESENT)])
    text(s, 1.00, 1.50, 19.0, 0.95,
         [('슬롯을 채우기 «전에» — 관문 6개', 44, True, C_TITLE)])

    # ══ 왼쪽 : 글 ═════════════════════════════════════════════
    LX, LW = 1.00, 7.10
    y = 2.72
    for head, lines in NOTES:
        text(s, LX, y, LW, 0.32, [(head, 16, True, C_TITLE)])
        y += 0.36
        for ln in lines:
            text(s, LX + 0.30, y, LW - 0.30, 0.28, [(ln, 13, False, C_BODY)])
            y += 0.27
        y += 0.20

    # ══ 오른쪽 : 위 → 아래 한 줄기 ════════════════════════════
    GX, GW = 8.60, 4.40          # 관문 상자
    OX, OW = 13.55, 5.45         # 결과 상자
    SPINE = GX + GW / 2          # 10.80
    BH = 0.50

    #  발화
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, GX + 0.80, 2.72, GW - 1.60, 0.52, C_CARD, C_CARD_LN)
    text(s, GX + 0.80, 2.85, GW - 1.60, 0.30,
         [('사용자 발화 한 마디', 15, True, C_TITLE)], align=PP_ALIGN.CENTER)
    prev_bottom = 3.24

    def gate_box(y, mark, label, hot=False):
        shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, GX, y, GW, BH,
              C_PANEL if hot else C_CARD, C_PRESENT if hot else C_CARD_LN,
              1.75 if hot else 0.75)
        text(s, GX + 0.28, y + 0.12, 0.42, 0.28,
             [(mark, 15, True, C_PRESENT if hot else C_MUTED)])
        text(s, GX + 0.76, y + 0.12, GW - 1.0, 0.28,
             [(label, 15, hot, C_TITLE)])

    def result(y, txt, col, h=BH):
        shape(s, MSO_SHAPE.RIGHT_ARROW, GX + GW + 0.08, y + h / 2 - 0.11, 0.36, 0.22, C_LINE)
        text(s, OX, y + h / 2 - 0.15, OW, 0.30, [(txt, 15, False, col)])

    #  ①~④
    ys = []
    y = 3.56
    for mark, label, res, col in GATES:
        seg(s, SPINE, prev_bottom, SPINE, y)
        gate_box(y, mark, label)
        result(y, res, col)
        prev_bottom = y + BH
        ys.append(y)
        y += 0.80

    #  ⑤ pre_check — 네 갈래
    y5 = y
    seg(s, SPINE, prev_bottom, SPINE, y5)
    #  ⚠ 여기 「낱말 116개」를 넣었다가 상자(h=0.50) 밖으로 삐져나와 테두리를 뚫었다.
    #    왼쪽 글에 이미 있는 말이라 그냥 뺀다 — 그림 안에는 «이름과 결과»만 둔다는 규칙대로.
    gate_box(y5, '⑤', 'pre_check', hot=True)

    FH, FGAP = 0.48, 0.08
    fys = [y5 + i * (FH + FGAP) for i in range(4)]
    SUB = OX - 0.52
    seg(s, GX + GW + 0.08, y5 + BH / 2, SUB, y5 + BH / 2, C_LINE, 2.2)
    seg(s, SUB, y5 + BH / 2, SUB, fys[-1] + FH / 2, C_LINE, 2.2)
    for (name, what, col), fy in zip(FORKS, fys):
        seg(s, SUB, fy + FH / 2, OX - 0.10, fy + FH / 2, col, 2.2)
        text(s, OX, fy + 0.02, 1.35, 0.28, [(name, 15, True, col)])
        text(s, OX + 1.45, fy + 0.04, OW - 1.45, 0.26, [(what, 13.5, False, C_BODY)])

    prev_bottom = y5 + BH
    #  ⑥
    y6 = fys[-1] + FH + 0.42
    seg(s, SPINE, prev_bottom, SPINE, y6)
    gate_box(y6, '⑥', '공격인가')
    result(y6, '되돌린다  ·  남용 +1', C_BLOCK)

    #  끝 — 슬롯으로
    yend = y6 + BH + 0.42
    seg(s, SPINE, y6 + BH, SPINE, yend)
    shape(s, MSO_SHAPE.ROUNDED_RECTANGLE, GX + 0.80, yend, GW - 1.60, 0.52,
          C_PANEL, C_PRESENT, 1.5)
    text(s, GX + 0.80, yend + 0.13, GW - 1.60, 0.30,
         [('슬롯 7칸 채우기로', 15, True, C_TITLE)], align=PP_ALIGN.CENTER)

    text(s, GX, yend + 0.66, OW + OX - GX, 0.30,
         [('여기까지 전부 코드 · 턴당 0원 · 약 1밀리초 — LLM 을 한 번도 안 부릅니다',
           13.5, False, C_MUTED)])
    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    build(prs)
    prs.save(OUT)
    W, H = prs.slide_width / IN, prs.slide_height / IN
    over = [(sh.shape_type, round((sh.left + sh.width) / IN, 2), round((sh.top + sh.height) / IN, 2))
            for sh in Presentation(OUT).slides[-1].shapes
            if (sh.left + sh.width) / IN > W + 0.01 or (sh.top + sh.height) / IN > H + 0.01]
    print(f'  {n0}장 → {len(Presentation(OUT).slides)}장')
    print(f'  슬라이드 밖: {len(over)}개  {"✅" if not over else over}')
    print(f'  저장 → {OUT}')


main()
