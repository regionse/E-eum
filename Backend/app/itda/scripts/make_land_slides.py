# -*- coding: utf-8 -*-
"""착지 판정 2장 — ③ 언제 답하나 · ④ 어떻게 골랐나(실측). (2026-08-07)

사용자 요청
  「1. can_land 라는 전제하에 이루어지는 거 — 어떻게 슬롯에서부터 시작되는가.
   2. 그리고 그걸 어떻게 «비교»하는가 — 이게 사실 우리가 가진 최대의 카드라고 생각해.
      (이건 실측 기록들, 없으면 지금 비교하는 게 나음)」

②가 최대 카드라는 판단에 동의한다. 그런데 **덱에 그 표가 없었다.**
실측기록_잇다.md §3-1 에 9개 신호를 14케이스로 잰 표가 그대로 있는데,
발표 자료에는 「n=52 p25=0.135」 한 줄로만 스쳐 지나갔다(그마저 지금은 지웠다).

⚠ 숫자는 전부 실측기록_잇다.md §3-1 ~ §3-1-c 원문에서 «그대로» 옮겼다. 지어낸 값 없다.
   · 9개 신호 정확도 (라벨 14케이스, 2026-08-06)
   · Vtop%(항목) vs Ctop%(묶음) — 용접 0.313 → 1.000
   · 안전 구간 폭 — C 단독 0.034 / V AND C 0.317
   · CLUSTER_H 스윕 — 0.05~0.30 전부 14/14, 0.32 에서 13/14
⚠ 무게 상수는 itda_core.py 1263~1267행 실측:
   LAND_W_USER 1.0 · LAND_W_CODE 0.5 · LAND_NEED 2.0 · LAND_NEED_LATE 1.5 · RELAX_AFTER 12

쓰는 법
  python make_land_slides.py [바탕.pptx] [출력.pptx]
"""
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR            # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v10.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v11.pptx'

IN = 914400
FONT = 'Arial'
C_HEAD    = RGBColor(0x8B, 0x84, 0xA0)
C_PRESENT = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE   = RGBColor(0x3B, 0x20, 0x63)
C_BODY    = RGBColor(0x55, 0x50, 0x6B)
C_CARD    = RGBColor(0xFF, 0xFF, 0xFF)
C_CARD_LN = RGBColor(0xE7, 0xE2, 0xF0)
C_PANEL   = RGBColor(0xEF, 0xEB, 0xF6)
C_PANEL_LN = RGBColor(0xDE, 0xD8, 0xEA)
C_MUTED   = RGBColor(0x8B, 0x84, 0xA0)
C_LINE    = RGBColor(0xC9, 0xC1, 0xDC)
C_WIN     = RGBColor(0x4E, 0x7D, 0x5E)
C_HOT     = RGBColor(0xA8, 0x36, 0x4A)
C_ORANGE  = RGBColor(0xC2, 0x6B, 0x2E)


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.18):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.color.rgb = col
    return tb


def box(sl, x, y, w, h, fill, line, lw=0.75):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.text = ''
    try:
        sh.adjustments[0] = 0.02
    except Exception:
        pass
    return sh


def bar(sl, x, y, w, h, col):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = col
    sh.line.fill.background()
    sh.shadow.inherit = False
    sh.text_frame.text = ''
    return sh


def header(s, sub, title):
    text(s, 1.00, 0.92, 9.0, 0.30, [(sub, 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRESENT)])
    text(s, 1.00, 1.50, 19.0, 0.95, [(title, 44, True, C_TITLE)])


# ════════════════════════════════════════════════════════════════
#  ③ 슬롯 → 착지 판정
# ════════════════════════════════════════════════════════════════
def slide_land(prs):
    s = blank_like(prs)
    header(s, 'PART 03 · 잇다 — 착지 판정', '칸이 «얼마나» 차야 답하나')

    # ── 왼쪽 : 글 ─────────────────────────────────────────
    LX, LW = 1.00, 8.30
    y = 2.68
    blocks = [
        ('세 축만 봅니다', [
            '관심분야 · 활동유형 · 다루는대상 — 직업을 «가르는» 축입니다.',
            '제약·강점성향은 안 셉니다. 「시간이 없다」로 직업이 좁혀지진 않으니까요.',
        ]),
        ('개수가 아니라 «무게»입니다', [
            '사용자가 자기 입으로 말한 축   1.0',
            '코드가 단서로 추론해 채운 축   0.5      ← 공짜로 얻은 건 반값',
            '합이 2.0 이상이어야 답합니다. 12턴을 넘으면 1.5 로 낮춥니다.',
        ]),
        ('왜 «반값»인가', [
            '「할머니」를 읽고 코드가 채운 축은 사용자가 «고른» 게 아닙니다.',
            '추론으로 얻은 축이 사용자 발화와 같은 무게면, 코드가 혼자',
            '조건을 채우고 답을 내버립니다.',
        ]),
        ('⚠ 한 문장에 다 차도 «한 번 미룹니다»', [
            '「빵 만드는 거 좋아해요」 한 마디에 [제빵] 카드가 나갔습니다.',
            '조건은 찼지만 «대화를 한 적이 없습니다».',
        ]),
    ]
    for head, lines in blocks:
        text(s, LX, y, LW, 0.34, [(head, 18, True, C_TITLE)])
        y += 0.40
        for ln in lines:
            text(s, LX + 0.28, y, LW - 0.28, 0.30, [(ln, 14, False, C_BODY)])
            y += 0.30
        y += 0.26

    # ── 오른쪽 : 저울 두 경우 ─────────────────────────────
    RX, RW = 10.00, 9.00
    for idx, (title, items, total, verdict, vcol) in enumerate([
        ('사용자가 두 축을 «말했다»',
         [('관심분야  “만드는 게 재밌었어요”', 1.0, C_PRESENT),
          ('활동유형  “돌보느라”', 1.0, C_PRESENT)],
         2.0, '2.0 ≥ 2.0   →   답한다', C_WIN),
        ('한 축은 코드가 «채웠다»',
         [('활동유형  “돌보느라”', 1.0, C_PRESENT),
          ('다루는대상  코드가 추론', 0.5, C_MUTED)],
         1.5, '1.5 < 2.0   →   더 묻는다', C_ORANGE),
    ]):
        by = 2.68 + idx * 3.62
        box(s, RX, by, RW, 3.30, C_CARD if idx == 0 else C_PANEL,
            C_CARD_LN if idx == 0 else C_PANEL_LN)
        text(s, RX + 0.40, by + 0.28, RW - 0.80, 0.32, [(title, 17, True, C_TITLE)])
        yy = by + 0.78
        for label, w, col in items:
            text(s, RX + 0.40, yy, 5.40, 0.30, [(label, 14, False, C_BODY)])
            bar(s, RX + 6.00, yy + 0.05, 1.55 * w, 0.20, col)
            text(s, RX + 7.75, yy - 0.02, 1.0, 0.30, [(f'{w:.1f}', 15, True, col)])
            yy += 0.46
        #  합계 줄
        s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                               Emu(int((RX + 0.40) * IN)), Emu(int(yy * IN)),
                               Emu(int((RX + RW - 0.40) * IN)), Emu(int(yy * IN))
                               ).line.color.rgb = C_LINE
        text(s, RX + 0.40, yy + 0.18, 5.40, 0.32, [('합', 15, True, C_TITLE)])
        text(s, RX + 6.00, yy + 0.14, 2.75, 0.36, [(f'{total:.1f}', 20, True, vcol)])
        text(s, RX + 0.40, yy + 0.72, RW - 0.80, 0.34, [(verdict, 17, True, vcol)])
    return s


# ════════════════════════════════════════════════════════════════
#  ④ 어떻게 골랐나 — 실측 비교  (실측기록 §3-1 ~ §3-1-c)
# ════════════════════════════════════════════════════════════════
SIGNALS = [
    ('C ent',  '중분류로 «묶어서» 잰 엔트로피', 14, True,  '안 쓰고 있던 것'),
    ('Ctop %', '최대 «묶음»의 확률 질량',       13, False, ''),
    ('R max',  '리랭커 최고점 절대값',          13, False, ''),
    ('V ent',  '벡터 점수 엔트로피',            12, True,  '쓰고 있던 것'),
    ('Vtop %', '1위 «항목»의 확률',            11, False, ''),
    ('Rtop %', '리랭커 1위 확률',              11, False, ''),
    ('R ent',  '리랭커 점수 엔트로피',          11, False, ''),
    ('ncl',    '묶음 개수',                    10, False, ''),
    ('margin', '1·2위 차',                    10, False, '이미 죽은 경로'),
]


def slide_compare(prs):
    s = blank_like(prs)
    header(s, 'PART 03 · 잇다 — 어떻게 골랐나', '신호 9개를 «재보고» 골랐습니다')
    text(s, 1.00, 2.42, 12.0, 0.30,
         [('되묻기가 맞는지 카드가 맞는지 — 라벨 14케이스에 신호 9개를 다 대봤습니다',
           15, False, C_MUTED)])

    # ── 왼쪽 : 정확도 표 ──────────────────────────────────
    LX, LW = 1.00, 10.20
    box(s, LX, 2.86, LW, 5.30, C_CARD, C_CARD_LN)
    y = 3.10
    for name, desc, score, mark, note in SIGNALS:
        col = C_WIN if score == 14 else (C_PRESENT if mark else C_BODY)
        text(s, LX + 0.34, y, 1.30, 0.30, [(name, 14.5, mark, col)])
        text(s, LX + 1.76, y, 4.30, 0.30, [(desc, 13, False, C_MUTED)])
        bar(s, LX + 6.20, y + 0.07, 0.176 * score, 0.16,
            col if mark or score == 14 else C_LINE)
        text(s, LX + 8.70, y - 0.02, 0.80, 0.30, [(f'{score}/14', 14, mark, col)])
        if note:
            text(s, LX + 6.20, y + 0.28, 3.60, 0.24, [('← ' + note, 11.5, False, col)])
            y += 0.22
        y += 0.34

    text(s, LX + 0.34, 8.28, LW, 0.32,
         [('정답표를 신호보다 «먼저» 붙였습니다. 측정한 뒤에 고치지 않았습니다.',
           15, True, C_PRESENT)])

    # ── 오른쪽 : 세 가지 발견 ────────────────────────────
    RX, RW = 11.90, 7.10
    finds = [
        ('① 쓰고 있던 게 4등이었습니다',
         ['V ent (쓰던 것)     12 / 14',
          'C ent (안 쓰던 것)  14 / 14',
          '「항목」이 아니라 «묶음»으로 재야 했습니다.',
          '용접 — 항목 0.313  →  묶음 1.000'], C_PRESENT),
        ('② 정확도보다 «마진»입니다',
         ['C 단독        안전 구간 0.034',
          'V AND C       안전 구간 0.317   ← 10배',
          '14케이스에서 「몇 개 맞았나」는 쉽게 뒤집힙니다.',
          '문턱을 «어디에 둬도 되는 폭»이 진짜 견고함입니다.'], C_WIN),
        ('③ 그래서 문턱도 훑었습니다',
         ['0.05 ~ 0.30   전부 14 / 14',
          '0.32          13 / 14   ← 무너짐',
          '0.30 은 무너지기 «바로 직전»이라 0.15 를 씁니다.'], C_ORANGE),
    ]
    y = 2.86
    for head, lines, col in finds:
        h = 0.62 + len(lines) * 0.30
        box(s, RX, y, RW, h, C_PANEL, C_PANEL_LN)
        text(s, RX + 0.34, y + 0.24, RW - 0.68, 0.32, [(head, 16.5, True, col)])
        yy = y + 0.66
        for ln in lines:
            text(s, RX + 0.34, yy, RW - 0.68, 0.28, [(ln, 13.5, False, C_BODY)])
            yy += 0.30
        y += h + 0.26

    text(s, RX, y + 0.06, RW, 0.60,
         [('만점을 받고도 «되돌린» 것이 있습니다 —', 14, True, C_HOT),
          ('정답표가 틀렸기 때문입니다. 다음 장.', 14, True, C_HOT)])
    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    slide_land(prs)
    slide_compare(prs)
    prs.save(OUT)
    W, H = prs.slide_width / IN, prs.slide_height / IN
    out2 = Presentation(OUT)
    over = [(i, round((sh.left + sh.width) / IN, 2), round((sh.top + sh.height) / IN, 2))
            for i in (-2, -1) for sh in out2.slides[i].shapes
            if (sh.left + sh.width) / IN > W + 0.01 or (sh.top + sh.height) / IN > H + 0.01]
    print(f'  {n0}장 → {len(out2.slides)}장  (착지 판정 2장 추가)')
    print(f'  슬라이드 밖: {len(over)}개  {"✅" if not over else over}')
    print(f'  저장 → {OUT}')


main()
