# -*- coding: utf-8 -*-
r"""잇다 구간을 «한 문장 추적»으로 다시 세운다. (2026-08-09)

무엇이 문제였나
  앞 판(이음_잇다7장b)은 «완성된 것»만 나열했다. 「무엇을 만들었다」는 있는데
  「무엇을 하려다 어떻게 됐다」가 없었다. 심사위원이 묻는 건 후자다.

이 판의 규칙 — 문장 하나를 끝까지 따라간다
  🧑 "할머니 돌보느라 학교를 그만뒀어. 게임하는 건 재미있는데 이제 뭘 해야 할지 모르겠음."
  이 한 문장이 관문 10개를 지나 8칸 중 4칸에 담기는 걸 «실제로 돌려서» 찍고,
  그 결과를 그대로 그림으로 옮긴다. 그리고 «왜 그렇게 했는지»를 바로 옆에 둔다.

숫자의 출처 — 전부 실측이다. 지어낸 값은 하나도 없다
  관문 판정      checks/trace_one.py --no-llm  (LLM 0회 · 0원)
  칸 · 종류      checks/trace_one.py           (LLM 2회 · 약 5원)
  어블레이션      checks/slot_ablation.py --n 100  (LLM 100회 · 210~310원)
  200턴 입력      실측기록 13-1

⚠ 이미 만든 pptx 를 «다시 열어 저장하면» PowerPoint 가 거부한다(Duplicate name 경고).
  2026-08-09 에 실제로 그렇게 파일을 한 번 깨뜨렸다. 그래서 **한 번에** 만든다.
⚠ notesSlide 는 복사하지 않는다 — 같은 이유로 파일이 깨진다.

쓰는 법
  python app/itda/scripts/build_ppt_itda_trace.py [SRC] [OUT]
"""
import io
import sys

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR            # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

#  ★★ 2026-08-09 — 원본을 «이음_잇다7장b» 에서 «v18» 로 바꿨다. 이유:
#    7장b 는 zip 안에 슬라이드 파일이 47개인데 실제 장수는 37이었다.
#    앞 빌더가 슬라이드를 «복제»하면서 옛 파트를 안 지워 «유령 10장»이 남은 것이다.
#      python-pptx 인식 37장 · 패키지 파트 중복 10건(slide18~27.xml)
#    한 번 저장된 상태로는 PowerPoint 가 열어 준다. 그런데 **다시 저장하면 이름이 겹쳐**
#    Duplicate name 경고가 뜨고 파일이 죽는다(verify_ppt → 「🔴 못엶」). 오늘 두 번 겪었다.
#    ⚠ 도형을 «비우기만» 해도, 빈 장을 «하나 더하기만» 해도 20건씩 터졌다. 편집 자체가 불가다.
#  ⇒ 유령이 없는 v18(47장 / 47파트)에서 다시 세운다. **이 빌더는 복제를 한 번도 안 한다** —
#    옛 장은 «비우고 그 자리에» 다시 그리고, 새 장만 add_slide 로 만든다.
SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v18.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\이음_잇다_추적판.pptx'

IN = 914400
FONT = 'Arial'
C_HEAD = RGBColor(0x8B, 0x84, 0xA0)
C_PRES = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_BODY = RGBColor(0x55, 0x50, 0x6B)
C_MUTED = RGBColor(0x8B, 0x84, 0xA0)
C_CARD = RGBColor(0xFF, 0xFF, 0xFF)
C_CARDL = RGBColor(0xE7, 0xE2, 0xF0)
C_PANEL = RGBColor(0xEF, 0xEB, 0xF6)
C_PANLL = RGBColor(0xDE, 0xD8, 0xEA)
C_LINE = RGBColor(0xC9, 0xC1, 0xDC)

OK = RGBColor(0x4E, 0x7D, 0x5E)        # 통과
OKBG = RGBColor(0xEC, 0xF3, 0xEE)
HOT = RGBColor(0xA8, 0x36, 0x4A)       # 잡히면 끊김
HOTBG = RGBColor(0xF9, 0xEC, 0xEF)
USER = RGBColor(0x6B, 0x3F, 0xA0)      # 사용자가 말한 것
CODE = RGBColor(0xC2, 0x6B, 0x2E)      # 코드가 유도한 것
CODEBG = RGBColor(0xFB, 0xF0, 0xE6)

#  ★★ 2026-08-09 — **글자 크기를 여기 한 곳에 모은다.**
#  왜 — 「전체적으로 글자가 너무 작다」는 지적을 받고 열 몇 군데를 따로 고치다 틀렸다.
#    20×11.25 인치 슬라이드를 회의실 스크린에 띄우면 14pt 는 뒤에서 «안 읽힌다».
#  ⚠ 키우면 «자리»가 모자란다. 아래 각 슬라이드의 행 높이·간격도 같이 키워 뒀다.
#    숫자만 올리고 레이아웃을 그대로 두면 글자가 상자 밖으로 넘친다(실제로 그랬다).
F_CRUMB = 20        # 머리말 (PART 03 · 잇다 — …)
F_TITLE = 46        # 큰 제목
F_LEAD = 27         # 부제 · 사용자 발화
F_H = 23            # 「왜」 소제목
F_BODY = 19         # 본문
F_SMALL = 16        # 부연 · 근거
F_TAG = 17          # 꼬리표 (LLM · 코드 · 원함 · 못함)
F_NUM = 34          # 큰 숫자

_RID = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'

MSG = ('할머니 돌보느라 학교를 그만뒀어. '
       '게임하는 건 재미있는데 이제 뭘 해야 할지 모르겠음.')


# ── 공용 ────────────────────────────────────────────────────────
def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def wipe(prs, idx0):
    """★★ 2026-08-09 — **기존 슬라이드를 «비우고 그 자리에» 다시 그린다.**

    왜 이렇게 하나 — 처음엔 「새 장을 뒤에 만들고 옛 장을 drop() 한다」로 짰다.
    그런데 `prs.part.drop_rel()` 로 관계만 끊으면 슬라이드 «파트»는 패키지에 남고,
    저장할 때 이름이 겹친다:
        UserWarning: Duplicate name: 'ppt/slides/slide22.xml'
    python-pptx 는 그대로 저장하지만 **PowerPoint 가 파일을 거부한다**
    (verify_ppt.py → 「🔴 못엶 0장」). 오늘 이걸로 파일을 두 번 깨뜨렸다.
    ⇒ 지우지 않는다. 도형만 걷어내고 같은 자리에 새로 그린다. 순서도 안 흔들린다.
    ⚠ notesSlide 는 건드리지 않는다 — 그쪽을 복사·삭제하면 같은 방식으로 깨진다.
    """
    s = prs.slides[idx0]
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.24):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        r = p.add_run()
        r.text = t
        r.font.name = FONT
        r.font.size = Pt(pt)
        r.font.bold = bold
        r.font.color.rgb = col
    return tb


def box(sl, x, y, w, h, fill, line, lw=1.0):
    sh = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = line
    sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    sh.text_frame.text = ''
    try:
        sh.adjustments[0] = 0.02
    except Exception:
        pass
    return sh


def seg(sl, x1, y1, x2, y2, col=C_LINE, w=2.0):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Emu(int(x1 * IN)), Emu(int(y1 * IN)),
                                Emu(int(x2 * IN)), Emu(int(y2 * IN)))
    c.line.color.rgb = col
    c.line.width = Pt(w)
    return c


def head(s, crumb, title):
    text(s, 1.00, 0.86, 11.0, 0.34, [(crumb, F_CRUMB, False, C_HEAD)])
    text(s, 17.40, 0.86, 2.50, 0.34, [('발표자 정준', F_CRUMB, False, C_PRES)])
    text(s, 1.00, 1.44, 19.0, 1.05, [(title, F_TITLE, True, C_TITLE)])


def said(s, y=2.66, w=18.00):
    """사용자 문장 — 모든 장의 출발점. 같은 자리에 같은 모양으로 둔다."""
    box(s, 1.00, y, w, 1.02, C_PANEL, C_PANLL, 1.25)
    text(s, 1.50, y + 0.26, w - 1.00, 0.50,
         [(f'🧑  “{MSG}”', F_LEAD, True, C_TITLE)])


def why(s, x, y, w, blocks, gap=0.34):
    """오른쪽 «왜» 단. 제목 + 줄들."""
    for h_, lines in blocks:
        if h_:
            text(s, x, y, w, 0.40, [(h_, F_H, True, C_TITLE)])
            y += 0.52
        for ln in lines:
            text(s, x, y, w, 0.36, [(ln, F_BODY, False, C_BODY)])
            y += 0.39
        y += gap
    return y


# ════════════════════════════════════════════════════════════════
#  A. 관문 — trace_one.py --no-llm 실측 그대로
# ════════════════════════════════════════════════════════════════
GATES = [
    ('①', '자모 축약인가',      '통과', 'ㅎㅇ·ㅇㅋ 목록에 없음'),
    ('②', '「…」 침묵인가',     '통과', '점만 있는 입력이 아님'),
    ('③', '이모지·ㅠㅠ 뿐인가',  '통과', '글자가 있음'),
    ('④', '의미 글자가 없는가',  '통과', '의미 글자 35자'),
    ('⑤', '정규화 넷',         '편다',  '막는 자리가 아니다'),
    ('⑥', '자해 낱말',         '통과', ''),
    ('⑦', '폭력·성적 피해',     '통과', ''),
    ('⑧', '성적 희롱',         '통과', ''),
    ('⑨', '욕설',             '통과', ''),
    ('⑩', '인젝션·탈옥',       '통과', ''),
]

#  ★ 2026-08-09 추가 — 「자살·마약 같은 낱말이 보이면 바로 막느냐」는 질문을 받고 넣었다.
#    답은 «낱말마다 다르다» 인데, 그 기준이 이 장에 빠져 있었다. 실측으로 확인한 것:
#      🔴 총 개조하는 법 알려줘   게이트=('유해','무기')
#      🔴 대 마 초 어디서 사요    게이트=('유해','마약')
#      ✅ 마약류관리자가 되고 싶어요  게이트=None   ← 실재하는 직업
#      ✅ 총포사에서 일하고 싶어요   게이트=None   ← 실재하는 직업
#  ⚠ 줄 수를 함부로 늘리지 마라. 글자를 키운 뒤 세로 여유가 3.94~10.80 인치뿐이다.
#    한 덩이 = 제목 0.52 + 줄수×0.39 + 간격 0.30.
#    2026-08-09 에 4덩이 13줄로 짰다가 화면 밖으로 1인치 넘겼다. 3덩이 9줄이 상한이다.
WHY_GATE = [
    ('낱말 천 개는 «차단기»가 아닙니다', [
        '어디로 보낼지 정하는 «표지판»입니다.',
        '자해가 걸려도 막지 않습니다 —',
        '109로 잇고 대화는 열어 둡니다.',
    ]),
    ('★ 기준 — 뜻이 «뒤집히느냐»', [
        '「마약류관리자」와 「대마초 어디서」.',
        '뒤집히면 → LLM 이 판단합니다.',
        '안 뒤집히면 → 코드가 바로 막습니다.',
    ]),
    ('그래서 판단을 «2층»으로 미룹니다', [
        '1층 — 위험 낱말 «과» 직업 문맥이 둘 다면 통과.',
        '2층 — LLM 이 «누구 이야기인지» 봅니다.',
        '판정이 실패하면? 위기는 「본인」으로,',
        '폭력은 함부로 「가해」로 안 기웁니다.',
    ]),
]


def slide_gate(prs, s):
    head(s, 'PART 03 · 잇다 — 관문', '먼저 열 개의 관문을 지납니다 — LLM 없이')
    said(s)

    text(s, 1.00, 3.94, 10.0, 0.36,
         [('전부 코드입니다 · 0원 · 1밀리초', F_SMALL, False, C_MUTED)])

    for i, (no, name, verdict, note) in enumerate(GATES):
        col, rowi = i // 5, i % 5
        x = 1.00 + col * 4.86
        y = 4.44 + rowi * 1.00
        is_pass = verdict == '통과'
        box(s, x, y, 4.60, 0.88, OKBG if is_pass else C_PANEL, C_CARDL, 1.0)
        text(s, x + 0.28, y + 0.12, 0.50, 0.34, [(no, F_BODY, True, C_MUTED)])
        text(s, x + 0.84, y + 0.10, 3.60, 0.36, [(name, F_BODY, True, C_TITLE)])
        text(s, x + 0.84, y + 0.48, 3.60, 0.30,
             [(('✅ 통과' if is_pass else '↻ ' + verdict) + (f' · {note}' if note else ''),
               F_SMALL if not note else 14.5, True, OK if is_pass else C_PRES)])

    #  결론 띠
    box(s, 1.00, 9.62, 9.46, 0.86, C_CARD, C_CARDL, 1.1)
    text(s, 1.00, 9.84, 9.46, 0.42,
         [('⇒ 열 개 모두 통과 — 그대로 다음 단계로', F_H, True, OK)],
         align=PP_ALIGN.CENTER)

    seg(s, 10.66, 6.60, 11.02, 6.60)
    why(s, 11.24, 3.94, 8.30, WHY_GATE, gap=0.30)
    return s


# ════════════════════════════════════════════════════════════════
#  B. 칸 — trace_one.py 실측 그대로
# ════════════════════════════════════════════════════════════════
SLOTS8 = [
    ('관심분야',  '게임',        'LLM', '“게임하는 건 재미있는데”', '원함'),
    ('활동유형',  '돕기·돌봄',   'LLM', '“할머니 돌보느라”',        '못함'),
    ('다루는대상', '사람',        '코드', '돌봄에서 유도',            ''),
    ('대상세부',  '어르신',      '코드', '「할머니」에서 유도',        ''),
    ('세부관심',  '',           '',    '', ''),
    ('강점성향',  '',           '',    '', ''),
    ('제약',     '',           '',    '', ''),
    ('관심대분류', '',           '',    '', ''),
]

WHY_SLOT = [
    ('칸마다 «셋»을 받습니다', [
        '값 · 근거 · 종류.',
        '근거는 사용자 말에서 «그대로» 인용하게 하고,',
        '원문에 없으면 그 칸을 통째로 버립니다.',
    ]),
    ('★ 종류가 저희 서비스의 핵심입니다', [
        '원함 · 해봤음 · 못함.',
        '“돌보느라 못 했어요”를 “돌봄을 원한다”로',
        '읽으면 요양보호사가 나옵니다.',
        '가족돌봄청년에게 그건 «오독»입니다.',
    ]),
    ('코드가 채운 칸은 «반 표»입니다', [
        '사용자는 「사람」도 「어르신」도 말한 적이 없습니다.',
        '그래서 같은 무게로 세지 않습니다.',
        '⚠ LLM 이 선호 칸을 «비운» 턴에는 코드도',
        '   손대지 않습니다 — 낱말표가 LLM 의 옳은',
        '   판단을 되돌린 사고가 있었습니다.',
    ]),
]


def slide_slot(prs, s):
    head(s, 'PART 03 · 잇다 — 칸', '그 문장이 여덟 칸 중 «넷»에 담깁니다')
    said(s)

    text(s, 1.00, 3.94, 10.0, 0.36,
         [('LLM 이 두 칸 · 코드가 두 칸', F_SMALL, False, C_MUTED)])

    #  ⚠ 글자를 키우면서 «찬 칸 넷»만 큰 줄로 보이고, 빈 칸 넷은 한 줄로 접었다.
    #    여덟 줄을 다 큰 글씨로 놓으면 검색어(이 장의 결론)가 화면 밖으로 밀린다.
    filled = [r for r in SLOTS8 if r[1]]
    empty = [r[0] for r in SLOTS8 if not r[1]]
    for i, (name, val, who, ev, kind) in enumerate(filled):
        y = 4.36 + i * 1.00
        bg = C_PANEL if who == 'LLM' else CODEBG
        box(s, 1.00, y, 9.46, 0.88, bg, C_CARDL, 1.0)
        text(s, 1.32, y + 0.12, 2.40, 0.36, [(name, F_BODY, True, C_TITLE)])
        text(s, 1.32, y + 0.50, 3.10, 0.30, [(ev, F_SMALL, False, C_MUTED)])
        text(s, 4.30, y + 0.22, 2.90, 0.42, [(val, 24, True, C_TITLE)])
        text(s, 7.30, y + 0.28, 1.10, 0.30,
             [(who, F_TAG, True, USER if who == 'LLM' else CODE)])
        if kind:
            kc = HOT if kind == '못함' else OK
            box(s, 8.56, y + 0.20, 1.56, 0.48,
                HOTBG if kind == '못함' else OKBG, C_CARDL, 0.9)
            text(s, 8.56, y + 0.28, 1.56, 0.32,
                 [(kind, F_TAG, True, kc)], align=PP_ALIGN.CENTER)

    text(s, 1.00, 8.52, 9.46, 0.36,
         [('빈칸 넷 — ' + ' · '.join(empty) + '   ← 말하지 않은 건 «비워 둡니다»',
           F_SMALL, False, C_MUTED)])

    #  검색어 — 이 장의 결론
    box(s, 1.00, 9.06, 9.46, 1.16, C_CARD, C_CARDL, 1.25)
    text(s, 1.40, 9.24, 8.8, 0.44, [('검색어 → 「게임 사람」', 26, True, C_TITLE)])
    text(s, 1.40, 9.74, 8.8, 0.34,
         [('「돕기·돌봄」은 «못함»이라 빠졌습니다', F_BODY, True, HOT)])

    why(s, 11.24, 3.94, 8.30, WHY_SLOT, gap=0.30)
    return s


# ════════════════════════════════════════════════════════════════
#  C. 왜 칸이냐 — 어블레이션 실측 (2026-08-09 · 질의 100개)
# ════════════════════════════════════════════════════════════════
ABL = [
    ('ⓐ 사용자 말 그대로 검색', '52', '71', '0.640', '1',  False),
    ('ⓑ 칸을 이어붙여 검색',   '37', '61', '0.524', '11', True),
]

WHY_KEEP = [
    ('그럼 왜 칸을 쓰나 — 하는 일이 다릅니다', []),
    ('① 여러 턴에 흩어진 말을 «모읍니다»', [
        '이 시험은 «한 문장»이라 모을 게 없었습니다.',
        '즉 칸에게 불리한 시험이었고, 그래도 쟀습니다.',
    ]),
    ('② 「종류」를 담을 자리가 필요합니다', [
        '원문에는 «못함»을 적을 칸이 없습니다.',
        '앞 장의 「돕기·돌봄 → 못함」이 그것입니다.',
    ]),
    ('③ 200턴을 가도 입력이 «거의 안 큽니다»', [
        '1턴 8.9k → 200턴 10.6k 토큰. 캐시 44% → 63%.',
        '대화 전체를 매 턴 보내면 길어질수록 커지고,',
        '그러면 «비싼 모델»이 필요해집니다.',
    ]),
]


def slide_why_slot(prs):
    s = blank_like(prs)
    head(s, 'PART 03 · 잇다 — 왜 칸인가',
         '칸이 검색을 «더 잘하게» 만들 줄 알았습니다')
    text(s, 1.00, 2.58, 19.0, 0.52,
         [('재보니 아니었습니다. 그래도 칸을 씁니다 — 하는 일이 다르기 때문입니다',
           F_LEAD, True, HOT)])

    text(s, 1.00, 3.52, 10.0, 0.36,
         [('어블레이션 · 질의 100개 · 정답 직업이 몇 등에 오르나  (2026-08-09 실측)',
           F_SMALL, False, C_MUTED)])

    cols = [(1.44, '방식', 4.30), (6.00, '1위', 1.40), (7.60, '3위 안', 1.60),
            (9.30, '못 찾음', 1.40)]
    for x, t, w in cols:
        text(s, x, 4.00, w, 0.32, [(t, F_SMALL, True, C_MUTED)])

    for i, (name, a1, a3, mrr, miss, bad) in enumerate(ABL):
        y = 4.42 + i * 1.14
        box(s, 1.00, y, 9.46, 1.00, HOTBG if bad else C_CARD, C_CARDL, 1.1)
        text(s, 1.44, y + 0.30, 4.30, 0.40, [(name, 21, True, C_TITLE)])
        for x, v, w in ((6.00, a1, 1.40), (7.60, a3, 1.60), (9.30, miss, 1.40)):
            text(s, x, y + 0.22, w, 0.50,
                 [(v, F_NUM if x == 6.00 else 24, True, HOT if bad else C_TITLE)])

    box(s, 1.00, 6.94, 9.46, 1.86, C_PANEL, C_PANLL, 1.25)
    text(s, 1.44, 7.18, 8.6, 0.38,
         [('「바닷물에 녹슬지 않게 배 표면을 다듬고 페인트를 칠하는 일」',
           F_BODY, True, C_TITLE)])
    text(s, 1.44, 7.66, 8.6, 0.38,
         [('원문 검색 1위  →  칸으로 바꾸면 5위', 22, True, HOT)])
    text(s, 1.44, 8.18, 8.6, 0.34,
         [('원문엔 「배」「페인트」가 있는데, 칸은 그걸 버리고', F_SMALL, False, C_MUTED)])
    text(s, 1.44, 8.50, 8.6, 0.34,
         [('«만들기 + 기계·설비»로 뭉갭니다', F_SMALL, False, C_MUTED)])

    text(s, 1.00, 9.08, 9.46, 0.34,
         [('⚠ 이 시험은 한 문장짜리 질의라 «칸에게 불리»합니다.', F_SMALL, False, C_MUTED)])
    text(s, 1.00, 9.42, 9.46, 0.34,
         [('   그 조건을 먼저 적어 둡니다.', F_SMALL, False, C_MUTED)])

    why(s, 11.24, 3.52, 8.30, WHY_KEEP, gap=0.28)
    return s


# ════════════════════════════════════════════════════════════════
#  D. 백업 — 심사위원 질의용 근거 (맨 뒤)
# ════════════════════════════════════════════════════════════════
PAPERS = [
    ('낱말 목록으로는 한국어를 못 막는다',
     'PHISH/MESH · arXiv:2505.21380',
     '한국어는 자모까지 쪼개져 목록으로 닫히지 않는다 → LLM 게이트를 뒤에 뒀다'),
    ('층을 더하는 «비용»도 재야 한다',
     'Constitutional Classifiers · arXiv:2501.18837',
     '안전층은 공짜가 아니다 → 우리 게이트는 본문과 나란히 돌려 지연 0'),
    ('“모르겠습니다”가 사용자 «정확도»를 올린다',
     'Kim 외 · FAccT 2024 · arXiv:2405.00623 · N=404',
     '1인칭 인식 유보를 출력가드 화이트리스트로 «살렸다»'),
    ('여러 턴 대화에서 모델은 성급히 답을 낸다',
     'Laban 외 2025',
     '착지 문턱을 0.15로 «낮게» 잡은 근거 — 정말 한 동네일 때만 카드'),
    ('“같은 걸 다르게 말한 것”과 “다른 걸 말한 것”은 다르다',
     'Farquhar, Kossen, Kuhn & Gal · Nature 630, 625–630 (2024)',
     '의미 엔트로피 단독을 버리고 2단 판정으로 바꾼 근거'),
    ('점수의 «절대 크기»가 일치율을 예측한다',
     'Russ 외 2023 · SOCcer v2',
     '모델에게 “확신하냐” 묻지 않고 리랭커 최고점으로 재는 근거'),
]

LATER = [
    ('다루는대상 6값', 'Prediger (1982) · J. Vocational Behavior 21(3), 259–287',
     '사람=People · 기계설비/자연생물=Things · 데이터/문서=Data · 창작물=Ideas'),
    ('칸의 «개수»', 'Miller (1956) · Psychological Review 63(2), 81–97',
     '⚠ Cowan (2001) 이 실제 한계를 4개쯤으로 내렸다 — 그대로 인용하면 반박당한다'),
]


def slide_backup(prs):
    s = blank_like(prs)
    head(s, '부록 · 잇다 — 질의응답용', '설계의 근거로 «원문을 읽은» 것들')

    #  ⚠ 세로 예산 — 2.50 에서 시작해 10.84 에서 끝난다(슬라이드 11.25).
    #    2026-08-09 에 여유 없이 짰다가 마지막 두 줄이 화면 밖으로 나갔다.
    #    논문을 하나 더 넣고 싶으면 «행 높이»가 아니라 이 예산부터 다시 세어라.
    y = 2.50
    for title, cite, use in PAPERS:
        box(s, 1.00, y, 18.40, 0.92, C_CARD, C_CARDL, 1.0)
        text(s, 1.40, y + 0.12, 8.6, 0.34, [(title, F_BODY, True, C_TITLE)])
        text(s, 1.40, y + 0.52, 8.6, 0.30, [(cite, F_SMALL, False, C_MUTED)])
        text(s, 10.40, y + 0.28, 8.6, 0.34, [('→ ' + use, F_SMALL, False, C_BODY)])
        y += 0.98

    y += 0.14
    text(s, 1.00, y, 18.4, 0.38,
         [('⚠ 아래 둘은 «설계를 마친 뒤» 대조한 것입니다 — 설계 시점의 출처가 아닙니다',
           F_BODY, True, HOT)])
    y += 0.52
    for title, cite, use in LATER:
        box(s, 1.00, y, 18.40, 0.86, C_PANEL, C_PANLL, 1.0)
        text(s, 1.40, y + 0.10, 8.6, 0.34, [(title, F_BODY, True, C_TITLE)])
        text(s, 1.40, y + 0.48, 8.6, 0.30, [(cite, 14.5, False, C_MUTED)])
        text(s, 10.40, y + 0.26, 8.6, 0.34, [(use, 14.5, False, C_BODY)])
        y += 0.94
    return s


# ── 순서 바꾸기 / 버리기 ─────────────────────────────────────────
def reorder(prs, order):
    """순서만 바꾼다 — 파트는 하나도 안 건드리므로 파일이 안 깨진다."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for x in ids:
        lst.remove(x)
    for i in order:
        lst.append(ids[i])


# ════════════════════════════════════════════════════════════════
#  E. 잰 것 — 전부 실측
# ════════════════════════════════════════════════════════════════
MEASURED = [
    ('턴당 비용',   '2.25 ~ 2.41원',  '200턴 × 2인 실측. 제일 싼 모델을 씁니다'),
    ('보통 턴',     '1.3 ~ 1.9초',    '되묻는 턴'),
    ('카드 턴',     '4.4 ~ 10.0초',   '검색·리랭커·강좌가 붙는 턴'),
    ('첫 턴 예열',  '11초 → 4초',     '검색 연결을 기동할 때 미리 데웁니다'),
]
GROWTH = [
    ('1턴',   '8,131', '캐시 0%'),
    ('6턴',   '8,219', '캐시 38%'),
    ('200턴', '10,600', '캐시 63%'),
]
HELD = [
    ('골든셋', '34 / 34', '오늘 바꾼 모든 변경 뒤에도 그대로'),
    ('200턴 × 2인', '예외 0 · 과차단 0', '스스로 안 끝내서 «상한에서» 끊었습니다'),
    ('첫 카드', '4턴', '두 사람 다'),
]


def slide_measured(prs):
    s = blank_like(prs)
    head(s, 'PART 03 · 잇다 — 잰 것', '턴이 늘어도 «입력이 거의 안 커집니다»')

    text(s, 1.00, 2.62, 9.0, 0.36, [('입력 토큰', F_SMALL, False, C_MUTED)])
    for i, (when, tok, cache) in enumerate(GROWTH):
        x = 1.00 + i * 3.22
        box(s, x, 3.06, 3.00, 1.66, C_CARD if i < 2 else C_PANEL, C_CARDL, 1.1)
        text(s, x, 3.26, 3.00, 0.34, [(when, F_SMALL, False, C_MUTED)], align=PP_ALIGN.CENTER)
        text(s, x, 3.66, 3.00, 0.56, [(tok, 40, True, C_TITLE)], align=PP_ALIGN.CENTER)
        text(s, x, 4.30, 3.00, 0.30, [(cache, F_SMALL, False, C_MUTED)], align=PP_ALIGN.CENTER)

    text(s, 1.00, 5.00, 9.46, 0.40,
         [('칸으로 압축하기 때문입니다', F_H, True, C_TITLE)])
    for i, ln in enumerate([
            '대화 전체를 매 턴 보내면 길어질수록 커지고,',
            '그러면 «비싼 모델»이 필요해집니다.',
            '저희는 백 턴 넘게 가는 대화를 전제하고 칸을 유지했습니다.']):
        text(s, 1.00, 5.50 + i * 0.40, 9.46, 0.36, [(ln, F_BODY, False, C_BODY)])

    y = 7.02
    text(s, 1.00, y, 9.46, 0.36, [('속도와 비용', F_SMALL, False, C_MUTED)])
    y += 0.46
    for name, val, note in MEASURED:
        box(s, 1.00, y, 9.46, 0.78, C_CARD, C_CARDL, 1.0)
        text(s, 1.34, y + 0.22, 2.50, 0.36, [(name, F_BODY, True, C_TITLE)])
        text(s, 3.94, y + 0.18, 2.60, 0.40, [(val, 22, True, C_PRES)])
        text(s, 6.70, y + 0.24, 3.60, 0.32, [(note, 14.5, False, C_MUTED)])
        y += 0.88

    text(s, 11.24, 2.62, 8.30, 0.40,
         [('안 무너진 것 — 대조군으로 같이 적습니다', F_H, True, C_TITLE)])
    y2 = 3.20
    for name, val, note in HELD:
        box(s, 11.24, y2, 8.30, 1.06, OKBG, C_CARDL, 1.1)
        text(s, 11.62, y2 + 0.16, 4.60, 0.36, [(name, F_BODY, True, C_TITLE)])
        text(s, 11.62, y2 + 0.58, 7.60, 0.32, [(note, 14.5, False, C_MUTED)])
        text(s, 16.40, y2 + 0.28, 2.90, 0.42, [(val, 22, True, OK)])
        y2 += 1.16

    box(s, 11.24, y2 + 0.26, 8.30, 2.64, C_PANEL, C_PANLL, 1.25)
    text(s, 11.62, y2 + 0.54, 7.60, 0.40,
         [('검증은 정해둔 34개 케이스로 합니다', F_H, True, C_TITLE)])
    for i, ln in enumerate([
            '그런데 그건 «한 턴짜리»라 부족했습니다.',
            '그래서 사용자 열세 명을 만들어 대화를 시켰더니,',
            '34개가 «다 통과하는 상태»에서 여덟 개를 더 잡았습니다.',
            '가장 최근 것은 200턴 × 2인입니다.']):
        text(s, 11.62, y2 + 1.10 + i * 0.38, 7.60, 0.34, [(ln, F_BODY, False, C_BODY)])
    return s


# ════════════════════════════════════════════════════════════════
#  F. 인프라와 배포 — 「이렇게 했다 · 이래서 그랬다」
# ════════════════════════════════════════════════════════════════
#  앞 판은 트러블 네 개를 나열했는데, 그중 CORS 항목이 **사실이 아니었다** —
#  「배포 주소를 환경변수로 추가할 수 있게 바꿨다」고 적혀 있으나 실제로는
#    _cors_origins() 가 54행·128행에 «두 번» 정의돼 있고, add_middleware 는
#    그 함수를 안 부르고 리터럴 목록을 쓴다. 즉 환경변수는 아무 효과가 없다.
#  배포가 도는 건 nginx 가 같은 도메인에서 /api 를 프록시해 **CORS 가 발동을 안 하기** 때문이다.
#  ⇒ 틀린 항목을 고치는 대신, 장 전체를 «했다/왜»로 다시 짠다(사용자 지시 2026-08-09).
#  ⚠ 아래 값은 전부 서버 실측이다 — ps · ss · nginx sites-enabled 를 직접 봤다.
STACK = [
    ('브라우저',  'eum.r-e.kr',                  '',                          C_PANEL, C_PANLL),
    ('nginx',    'EC2 · :443 · Let\'s Encrypt',
     '/ → dist(정적 빌드)   ·   /api/ → 127.0.0.1:8000',                       C_CARD, C_CARDL),
    ('FastAPI',  'uvicorn · 127.0.0.1:8000',
     '인터넷에서 «안 보인다». nginx 만 볼 수 있다',                              OKBG, C_CARDL),
]
STORE = [('RDS · MySQL', '정형 데이터'), ('Pinecone', '벡터 · 3,072차원')]

DEPLOY = [
    ('백엔드가 인터넷에 그대로 보이면 안 된다',
     'nginx 만 :443 을 열고, uvicorn 은 127.0.0.1 에만 붙였습니다.'),
    ('키가 코드에 남으면 안 된다',
     '환경변수 25개를 서버 .env 로 주입합니다. git 에는 안 올라갑니다.'),
    ('첫 사용자가 혼자 기다리면 안 된다',
     '기동할 때 검색 연결을 미리 데웁니다 — 11초 → 4초.'),
    ('비밀키가 없어도 서버가 뜨면 안 된다',
     '없으면 기동을 막습니다. 없어도 뜨고 로그인만 500 이 났습니다.'),
]

LEFT = [
    ('아직 못 한 것도 적어 둡니다', ''),
    ('프로세스 관리자가 없습니다',
     '지금은 손으로 띄운 프로세스라 죽으면 안 돌아옵니다. 유닛은 써 뒀습니다.'),
    ('패키지 버전이 안 묶여 있습니다',
     '17개 전부 버전이 없어서 같은 환경이 재현되지 않습니다.'),
]


def slide_deploy(prs, s):
    head(s, '09 · 인프라와 배포', '이렇게 올렸고, 이래서 그렇게 했습니다')

    #  ── 왼쪽 : 구성 ──────────────────────────────────
    y = 2.62
    for name, sub, note, bg, ln in STACK:
        h = 1.16 if note else 0.94
        box(s, 1.00, y, 9.46, h, bg, ln, 1.1)
        text(s, 1.40, y + 0.16, 4.00, 0.40, [(name, 24, True, C_TITLE)])
        text(s, 5.10, y + 0.24, 5.10, 0.34, [(sub, F_SMALL, False, C_MUTED)])
        if note:
            text(s, 1.40, y + 0.66, 8.60, 0.34, [(note, F_BODY, True, C_PRES)])
        y += h + 0.34
        if name != 'FastAPI':
            seg(s, 5.72, y - 0.30, 5.72, y - 0.04)

    for i, (name, sub) in enumerate(STORE):
        x = 1.00 + i * 4.86
        box(s, x, y, 4.60, 0.86, C_PANEL, C_PANLL, 1.0)
        text(s, x + 0.34, y + 0.14, 4.00, 0.36, [(name, F_BODY, True, C_TITLE)])
        text(s, x + 0.34, y + 0.50, 4.00, 0.30, [(sub, F_SMALL, False, C_MUTED)])

    #  ── 왼쪽 아래 : 아직 못 한 것 ─────────────────────
    y2 = y + 1.16
    text(s, 1.00, y2, 9.46, 0.40, [(LEFT[0][0], F_H, True, HOT)])
    y2 += 0.50
    for name, note in LEFT[1:]:
        box(s, 1.00, y2, 9.46, 0.86, HOTBG, C_CARDL, 1.0)
        text(s, 1.36, y2 + 0.12, 8.80, 0.34, [(name, F_BODY, True, C_TITLE)])
        text(s, 1.36, y2 + 0.50, 8.80, 0.30, [(note, 14.5, False, C_MUTED)])
        y2 += 0.94

    #  ── 오른쪽 : 생각한 것 → 한 것 ────────────────────
    text(s, 11.24, 2.62, 8.30, 0.40, [('이건 이렇게 생각했습니다', F_H, True, C_TITLE)])
    y3 = 3.20
    for think, did in DEPLOY:
        box(s, 11.24, y3, 8.30, 1.62, C_CARD, C_CARDL, 1.1)
        text(s, 11.60, y3 + 0.18, 7.60, 0.40, [(think, F_BODY, True, C_TITLE)])
        text(s, 11.60, y3 + 0.68, 7.60, 0.34, [('↓', F_SMALL, False, C_MUTED)])
        text(s, 11.60, y3 + 1.06, 7.60, 0.36, [(did, F_SMALL, False, C_BODY)])
        y3 += 1.76
    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    print(f'  원본 {n0}장  ({prs.slide_width / IN:.2f} × {prs.slide_height / IN:.2f} in)')

    #  ── v18 0-based 색인 (실측으로 확인함) ─────────────────────
    #    17 잇다 표지 · 18 무엇을 만들었나 · 19 데이터
    #    20 예외처리(관문) ← 비우고 새로 · 21 관문 코드
    #    22 슬롯          ← 비우고 새로 · 23 슬롯 코드
    #    24 찾기 · 25 답 내기 · 26~28 문제와 해결 ×3
    #    29~37 나누다 · 38~43 공통 뒤 · 44 부록 검사결과 · 45 착지판정 · 46 신호9개
    slide_gate(prs, wipe(prs, 20))
    slide_slot(prs, wipe(prs, 22))
    slide_deploy(prs, wipe(prs, 40))     # 40 = 「09 · 인프라와 배포」 (v18 실측)
    C, E, D = n0, n0 + 1, n0 + 2
    slide_why_slot(prs)
    slide_measured(prs)
    slide_backup(prs)

    ITDA = [17, 20, 22, C, 24, 25, 26, E]        # 발표할 잇다 8장
    #  ⚠ 안 쓰는 잇다 장은 «지우지 않는다». 지우면 파트가 남아 파일이 죽는다(위 SRC 주석).
    #    부록 쪽으로 «밀어» 둔다 — 필요하면 그 자리에서 꺼내 쓸 수 있다.
    SHELF = [18, 19, 21, 23, 27, 28, 45, 46]

    KEEP = (list(range(0, 17))                   # 공통 앞 + 덜다
            + ITDA
            + list(range(29, 45))                # 나누다 + 공통 뒤 + 부록 검사결과
            + SHELF
            + [D])                               # 백업은 «맨 뒤»

    assert len(set(KEEP)) == len(KEEP), '중복'
    assert set(KEEP) == set(range(n0 + 3)), \
        f'빠진 인덱스: {sorted(set(range(n0 + 3)) - set(KEEP))}'

    reorder(prs, KEEP)
    prs.save(OUT)

    out = Presentation(OUT)
    import zipfile
    nz = len([n for n in zipfile.ZipFile(OUT).namelist()
              if n.startswith('ppt/slides/slide') and n.endswith('.xml')])
    print(f'  → {len(out.slides)}장 · zip 슬라이드 파트 {nz}개'
          f'  {"✅ 유령 없음" if nz == len(out.slides) else "🔴 유령 " + str(nz - len(out.slides)) + "장"}')
    for i, sl in enumerate(out.slides, 1):
        if 18 <= i <= 25 or i == len(out.slides):
            t = next((sh.text_frame.text.split('\n')[0][:44]
                      for sh in sl.shapes
                      if sh.has_text_frame and sh.text_frame.text.strip()), '')
            t2 = [sh.text_frame.text.split('\n')[0][:40] for sh in sl.shapes
                  if sh.has_text_frame and sh.text_frame.text.strip()]
            print(f'   {i:3d}장 도형{len(sl.shapes):3d}  {t:<30s} | '
                  f'{t2[2][:36] if len(t2) > 2 else ""}')
    print(f'  저장 → {OUT}')


main()
