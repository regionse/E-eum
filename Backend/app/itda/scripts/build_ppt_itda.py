# -*- coding: utf-8 -*-
"""잇다 발표 슬라이드 생성 — 기존 디자인을 «복제»해서 글자만 바꾼다.

왜 이렇게 하나
  이 덱은 레이아웃 자리표시자가 없고 슬라이드마다 도형을 직접 놓은 구조다.
  그래서 python-pptx 로 새 슬라이드를 «만들면» 디자인이 통째로 날아간다.
  ⇒ 기존 슬라이드의 XML 을 그대로 복제하고 텍스트만 교체한다.

쓰는 법
  python build_ppt_itda.py <원본.pptx> <출력.pptx>

원고: 발표원고_잇다.md  (내용을 고치려면 아래 SPEC 을 고친다)
"""
import sys
import io
import copy

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션 수정안_잇다완성.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션 수정안_잇다v2.pptx'

#  원본에서 «잇다» 구간 (표지 16 ~ 문제와해결 26). 이 구간을 통째로 갈아끼운다.
ITDA_FROM, ITDA_TO = 16, 26
#  본보기로 쓸 슬라이드 (1-based, 원본 기준)
TPL_COVER = 16      # 잇다 표지 — 큰 제목 + 항목 3개
TPL_A = 18          # 좌 3카드 + 우 강조패널
TPL_B = 26          # 2단 카드 (TROUBLE)
TPL_DEMO = 25       # 시연


# ── 슬라이드 복제 ───────────────────────────────────────────────
def clone(prs, idx0):
    """idx0(0-based) 슬라이드를 복제해 «맨 뒤»에 붙이고 새 슬라이드를 준다."""
    src = prs.slides[idx0]
    dst = prs.slides.add_slide(src.slide_layout)
    #  add_slide 가 넣은 빈 자리표시자를 제거한다
    for sh in list(dst.shapes):
        sh._element.getparent().remove(sh._element)
    for sh in src.shapes:
        dst.shapes._spTree.append(copy.deepcopy(sh._element))
    #  그림·차트 등 관계(rel)도 따라가야 한다
    for k, rel in src.part.rels.items():
        if rel.reltype.endswith('slideLayout'):
            continue
        if rel.is_external:
            dst.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            dst.part.rels.get_or_add(rel.reltype, rel._target)
    return dst


def reorder(prs, order):
    """order(0-based 인덱스 목록) 순서로 슬라이드를 재배열한다."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for x in ids:
        lst.remove(x)
    for i in order:
        lst.append(ids[i])


_RID = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'


def drop(prs, idx0s):
    """슬라이드 삭제 (0-based).

    ⚠ `sldIdLst` 에서 «빼기만» 하면 안 된다. 슬라이드 부품(part)이 파일 안에 그대로
      남아서, PowerPoint 는 45장으로 보는데 파일에는 56장이 들어 있는 상태가 된다.
      실제로 처음에 그렇게 만들었다. **관계(rel)까지 끊어야** 저장할 때 빠진다.
    """
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for i in sorted(idx0s, reverse=True):
        prs.part.drop_rel(ids[i].get(_RID))
        lst.remove(ids[i])


# ── 텍스트 교체 (서식 보존) ─────────────────────────────────────
def put(shape, text):
    """도형의 글자를 바꾼다. **첫 run 의 서식을 살리고** 나머지 run 은 지운다.

    줄바꿈('\n')은 문단으로 나눈다 — 첫 문단 서식을 그대로 복사해서 쓴다.
    """
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    paras = tf.paragraphs
    p0 = paras[0]
    if not p0.runs:                      # 빈 도형이면 건드리지 않는다
        return
    lines = str(text).split('\n')

    #  첫 문단 — run 하나만 남기고 글자 교체
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    p0.runs[0].text = lines[0]

    #  나머지 문단 제거
    for p in list(paras[1:]):
        p._p.getparent().remove(p._p)

    #  추가 줄을 문단으로 복제해 붙인다 (서식 유지)
    for ln in lines[1:]:
        newp = copy.deepcopy(p0._p)
        p0._p.getparent().append(newp)
        from pptx.text.text import _Paragraph
        np = _Paragraph(newp, tf)
        for r in list(np.runs[1:]):
            r._r.getparent().remove(r._r)
        np.runs[0].text = ln


def fill(slide, mapping):
    """{도형 인덱스: 글자} 로 채운다. None 이면 건너뛴다."""
    shapes = list(slide.shapes)
    for i, txt in mapping.items():
        if txt is None or i >= len(shapes):
            continue
        put(shapes[i], txt)


# ── 본보기별 도형 인덱스 ────────────────────────────────────────
#  A(18): 0 헤더 · 1 발표자 · 2 제목
#         3 카드1배경 · 4 카드1제목 · 5 카드1본문
#         6 카드2배경 · 7 카드2제목 · 8 카드2본문
#         9 카드3배경 ·10 카드3제목 ·11 카드3본문
#        12 우패널배경 ·13 우제목 ·14 우본문
def slide_a(prs, head, title, c1, c2, c3, right):
    s = clone(prs, TPL_A - 1)
    fill(s, {0: head, 2: title,
             4: c1[0], 5: c1[1], 7: c2[0], 8: c2[1], 10: c3[0], 11: c3[1],
             13: right[0], 14: right[1]})
    return s


#  B(26): 0 헤더 · 1 발표자 · 2 제목
#         3 좌배경 · 4 좌번호 · 5 좌제목 · 6 선 · 7,8,9 좌본문
#        10 우배경 ·11 우번호 ·12 우제목 ·13 선 ·14,15,16 우본문
def slide_b(prs, head, title, left, right):
    s = clone(prs, TPL_B - 1)
    fill(s, {0: head, 2: title,
             4: left[0], 5: left[1], 7: left[2], 8: left[3], 9: left[4],
             11: right[0], 12: right[1], 14: right[2], 15: right[3], 16: right[4]})
    return s


# ════════════════════════════════════════════════════════════════
#  내용
# ════════════════════════════════════════════════════════════════
H = 'PART 03 · 잇다'


def build(prs):
    made = []

    # ── 표지 ─────────────────────────────────────────────────
    #  0 PART · 1 「잇다」 · 2 부제 · 3 패널배경 · 4 발표자 · 5 선 · 6,7,8 항목
    s = clone(prs, TPL_COVER - 1)
    fill(s, {2: '흔한 재료로, 흔치 않게',
             6: '1 · 답할까 물어볼까 — 후보의 «흩어짐»으로',
             7: '2 · 맞고 있는 사람을 가해자로 만들던 것',
             8: '3 · 논문 16편 — 못 쓰는 부분은 «번역»했다'})
    made.append(s)

    # 1. 기능 — 원본 17 재사용
    made.append(clone(prs, 17 - 1))

    # 2. 시스템 구조
    made.append(slide_a(
        prs, f'{H} — 시스템 구조', '대화 한 턴이 도는 길',
        ('① 입력 방어 — LLM 0회 · 0원',
         '거부·번복 / 정책거절 / 남용 / 불법신호 계측 / 위기·폭력 / 칩 선택 / 인젝션 / 지름길 3종. 여기서 걸리면 LLM 을 한 번도 안 부른다.'),
        ('② 대화 판정 — LLM #1',
         'Gemini · responseSchema 로 슬롯 6축 + action + 검색어 + 답변을 한 번에 받는다. 자유 텍스트를 파싱하지 않는다.'),
        ('③ 슬롯 후처리 — 8단계, 전부 코드',
         '근거 대조 → 무의미어 제거 → 빼기 → 쌓기 → 다루는대상 보완 → 「누구를」 보완 → 턴 기록 → 사다리 리셋.'),
        #  ⚠ 우패널 제목은 «한 줄»이어야 한다 — 상자 W7.4" · 34pt 라 15자를 넘으면
        #    두 줄이 되어 H1.0" 를 넘친다(overflow.py 로 확인).
        ('검색 → 판정 → 출력 검사',
         '벡터 + 키워드 → RRF → 크로스인코더.\n2단 판정이 방향 칩 / 세부 칩 / 카드를 가른다.\n나가는 말은 코드가 한 번 더 본다.\n\n턴당 LLM 1.1회 · 턴당 ≈1.2원 · 캐시 60%')))

    # 3. 기술 스택
    made.append(slide_a(
        prs, f'{H} — 기술 스택', '재료는 평범합니다',
        ('LLM · 임베딩',
         'Gemini 3.1 Flash-Lite · REST 직접(urllib) · responseSchema 로 JSON 강제. 임베딩은 gemini-embedding-2 · 3,072차원.'),
        ('검색 · 융합 · 리랭크',
         'Pinecone(ns 3개 · 코사인) + MySQL FULLTEXT(ngram) → RRF k=60 → Jina v3.5 → Pinecone bge-v2-m3 폴백.'),
        ('서버 · 데이터',
         'FastAPI · SQLAlchemy async · EC2 + RDS. NCS 직업 1,094(태깅 1,064) · 자격증 613 + 시험일정 2,655 · 강좌 8,371.'),
        ('여기까지는 특별한 게 없습니다',
         '하이브리드 검색도 RRF 도 누구나 씁니다.\n덜다도 씁니다.\n\n다음 장이 본론입니다.\n\n※ 잇다는 LangGraph 를 안 씁니다 — 덜다만.')))

    # 4. ★ 대조표
    made.append(slide_b(
        prs, f'{H} — 쓰는 법', '같은 재료를 어디에 놓느냐가 다릅니다',
        ('보통 이렇게', '업계에서 흔한 방식',
         '안전 — LLM 을 먼저 부르고 나중에 거르거나, moderation API 를 따로 호출한다 (비용 + 지연 추가).',
         '구조화 출력 — 모델이 낸 JSON 을 그대로 믿는다. 마음이 바뀌면 덮어쓰거나 무한 누적한다.',
         '임계값 — 감으로 찍고 다시 안 본다. 폐기한 시도는 지운다. 월 청구서는 알아도 턴당 단가는 모른다.'),
        ('우리는', '같은 재료, 다른 자리',
         '안전 — LLM 앞에서 코드가 끊는다. 걸리면 호출 0회 · 0원. 위기 발화는 차단하지 않고 ☎109 로 «잇고», 남용 카운터에 절대 안 센다.',
         '구조화 출력 — {값, 근거} 를 함께 받아 근거를 원문과 대조한다. 없으면 버린다. 마음이 바뀌면 «빼기» 연산으로 덜어낸다.',
         '임계값 — n=52 실측의 p25=0.135 → RR_LOW_CONF=0.13. 폐기한 시도도 기록한다. 턴당 1.2원 · LLM 1.1회를 잰다.')))

    # 5. 슬롯
    made.append(slide_a(
        prs, f'{H} — 슬롯', '무엇을 뽑고, 어떻게 믿나',
        ('7축 — 6개는 LLM 이, 1개는 코드가',
         '관심분야 · 활동유형 · 다루는대상 · 제약(다중값 ≤4) / 세부관심 · 강점성향(단일) / 대상세부는 코드가 채운다.'),
        ('근거를 함께 받아 원문과 대조한다',
         '{"값": "제빵", "근거": "빵 만드는 게 재밌었어요"} — 근거가 발화에 없으면 그 슬롯을 버린다. 다중값은 항목마다 따로 대조한다.'),
        ('그리고 코드가 8단계로 손본다',
         '근거 대조 → 무의미어 제거 → 빼기(DELETE) → 쌓기(덮지 않는다) → 보완 2종 → 턴 기록 → 사다리 리셋. 전부 LLM 0회.'),
        ('왜 프롬프트로 안 부탁하나',
         '모델은 매 턴 일관되게 해 주지 않는다.\n\n실측 — 스키마에 필드를 하나 더했더니\n골든셋 「지목」이 7/7 → 4/7 로 떨어졌다.')))

    # 6. 검색
    made.append(slide_a(
        prs, f'{H} — 검색', '왜 벡터만으로는 안 되나',
        ('직업 임베딩 본문이 «21자»다',
         '「직업명 · 중분류」뿐이라 뜻이 얇다. 그래서 벡터 단독으로는 이름을 못 잡는다 — 하이브리드의 이유다.'),
        ('벡터 + 키워드 → RRF (k=60)',
         '벡터는 뜻에 강하고 이름에 약하다. FULLTEXT 는 반대다. RRF 는 점수가 아니라 «순위»를 합쳐 공통으로 오르는 후보를 올린다.'),
        ('크로스인코더가 임베딩과 다른 점',
         '임베딩은 질의와 후보를 «따로» 벡터로 만들어 거리를 재고, 리랭커는 둘을 «붙여서 함께» 읽는다. 정확한 대신 느려 미리 못 만든다.'),
        ('질의는 세 갈래로 만든다',
         'LLM query + query_alts 2개\n+ 누적 슬롯을 «닻»으로 하나 더\n\n추가 호출 «없이» 같은 응답으로 받는다\n(Kostric & Balog, SIGIR 2024)')))

    # 7. 착지 판정
    made.append(slide_a(
        prs, f'{H} — 착지 판정', '언제 답하고, 언제 되묻나',
        ('① 조건이 찼나 — «무게»로 잰다',
         '사용자가 말한 축 1.0 · 코드가 추론한 축 0.5 · 합 2.0 이상이어야 답한다(12턴 넘으면 1.5). 「공짜 한 축」이 반값이 된다.'),
        ('한 문장에 다 찼으면 한 번 미룬다',
         '「빵 만드는 거 좋아해요」 한 마디에 [제빵] 카드가 나가던 것. 조건은 찼지만 «대화를 한 적이 없다».'),
        ('② 후보가 흩어졌나 — 2단으로',
         '1단 TOP%(1위 묶음 점유율) < 0.90 → 방향 칩. 2단 엔트로피 ≥ 0.70 → 세부 칩. 그 아래면 카드.'),
        ('③ 확신이 없으면 한 층 위로',
         '리랭커 최고점 < 0.13 이면 중분류를 «얹되» 잎사귀는 안 버린다.\n\n「보건복지 쪽입니다」는 「요양보호사 · 응시제한 없음 · 다음 시험 8월 27일」보다 쓸모없다.')))

    # 8. ★ 논문 번역
    made.append(slide_b(
        prs, f'{H} — 근거', '논문을 «번역»했습니다',
        ('PAPER', '학습 요소는 우리가 못 씁니다',
         'Dey 외 · ACL 2025 — 대화상태 갱신 앞에 «학습된 이진 분류기»를 둔다 (JGA 67.1 → 70.5).',
         'HAConvDR · ACL Findings 2024 — 이력의 유용성을 «검색에 준 영향»으로 학습해 가려낸다.',
         'SOCbot · Sturgis 외 2025 — shortlist → 충분판단 → 선택 → 설명, 네 단계로 직업을 코딩한다.'),
        ('OURS', '그래서 값싼 «대리»를 놓았습니다',
         '분류기가 없으니 «행동 신호(action)»로 대신한다 — 거절·이탈 턴에는 슬롯을 안 건드린다.',
         '학습으로 못 하므로 «RRF 의 합의»로 대신한다. 약하지만 공짜다.',
         '우리는 1·4만 하고 있었다 → «충분판단»을 추가했다. 논문 16편 · 학회 14곳을 이렇게 읽었다.')))

    # 9~12. TROUBLE
    made.append(slide_b(
        prs, f'{H} — 문제와 해결', '사용자를 막아버린 것',
        ('TROUBLE 01', '맞고 있는 사람을 가해자로 만들었다',
         '원인 — 「치매 할머니가 자꾸 저를 때려요」가 차단되고 남용 카운트까지 올랐다. 폭력 낱말표는 «누구의 행동인지»를 모른다.',
         '해결 — 낱말을 «차단»이 아니라 «라우팅»으로. 걸리면 LLM 에게 가해·피해·제3자만 묻고, 피해면 ☎129 와 나누다 지도로 «잇는다».',
         '배운 것 — 실측 21케이스 중 9건이 오차단이었다. 낱말이 할 수 있는 건 「여기 봐야 한다」까지다.'),
        ('TROUBLE 02', '낱말을 다듬어도 다음에 또 났다',
         '원인 — 「장애인 활동지원사」의 \'애인\', 「웹 개발 기초」의 \'발기\'. 공백을 지운 문자열에서 낱말을 찾고 있었다. 같은 버그가 코드에 세 번 기록돼 있었다.',
         '해결 — 낱말이 아니라 «매칭 방식». 우회는 어절 «사이»를 쪼개고 오탐은 어절 «가운데»서 붙는다 → 어절 시작에서만 인정한다.',
         '배운 것 — 낱말표를 한 글자도 안 건드렸는데 오차단 16건이 사라졌다. 「이번 건 고쳤다」와 「이 종류가 안 나게 했다」는 다르다.')))

    made.append(slide_b(
        prs, f'{H} — 문제와 해결', '도구의 한계를 논문이 이미 재놨다',
        ('TROUBLE 03', '조용한 걸 원했더니 «소음 측정»',
         '원인 — 「조용한 데서 하는 일」에 [소음진동측정·분석평가]가 나왔다. NevIR(EACL 2024): "IR models do not consider negation, performing the same or worse than a random ranking."',
         '해결 — 이기려 하지 않고 «안 넣기»로 했다. 조건 명사와 조건 표지가 둘 다 있을 때만 검색어에서 뺀다(「소음진동을 측정하는 일」은 안 빠지게).',
         '배운 것 — 논문이 지목한 세 모델 유형이 정확히 우리 세 층이었다. 못 하는 걸 시키는 대신 «그 일을 안 시키는 게» 맞을 때가 있다.'),
        ('TROUBLE 04', '절반쯤 틀리는 신호에 처벌을 얹었다',
         '원인 — 주제이탈로 판정되면 남용 카운터를 올렸다. 그런데 OOS 재현율은 Claude Haiku 0.619 · Mistral 0.453 · SetFit 0.462 다.',
         '해결 — 이탈은 «세지 않는다». 그리고 판별기를 이진에서 6분류로 넓혔다(진로·사정·못정함·되묻기·착지요청·이탈). ISO 24617-2 — 한 발화는 여러 차원에서 동시에 기능한다.',
         '배운 것 — «신호의 정확도가 처벌의 세기를 정해야 한다». 0.5짜리 신호로 세션을 닫으면 안 된다.')))

    made.append(slide_b(
        prs, f'{H} — 문제와 해결', '만들고, 재보고, 되돌렸다',
        ('TROUBLE 05', '콕 집어 말한 사람에게 되물었다',
         '원인 — 「용접 일을 하고 싶어요」에 엔트로피 0.887 → 되묻기. 그런데 후보가 전부 «용접»이었다. Farquhar 외(Nature 2024): 엔트로피는 「같은 것을 다르게 말한 것」과 「다른 것 여럿」을 못 가른다.',
         '해결 — 두 신호를 AND 로 붙여 14/14 만점을 얻었다. 그런데 «되돌렸다» — 정답표가 틀렸기 때문이다. 되묻기는 「어느 동네」와 「그 동네 어디」 두 종류였다.',
         '배운 것 — 만점이 정답의 증거가 아니다. 정답표를 먼저 의심해야 했다.'),
        ('TROUBLE 06', '좋아 보이는 아이디어 셋을 껐다',
         '가중 RRF — 순위가 «안 바뀌었다». k=60이면 1/61과 1/66이 거의 같아 RRF 는 사실상 «표 세기»다. 무게로는 못 이긴다.',
         '자기일관성 N=3 — 흔들림 1건이 «그대로»였고 비용만 3배(0.50 → 1.46원)였다. 껐다.',
         '배운 것 — 문턱도 «스윕»으로 정했다. 0.05~0.30 전부 14/14, 0.32 에서 무너짐. 정확도 하나가 아니라 「임계값을 어디에 둬도 되는 폭」을 본다.')))

    made.append(slide_b(
        prs, f'{H} — 문제와 해결', '조용히 나빠지는 것',
        ('TROUBLE 07', '리랭커가 죽었는데 며칠간 몰랐다',
         '원인 — Pinecone 무료 등급은 리랭크 월 500회다. 소진되자 조용히 폴백으로 돌았고 경고는 프로세스당 한 번만 찍혔다.',
         '무엇이 같이 꺼졌나 — 「우리 목록에 없어요」 판정 · 저확신 표시 · 강좌 무관 제거. 셋 다 리랭커 점수에 매달려 있었다. 응답은 계속 «정상처럼» 나갔다.',
         '배운 것 — 터지는 실패보다 «조용한 실패»가 무섭다. 500 은 보이지만 「200인데 내용이 나쁜 것」은 안 보인다. 폴백 자체를 세야 한다.'),
        ('TEST', '그래서 검사를 «생성»한다',
         '손으로 쓴 검사 325건 — 그중 283건은 LLM 0회 · 0원이라 언제든 다시 돌린다. 골든셋 34 · 코드단위 90 · 다중턴 8 · 오탐 35 · 우회 32 …',
         '생성한 검사 22,470건 — DB 의 직업·자격증·강좌명 3,745건 × 문장틀 6개. «내가 고른 문장이 아니다». 낱말표가 늘면 검사도 자동으로 같이 는다.',
         '그리고 「34/34 통과」를 안전 주장으로 안 쓴다 — 3의 법칙상 참 실패율 8.8%까지 허용되는 진술이다.')))

    # 16. 인증
    made.append(slide_a(
        prs, f'{H} — 인증', '토큰이 아니라 «소유권»이 문제였다',
        ('로그인 — 표준대로',
         'bcrypt 로 해싱한다. 평문은 어디에도 안 남긴다. PyJWT · HS256 · 24시간.'),
        ('토큰이 유효해도 다시 본다',
         '계정 상태(정지·휴면·탈퇴)를 매 요청 확인한다 — 발급 시점의 권한을 그대로 믿지 않는다.'),
        ('session_id 는 «프론트»가 만든다',
         '그래서 남의 session_id 를 알면 남의 대화 상태를 받아갈 수 있었다. 제약 슬롯에 학력·체력·비용 부담이 들어 있다.'),
        ('세션에 owner 를 찍는다',
         '저장·이어서하기뿐 아니라\n«대화 자체»에 검증을 걸었다.\n\n인증의 구멍은 토큰이 아니라\n«누구 것인가»를 안 물은 데 있었다.')))

    # 17. 시연 — 원본 25 재사용
    made.append(clone(prs, 25 - 1))

    return made


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    made = build(prs)
    print(f'  새 잇다 슬라이드 {len(made)}장 생성')

    #  ① 원본 잇다 구간(16~26)을 «관계까지» 지운다
    drop(prs, list(range(ITDA_FROM - 1, ITDA_TO)))
    n_old_before = ITDA_FROM - 1                 # 앞부분(1~15) 장수
    n_old_after = n0 - ITDA_TO                   # 뒷부분(27~끝) 장수
    print(f'  원본 잇다 {ITDA_TO - ITDA_FROM + 1}장 제거 · '
          f'앞 {n_old_before} + 뒤 {n_old_after} 유지')

    #  ② 남은 순서: [앞 15] + [뒤 15] + [새 잇다 15]  →  [앞] + [새] + [뒤] 로
    a = list(range(0, n_old_before))
    b = list(range(n_old_before, n_old_before + n_old_after))
    c = list(range(n_old_before + n_old_after,
                   n_old_before + n_old_after + len(made)))
    reorder(prs, a + c + b)

    prs.save(OUT)
    print(f'  저장 → {OUT}')
    print(f'  전체 {len(Presentation(OUT).slides)}장')


main()
