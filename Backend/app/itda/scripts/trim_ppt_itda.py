# -*- coding: utf-8 -*-
"""잇다 구간을 «프로젝트 소개» 톤으로 정리한다 — 13장 → 10장. (2026-08-07)

왜 하나 — 사용자 지시
  「'보통 이렇게' 같은 비교는 하지 마라. 프로젝트를 소개하는 거다, 전체적으로.
   논문 근거 같은 것도 하지 마라. 국비훈련생 프로젝트다. 그건 나중에 따로 하면 된다.」

  덱을 다시 보니 잇다만 13장이었다(덜다 8 · 나누다 9). 연구 발표 형태로 커져 있었다.
  ⇒ 「무엇을 만들었나 · 어떻게 도나 · 뭘 썼나 · 뭐가 막혔나」로 되돌린다.

무엇을 하나
  ① 삭제 3장   비교(쓰는 법) · 논문 근거 · TROUBLE 05·06 장
                 └ 05(엔트로피)는 Nature 논문이 근거의 전부라 빼면 남는 게 없다.
                   06(껐다)은 논문이 필요 없고 「재보고 껐다」가 강해서 아래로 «옮긴다».
  ② 논문 인용 제거   SIGIR · NevIR 영어 인용문 · OOS 재현율 수치 · ISO 24617-2 · 3의 법칙
  ③ 순서 교정   표지 → «기능» → 구조 → 스택   (기술 스택이 기능보다 앞에 있었다)
  ④ TROUBLE 번호 01·02·03·06·07 → 01~05 로 이어 붙임 (건너뛰면 빠뜨린 것처럼 보인다)
  ⑤ 표지 예고 3줄 재작성 — 지워진 내용을 가리키고 있었다

⚠ 잇다 «밖»은 건드리지 않는다. 덜다·나누다·공통 장은 그대로다.

쓰는 법
  python trim_ppt_itda.py <입력.pptx> <출력.pptx>
"""
import sys
import io
import copy

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v3.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v4.pptx'

_RID = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id'


def drop(prs, idx0s):
    """슬라이드 삭제 (0-based). 관계(rel)까지 끊어야 파일에서 실제로 빠진다."""
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for i in sorted(idx0s, reverse=True):
        prs.part.drop_rel(ids[i].get(_RID))
        lst.remove(ids[i])


def reorder(prs, order):
    lst = prs.slides._sldIdLst
    ids = list(lst)
    for x in ids:
        lst.remove(x)
    for i in order:
        lst.append(ids[i])


def put(shape, text):
    """도형의 글자를 바꾼다. 첫 run 의 서식을 살리고 나머지 run 은 지운다."""
    if shape is None or not shape.has_text_frame:
        return
    tf = shape.text_frame
    paras = tf.paragraphs
    p0 = paras[0]
    if not p0.runs:
        return
    lines = str(text).split('\n')
    for r in list(p0.runs[1:]):
        r._r.getparent().remove(r._r)
    p0.runs[0].text = lines[0]
    for p in list(paras[1:]):
        p._p.getparent().remove(p._p)
    for ln in lines[1:]:
        newp = copy.deepcopy(p0._p)
        p0._p.getparent().append(newp)
        from pptx.text.text import _Paragraph
        np = _Paragraph(newp, tf)
        for r in list(np.runs[1:]):
            r._r.getparent().remove(r._r)
        np.runs[0].text = ln


def S(prs, n1):
    """1-based 슬라이드 번호로 도형 목록을 준다."""
    return list(prs.slides[n1 - 1].shapes)


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    print(f'  원본 {n0}장')

    # ══════════════════════════════════════════════════════════
    #  ① 표지(18) — 예고 3줄이 «지워질 내용»을 가리키고 있었다
    #     옛것: 「모델이 낸 값을 원문과 대조 / 하이브리드 검색과 2단 판정 / 측정으로 뒤집은 설계 결정」
    #     → 소개 톤으로. 사용자가 무엇을 얻는지부터 말한다.
    # ══════════════════════════════════════════════════════════
    c = S(prs, 18)
    put(c[6], '1 · 물어보고 좁혀서 «직업 방향 카드»까지')
    put(c[7], '2 · 자격증·강좌를 붙여 «미래설계지도»로')
    put(c[8], '3 · 만들면서 부딪히고 고친 것들')

    # ══════════════════════════════════════════════════════════
    #  ② 검색(24) — SIGIR 인용 제거
    # ══════════════════════════════════════════════════════════
    put(S(prs, 24)[14],
        'LLM 이 만든 검색어 3개에\n'
        '누적 슬롯을 «닻»으로 하나 더.\n'
        '\n'
        '검색어를 따로 부르지 않는다 —\n'
        '답변과 «같은 응답»에 실려 온다.')

    # ══════════════════════════════════════════════════════════
    #  ③ TROUBLE 장 재편
    #     27 : 01·02  그대로 (논문 없음)
    #     28 : 03 + 04(옛 06 「껐다」)  ← 재작성. NevIR 영어 인용문·OOS 수치·ISO 제거
    #     29 : 삭제
    #     30 : 05(옛 07) + TEST        ← 번호만 고치고 「3의 법칙」 제거
    # ══════════════════════════════════════════════════════════
    #  TPL_B 인덱스: 0 헤더 · 2 제목 · 4 좌번호 · 5 좌제목 · 7,8,9 좌본문
    #                          11 우번호 ·12 우제목 ·14,15,16 우본문
    b = S(prs, 28)
    put(b[2], '엉뚱하게 찾고, 만들어놓고 껐습니다')
    # ── 좌: TROUBLE 03 — 부정 조건. 논문 인용문을 «현상 설명»으로 바꾼다
    put(b[4], 'TROUBLE 03')
    put(b[5], '조용한 걸 원했더니 «소음 측정»')
    put(b[7], '원인 — 「조용한 데서 하는 일」에 [소음진동측정·분석평가]가 1위로 나왔다. '
              '검색은 「조용한」을 «빼야 할 조건»이 아니라 «찾을 주제»로 읽는다.')
    put(b[8], '해결 — 이기려 하지 않고 «안 넣기»로 했다. 조건을 뜻하는 말과 그 표지가 '
              '둘 다 있을 때만 검색어에서 뺀다 — 「소음진동을 측정하는 일」은 안 빠지게.')
    put(b[9], '배운 것 — 도구가 못 하는 일이 있다. 억지로 시키는 것보다 '
              '«그 일을 안 시키는 게» 맞을 때가 있다.')
    # ── 우: TROUBLE 04 — 옛 06 「좋아 보이는 아이디어 셋을 껐다」
    put(b[11], 'TROUBLE 04')
    put(b[12], '좋아 보이던 것 셋을 껐다')
    put(b[14], '가중 RRF — 순위가 «안 바뀌었다». k=60 이면 1/61 과 1/66 이 거의 같아서 '
               'RRF 는 사실상 «표 세기»다. 무게로는 못 이긴다.')
    put(b[15], '자기일관성 N=3 — 세 번 물어 다수결. 흔들리던 1건이 «그대로»였고 '
               '비용만 3배(0.50 → 1.46원)였다. 껐다.')
    put(b[16], '배운 것 — 만들었다고 다 켜는 게 아니다. 문턱도 «훑어서» 정했다 — '
               '0.05~0.30 전부 14/14, 0.32 에서 무너졌다.')

    # ── 30: 번호 07 → 05 · 「3의 법칙」 제거
    d = S(prs, 30)
    put(d[4], 'TROUBLE 05')
    #  ⚠ 옛 문장: 「34/34 통과를 안전 주장으로 안 쓴다 — 3의 법칙상 참 실패율 8.8%…」
    #    통계 용어를 빼고, «실제로 겪은 일»로 바꾼다. 이쪽이 훨씬 세고 설명도 짧다.
    put(d[16], '그리고 「검사 통과 = 동작함」이 아니다. 실제로 골든셋 34/34 인데 '
               '카드가 «전부 죽어» 있던 적이 있다 — 검사가 그 경로를 안 탔기 때문이다.')

    # ══════════════════════════════════════════════════════════
    #  ④ 삭제 — 22(비교) · 26(논문) · 29(트러블 05·06)
    # ══════════════════════════════════════════════════════════
    drop(prs, [22 - 1, 26 - 1, 29 - 1])
    print('  삭제 3장 — 22 비교 · 26 논문 근거 · 29 TROUBLE 05·06')

    # 삭제 뒤 0-based 인덱스
    #   0..20  = 옛 1..21   (표지18=17 · 기능20=19 · 구조21=20 · 스택19=18)
    #   21..   = 옛 23,24,25,27,28,30,31...   (22·26·29 빠짐)
    #   옛23 슬롯→21 · 옛24 검색→22 · 옛25 착지→23 · 옛27 트A→24
    #   옛28 트B→25 · 옛30 트C→26 · 옛31~45 → 27~41
    order = (
        list(range(0, 17))     # 1~17   표지 ~ 덜다 끝
        + [17]                 # 18 잇다 표지
        + [19]                 # ★ 기능  (옛 20) — 스택보다 «앞»으로
        + [20]                 #    구조  (옛 21)
        + [18]                 #    기술 스택 (옛 19)
        + [21]                 #    슬롯  (옛 23)
        + [22]                 #    검색  (옛 24)
        + [23]                 #    착지 판정 (옛 25)
        + [24]                 #    TROUBLE 01·02 (옛 27)
        + [25]                 #    TROUBLE 03·04 (옛 28 · 재작성)
        + [26]                 #    TROUBLE 05+TEST (옛 30)
        + list(range(27, 42))  # 나누다 ~ 끝
    )
    assert len(order) == len(set(order)) == n0 - 3, \
        f'순서 목록 어긋남: {len(order)}개 · 중복 {len(order)-len(set(order))} · 기대 {n0-3}'
    reorder(prs, order)

    prs.save(OUT)
    print(f'  저장 → {OUT}')
    print(f'  전체 {len(Presentation(OUT).slides)}장  (잇다 10장)')


main()
