# -*- coding: utf-8 -*-
"""(3) 슬롯으로 «어떻게 찾고 어떻게 답을 내나» — 2장. (2026-08-07)

사용자가 정한 구조
  (1) 자연어를 받아 게이트로 거른다      ← 그림 완성
  (2) 통과한 것을 슬롯에 넣는다          ← 그림 완성
  (3) 그 슬롯으로 «어떻게 서치하고 결과를 낼지»   ← 이 파일
       3-a  어떻게 찾나
       3-b  어떻게 답을 내나

⚠ 이건 «프로젝트 소개»다. 내부 번호(①④⑤)나 「사실 이게 틀렸습니다」를 넣지 않는다.
  듣는 사람은 잇다를 처음 본다. 한 사람의 문장 하나가 카드가 되는 걸 따라가게 만든다.
  → 앞 그림들과 «같은 사람»(할머니 돌보는 사용자)으로 이어 붙인다.

⚠ 카드 내용은 «실제로 엔진을 돌려서» 받은 것이다. 지어내지 않았다 (2026-08-07 실행):
    3턴  「할머니 돌보느라 학교를 그만뒀어요」 → 「돌보는 건 익숙해요…」 → 「어르신 쪽이 편해요」
    → job 요양지원 (보건 · 06010108)
    → certs [] · no_cert_path=True  → 국민내일배움카드 안내로 갈린다
    → courses K-MOOC 「사회복지실천과 요양보호」 · 「와상환자를 위한 통합돌봄」

코드 조각도 «중요한 데만» 넣는다 — RRF 핵심 4줄. 전체를 붙이면 아무도 안 읽는다.
"""
import sys
import io

sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8')

from pptx import Presentation                                    # noqa: E402
from pptx.util import Emu, Pt                                    # noqa: E402
from pptx.dml.color import RGBColor                              # noqa: E402
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR            # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR                  # noqa: E402

SRC = sys.argv[1] if len(sys.argv) > 1 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v13.pptx'
OUT = sys.argv[2] if len(sys.argv) > 2 else \
    r'C:\Users\TFX255GS\Downloads\프레젠테이션_이음_v14.pptx'

IN = 914400
FONT, CODE = 'Arial', 'Consolas'
C_HEAD  = RGBColor(0x8B, 0x84, 0xA0)
C_PRES  = RGBColor(0x6B, 0x3F, 0xA0)
C_TITLE = RGBColor(0x3B, 0x20, 0x63)
C_BODY  = RGBColor(0x55, 0x50, 0x6B)
C_MUTED = RGBColor(0x8B, 0x84, 0xA0)
C_CARD  = RGBColor(0xFF, 0xFF, 0xFF)
C_CARDL = RGBColor(0xE7, 0xE2, 0xF0)
C_PANEL = RGBColor(0xEF, 0xEB, 0xF6)
C_PANLL = RGBColor(0xDE, 0xD8, 0xEA)
C_LINE  = RGBColor(0xC9, 0xC1, 0xDC)
C_CODEBG = RGBColor(0xF7, 0xF5, 0xFB)
C_GREEN = RGBColor(0x4E, 0x7D, 0x5E)
C_ORANGE = RGBColor(0xC2, 0x6B, 0x2E)


def blank_like(prs):
    s = prs.slides.add_slide(prs.slides[0].slide_layout)
    for sh in list(s.shapes):
        sh._element.getparent().remove(sh._element)
    return s


def text(sl, x, y, w, h, runs, align=PP_ALIGN.LEFT, ls=1.2, font=FONT):
    tb = sl.shapes.add_textbox(Emu(int(x * IN)), Emu(int(y * IN)),
                               Emu(int(w * IN)), Emu(int(h * IN)))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, (t, pt, bold, col) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = ls
        r = p.add_run(); r.text = t
        r.font.name = font; r.font.size = Pt(pt)
        r.font.bold = bold; r.font.color.rgb = col
    return tb


def box(sl, x, y, w, h, fill, line, lw=1.0, kind=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = sl.shapes.add_shape(kind, Emu(int(x * IN)), Emu(int(y * IN)),
                             Emu(int(w * IN)), Emu(int(h * IN)))
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line; sh.line.width = Pt(lw)
    sh.shadow.inherit = False; sh.text_frame.text = ''
    try:
        sh.adjustments[0] = 0.02
    except Exception:
        pass
    return sh


def seg(sl, x1, y1, x2, y2, col=C_LINE, w=2.0):
    c = sl.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,
                                Emu(int(x1 * IN)), Emu(int(y1 * IN)),
                                Emu(int(x2 * IN)), Emu(int(y2 * IN)))
    c.line.color.rgb = col; c.line.width = Pt(w)
    return c


def header(s, sub, title):
    text(s, 1.00, 0.92, 9.0, 0.30, [(sub, 18.75, False, C_HEAD)])
    text(s, 17.73, 0.92, 2.15, 0.30, [('발표자 정준', 18.75, False, C_PRES)])
    text(s, 1.00, 1.50, 19.0, 0.95, [(title, 42, True, C_TITLE)])


def steps(s, x, w, y0, blocks, pt_head=17, pt_body=13.5):
    y = y0
    for head, lines in blocks:
        text(s, x, y, w, 0.32, [(head, pt_head, True, C_TITLE)])
        y += 0.38
        for ln in lines:
            text(s, x + 0.26, y, w - 0.26, 0.28, [(ln, pt_body, False, C_BODY)])
            y += 0.28
        y += 0.22
    return y


# ════════════════════════════════════════════════════════════════
#  3-a  어떻게 찾나
# ════════════════════════════════════════════════════════════════
#  ⚠ 주석을 «별도 텍스트박스»로 위치 계산해 붙였다가 마지막 줄에서 코드와 겹쳤다
#    (한글·영문 폭이 달라 len × 상수로는 못 맞춘다). 한 줄을 통째로 한 run 으로 그린다.
RRF_CODE = [
    'score = {}',
    'for ids in ranked_id_lists:              # 벡터 순위 · 키워드 순위',
    '    for rank, cid in enumerate(ids):',
    '        score[cid] += 1.0 / (60 + rank + 1)',
]


def slide_a(prs):
    s = blank_like(prs)
    header(s, 'PART 03 · 잇다 — 찾기', '칸이 찼습니다. 이제 «찾습니다»')

    LX, LW = 1.00, 7.60
    steps(s, LX, LW, 2.66, [
        ('칸을 «찾을 말»로 바꿉니다', [
            '돕기·돌봄 · 사람 · 어르신  →  “돌봄 어르신 사람”',
        ]),
        ('두 가지 방법으로 «따로» 찾습니다', [
            '뜻으로 —  비슷한 의미를 찾습니다. 이름은 잘 못 잡습니다.',
            '이름으로 — 글자 그대로 찾습니다. 뜻은 못 봅니다.',
            '한쪽이 놓치는 걸 다른 쪽이 잡습니다.',
        ]),
        ('둘을 합칩니다 — 점수가 아니라 «순위»로', [
            '두 검색은 점수 단위가 다릅니다. 그냥 더하면 큰 쪽이 이깁니다.',
            '그래서 «몇 등이었나»만 봅니다. 양쪽에 다 나오면 위로 올라갑니다.',
        ]),
        ('마지막으로 다시 읽어 줄을 세웁니다', [
            '질문과 후보를 «붙여서 함께» 읽는 모델로 다시 정렬합니다.',
        ]),
    ])

    #  RRF 코드 — «중요한 데만»
    cy = 8.72
    box(s, LX, cy, LW, 1.62, C_CODEBG, C_PANLL, 0.75)
    text(s, LX + 0.30, cy + 0.18, 3.0, 0.26,
         [('합치는 부분 — match.py', 12, False, C_MUTED)])
    yy = cy + 0.50
    for ln in RRF_CODE:
        col = C_MUTED if ln.lstrip().startswith('#') else C_TITLE
        text(s, LX + 0.30, yy, LW - 0.60, 0.24, [(ln, 11.5, False, col)], font=CODE)
        yy += 0.25

    # ── 오른쪽 : 흐름 ──────────────────────────────────────
    RX, RW = 9.30, 9.70
    CX = RX + RW / 2

    def node(y, h, label, sub, fill, line, w=None, x=None):
        w = w or 5.20
        x = x if x is not None else CX - w / 2
        box(s, x, y, w, h, fill, line, 1.1)
        text(s, x, y + 0.16, w, 0.32, [(label, 17, True, C_TITLE)], align=PP_ALIGN.CENTER)
        if sub:
            text(s, x, y + 0.52, w, 0.26, [(sub, 12.5, False, C_MUTED)], align=PP_ALIGN.CENTER)
        return y + h

    y = 2.70
    node(y, 0.86, '“돌봄 어르신 사람”', '칸에서 만든 «찾을 말»', C_PANEL, C_PANLL)
    seg(s, CX, 3.56, CX, 3.92)

    #  두 갈래
    seg(s, CX - 2.40, 3.92, CX + 2.40, 3.92)
    seg(s, CX - 2.40, 3.92, CX - 2.40, 4.14)
    seg(s, CX + 2.40, 3.92, CX + 2.40, 4.14)
    node(4.14, 0.92, '뜻으로', '비슷한 의미를 찾는다', C_CARD, C_CARDL, w=4.30, x=CX - 4.55)
    node(4.14, 0.92, '이름으로', '글자 그대로 찾는다', C_CARD, C_CARDL, w=4.30, x=CX + 0.25)
    seg(s, CX - 2.40, 5.06, CX - 2.40, 5.32)
    seg(s, CX + 2.40, 5.06, CX + 2.40, 5.32)
    seg(s, CX - 2.40, 5.32, CX + 2.40, 5.32)
    seg(s, CX, 5.32, CX, 5.58)

    node(5.58, 0.86, '합치기', '양쪽에서 «공통으로» 오른 것을 위로', C_PANEL, C_PANLL)
    seg(s, CX, 6.44, CX, 6.70)
    node(6.70, 0.86, '다시 읽어 줄 세우기', '질문과 후보를 붙여서 «함께» 읽는다', C_PANEL, C_PANLL)
    seg(s, CX, 7.56, CX, 7.82)

    box(s, CX - 3.40, 7.82, 6.80, 1.30, C_CARD, C_GREEN, 1.25)
    text(s, CX - 3.40, 8.00, 6.80, 0.30,
         [('후보가 나왔습니다', 15, True, C_GREEN)], align=PP_ALIGN.CENTER)
    text(s, CX - 3.40, 8.38, 6.80, 0.56,
         [('요양지원 · 일상생활기능지원 · 아이돌봄 …', 15, True, C_TITLE),
          ('다음 장 — 이걸 «바로 줄까, 되물을까»', 12.5, False, C_MUTED)],
         align=PP_ALIGN.CENTER)
    return s


# ════════════════════════════════════════════════════════════════
#  3-b  어떻게 답을 내나
# ════════════════════════════════════════════════════════════════
def slide_b(prs):
    s = blank_like(prs)
    header(s, 'PART 03 · 잇다 — 답 내기', '바로 줄까, 한 번 더 물을까')

    LX, LW = 1.00, 7.60
    y = steps(s, LX, LW, 2.66, [
        ('후보가 «한 동네»에 모였나 봅니다', [
            '요양지원 · 일상생활기능지원 — 둘 다 돌봄입니다. 한 동네입니다.',
            '요양지원 · 웹개발 · 자동차정비 — 서로 다른 동네입니다.',
        ]),
        ('흩어져 있으면 «되묻습니다»', [
            '「돌봄 쪽 / 만들기 쪽 중 어디가 편하세요?」',
            '고르시면 그 방향으로 다시 찾습니다.',
        ]),
        ('모였으면 «카드»를 만듭니다', [
            '직업 하나를 고르고, 그 직업에 필요한 것들을 붙입니다.',
        ]),
        ('자격증이 없는 직업이면', [
            '억지로 붙이지 않습니다. 대신 국민내일배움카드를 안내합니다.',
            '요양지원이 실제로 그런 경우입니다.',
        ]),
    ])
    text(s, LX, y + 0.10, LW, 0.30,
         [('오른쪽 카드는 «실제로 3턴 대화해서» 받은 것입니다', 13.5, True, C_PRES)])

    # ── 오른쪽 : 진짜 카드 ─────────────────────────────────
    RX, RW = 9.30, 9.70
    box(s, RX, 2.60, RW, 7.70, C_CARD, C_PANLL, 1.5)

    text(s, RX + 0.46, 2.88, RW - 0.92, 0.28, [('🃏  직업 방향 카드', 13, False, C_MUTED)])
    text(s, RX + 0.46, 3.24, RW - 0.92, 0.56, [('요양지원', 34, True, C_TITLE)])
    text(s, RX + 0.46, 3.86, RW - 0.92, 0.28, [('보건  ·  NCS 06010108', 13, False, C_MUTED)])

    seg(s, RX + 0.46, 4.26, RX + RW - 0.46, 4.26, C_PANLL, 1.25)
    text(s, RX + 0.46, 4.40, RW - 0.92, 0.72,
         [('“효율적인 진료 및 간호인력의 지원업무를 위하여 진료지원보조, 환자이송지원,',
           13, False, C_BODY),
          ('일상생활수행지원, 신체청결지원, 인지활동지원보조 등을 수행하는 일”', 13, False, C_BODY),
          ('NCS 직무정의 — 우리가 쓴 말이 아닙니다', 11.5, False, C_MUTED)])

    text(s, RX + 0.46, 5.42, 2.0, 0.28, [('왜 이 방향인가', 14, True, C_PRES)])
    text(s, RX + 0.46, 5.76, RW - 0.92, 0.56,
         [('“어르신을 돌보는 일을 편하게 느끼신다고 하셔서', 13.5, False, C_BODY),
          ('직접적인 신체 지원이 포함된 요양지원 업무를 추천합니다”', 13.5, False, C_BODY)])

    seg(s, RX + 0.46, 6.52, RX + RW - 0.46, 6.52, C_PANLL, 1.25)
    box(s, RX + 0.46, 6.66, RW - 0.92, 0.90, RGBColor(0xFB, 0xF0, 0xE6), C_ORANGE, 1.1)
    text(s, RX + 0.76, 6.82, RW - 1.52, 0.28,
         [('자격증  —  이 방향은 국가기술자격이 «없습니다»', 14, True, C_ORANGE)])
    text(s, RX + 0.76, 7.14, RW - 1.52, 0.28,
         [('억지로 붙이지 않고, 국민내일배움카드로 훈련비 지원을 안내합니다', 13, False, C_BODY)])

    text(s, RX + 0.46, 7.76, 3.0, 0.28, [('배울 수 있는 강좌', 14, True, C_PRES)])
    for i, (t, sub) in enumerate([
            ('사회복지실천과 요양보호', 'K-MOOC · 원광보건대 · 무료'),
            ('와상환자를 위한 통합돌봄', 'K-MOOC · 무료')]):
        yy = 8.10 + i * 0.56
        text(s, RX + 0.76, yy, RW - 1.52, 0.28, [(t, 13.5, True, C_TITLE)])
        text(s, RX + 0.76, yy + 0.26, RW - 1.52, 0.24, [(sub, 11.5, False, C_MUTED)])

    box(s, RX + 0.46, 9.38, RW - 0.92, 0.62, C_PANEL, C_PANLL, 1.25)
    text(s, RX + 0.46, 9.55, RW - 0.92, 0.30,
         [('이 카드를 «미래설계지도»에 저장 — 나중에 이어서 대화할 수 있습니다',
           14, True, C_TITLE)], align=PP_ALIGN.CENTER)
    return s


def main():
    prs = Presentation(SRC)
    n0 = len(prs.slides)
    slide_a(prs); slide_b(prs)
    prs.save(OUT)
    W, H = prs.slide_width / IN, prs.slide_height / IN
    o = Presentation(OUT)
    over = [(i, round((sh.left + sh.width) / IN, 2), round((sh.top + sh.height) / IN, 2))
            for i in (-2, -1) for sh in o.slides[i].shapes
            if (sh.left + sh.width) / IN > W + 0.01 or (sh.top + sh.height) / IN > H + 0.01]
    print(f'  {n0}장 → {len(o.slides)}장  (3-a · 3-b)')
    print(f'  슬라이드 밖: {len(over)}개  {"✅" if not over else over}')
    print(f'  저장 → {OUT}')


main()
